"""
Deck Export Service - Export decks to PowerPoint and PDF
"""

import io
import logging
import json
import base64
import os
import tempfile
import hashlib
import requests
import uuid
import subprocess
import sys
from typing import Dict, List, Any, Optional, Union
from datetime import datetime
from app.utils.formatters import DeckFormatter
from app.services.deck_storage_service import deck_storage

# Import centralized data validator for safe operations
from app.services.data_validator import (
    ensure_numeric, safe_divide, safe_multiply, safe_get_value
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

# Always import reportlab components we need
from reportlab.lib.pagesizes import letter, LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.platypus.flowables import Image
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.linecharts import HorizontalLineChart

try:
    from playwright.sync_api import sync_playwright
    import subprocess
    import sys
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False

logger = logging.getLogger(__name__)


def _ensure_playwright_browser():
    """Ensure Playwright Chromium browser is installed, install if missing"""
    if not PLAYWRIGHT_AVAILABLE:
        return False
    
    try:
        # Try to launch browser to check if it's installed
        with sync_playwright() as p:
            try:
                browser = p.chromium.launch(headless=True)
                browser.close()
                logger.info("[PLAYWRIGHT] Chromium browser is available")
                return True
            except Exception as e:
                error_msg = str(e)
                # Check if it's a missing browser error
                if "Executable doesn't exist" in error_msg or "browserType.launch" in error_msg:
                    logger.warning("[PLAYWRIGHT] Chromium browser not found, attempting to install...")
                    try:
                        # Install Chromium browser
                        result = subprocess.run(
                            [sys.executable, "-m", "playwright", "install", "chromium"],
                            capture_output=True,
                            text=True,
                            timeout=300  # 5 minute timeout
                        )
                        if result.returncode == 0:
                            logger.info("[PLAYWRIGHT] Successfully installed Chromium browser")
                            # Verify installation by trying again
                            browser = p.chromium.launch(headless=True)
                            browser.close()
                            return True
                        else:
                            logger.error(f"[PLAYWRIGHT] Failed to install Chromium: {result.stderr}")
                            return False
                    except subprocess.TimeoutExpired:
                        logger.error("[PLAYWRIGHT] Browser installation timed out")
                        return False
                    except Exception as install_error:
                        logger.error(f"[PLAYWRIGHT] Error installing browser: {install_error}")
                        return False
                else:
                    # Some other error
                    logger.error(f"[PLAYWRIGHT] Browser launch error: {error_msg}")
                    return False
    except Exception as e:
        logger.error(f"[PLAYWRIGHT] Error checking browser: {e}")
        return False


async def _ensure_playwright_browser_async():
    """Async version: Ensure Playwright Chromium browser is installed, install if missing"""
    if not PLAYWRIGHT_AVAILABLE:
        return False
    
    try:
        from playwright.async_api import async_playwright
        async with async_playwright() as p:
            try:
                browser = await p.chromium.launch(headless=True)
                await browser.close()
                logger.info("[PLAYWRIGHT] Chromium browser is available")
                return True
            except Exception as e:
                error_msg = str(e)
                if "Executable doesn't exist" in error_msg or "browserType.launch" in error_msg:
                    logger.warning("[PLAYWRIGHT] Chromium browser not found, attempting to install...")
                    try:
                        result = subprocess.run(
                            [sys.executable, "-m", "playwright", "install", "chromium"],
                            capture_output=True,
                            text=True,
                            timeout=300
                        )
                        if result.returncode == 0:
                            logger.info("[PLAYWRIGHT] Successfully installed Chromium browser")
                            browser = await p.chromium.launch(headless=True)
                            await browser.close()
                            return True
                        else:
                            logger.error(f"[PLAYWRIGHT] Failed to install Chromium: {result.stderr}")
                            return False
                    except Exception as install_error:
                        logger.error(f"[PLAYWRIGHT] Error installing browser: {install_error}")
                        return False
                else:
                    logger.error(f"[PLAYWRIGHT] Browser launch error: {error_msg}")
                    return False
    except Exception as e:
        logger.error(f"[PLAYWRIGHT] Error checking browser: {e}")
        return False


async def _ensure_playwright_browser_async():
    """Async version: Ensure Playwright Chromium browser is installed, install if missing"""
    if not PLAYWRIGHT_AVAILABLE:
        return False
    
    try:
        from playwright.async_api import async_playwright
        async with async_playwright() as p:
            try:
                browser = await p.chromium.launch(headless=True)
                await browser.close()
                logger.info("[PLAYWRIGHT] Chromium browser is available")
                return True
            except Exception as e:
                error_msg = str(e)
                if "Executable doesn't exist" in error_msg or "browserType.launch" in error_msg:
                    logger.warning("[PLAYWRIGHT] Chromium browser not found, attempting to install...")
                    try:
                        result = subprocess.run(
                            [sys.executable, "-m", "playwright", "install", "chromium"],
                            capture_output=True,
                            text=True,
                            timeout=300
                        )
                        if result.returncode == 0:
                            logger.info("[PLAYWRIGHT] Successfully installed Chromium browser")
                            browser = await p.chromium.launch(headless=True)
                            await browser.close()
                            return True
                        else:
                            logger.error(f"[PLAYWRIGHT] Failed to install Chromium: {result.stderr}")
                            return False
                    except Exception as install_error:
                        logger.error(f"[PLAYWRIGHT] Error installing browser: {install_error}")
                        return False
                else:
                    logger.error(f"[PLAYWRIGHT] Browser launch error: {error_msg}")
                    return False
    except Exception as e:
        logger.error(f"[PLAYWRIGHT] Error checking browser: {e}")
        return False

# Import ChartRendererService for high-quality chart pre-rendering
from app.services.chart_renderer_service import chart_renderer


class DeckExportService:
    """Export deck data to PowerPoint and PDF formats"""
    
    # Neo-noir color palette - matches web theme
    COLORS = {
        'primary': RGBColor(8, 8, 8),        # Deep black
        'secondary': RGBColor(20, 20, 20),    # Dark gray
        'accent': RGBColor(45, 45, 45),      # Medium gray
        'highlight': RGBColor(70, 70, 70),   # Light gray
        'warning': RGBColor(40, 40, 40),     # Warning gray
        'text_light': RGBColor(245, 245, 245), # Light text
        'text_dark': RGBColor(8, 8, 8),      # Dark text
        'background': RGBColor(250, 250, 250), # Clean white background
        'border': RGBColor(224, 224, 224),    # Light border
    }
    
    def __init__(self):
        self.prs = None
        self.pdf_buffer = None
        
    def _validate_and_format_data(self, data: Any) -> str:
        """Validate and format data for display"""
        if data is None or data == "":
            return "—"
        
        # Format numbers
        if isinstance(data, (int, float)):
            if data >= 1_000_000_000:
                return f"${data/1_000_000_000:.2f}B"
            elif data >= 1_000_000:
                return f"${data/1_000_000:.1f}M"
            elif data >= 1_000:
                return f"${data/1_000:.0f}K"
            else:
                return f"${data:.0f}"
        
        # Format percentages
        if isinstance(data, str) and data.endswith("%"):
            return data
        
        # Clean up strings
        if isinstance(data, str):
            # Remove None strings
            if data.lower() == "none" or data.lower() == "n/a":
                return "—"
            # Truncate long strings
            if len(data) > 100:
                return data[:97] + "..."
            return data
        
        return str(data)
    
    def _format_metric_value(self, value: Any, metric_type: str = "currency") -> str:
        """Format metric values based on type"""
        if value is None:
            return "—"
        
        if metric_type == "currency":
            return self._validate_and_format_data(value)
        elif metric_type == "percentage":
            if isinstance(value, (int, float)):
                return f"{value:.1f}%"
            return str(value)
        elif metric_type == "number":
            if isinstance(value, (int, float)):
                return f"{value:,.0f}"
            return str(value)
        elif metric_type == "text":
            return self._validate_and_format_data(value)
        
        return str(value)
        
    def export_to_pptx(self, deck_data: Dict[str, Any]) -> bytes:
        """Export deck to PowerPoint format"""
        try:
            # Create presentation
            self.prs = Presentation()
            self.prs.slide_width = Inches(16)  # 16:9 aspect ratio
            self.prs.slide_height = Inches(9)
            
            # Process each slide
            # Handle both dict with slides key and direct list of slides
            if isinstance(deck_data, list):
                slides = deck_data
            else:
                slides = deck_data.get("slides", deck_data.get("deck_slides", []))
            for slide_data in slides:
                self._add_pptx_slide(slide_data)
            
            # Save to bytes buffer
            buffer = io.BytesIO()
            self.prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error exporting to PPTX: {e}")
            raise
    
    def _add_pptx_slide(self, slide_data: Dict[str, Any]):
        """Add a slide to the PowerPoint presentation"""
        slide_type = slide_data.get("type", "content")
        content = slide_data.get("content", {})
        
        if slide_type == "title":
            self._add_title_slide(content)
        elif slide_type == "summary":
            self._add_summary_slide(content)
        elif slide_type == "company":
            self._add_company_slide(content)
        elif slide_type == "chart":
            # Chart data is in content.chart_data, not slide_data.chart
            chart_data = content.get("chart_data") or slide_data.get("chart")
            self._add_chart_slide(content, chart_data)
        elif slide_type == "comparison":
            self._add_comparison_slide(content)
        elif slide_type == "investment_thesis":
            self._add_thesis_slide(content)
        elif slide_type == "investment_comparison":
            self._add_investment_comparison_slide(content)
        elif slide_type == "cap_table" or slide_type == "sankey":
            # Handle cap table/Sankey diagram slides
            self._add_cap_table_slide(content)
        elif slide_type == "cap_table_comparison":
            # Handle cap table comparison with Sankey diagrams
            self._add_cap_table_comparison_slide(content)
        elif slide_type == "scenario_comparison":
            # Handle scenario comparison with charts
            self._add_scenario_comparison_slide(content)
        elif slide_type == "business_model_comparison":
            # Handle business model comparison
            self._add_business_model_comparison_slide(content)
        elif slide_type == "side_by_side":
            # Handle side-by-side comparison slides
            self._add_side_by_side_slide(content)
        elif slide_type == "citations":
            # Handle citation slides
            self._add_citations_slide(content)
        elif slide_type == "exit_scenarios_comprehensive":
            # Handle comprehensive exit scenarios with charts and breakpoints
            self._add_exit_scenarios_comprehensive_slide(content)
        elif slide_type == "probability_cloud":
            # Handle probability cloud visualization
            self._add_probability_cloud_slide(content)
        elif slide_type == "breakpoint_analysis":
            # Handle breakpoint analysis slide
            self._add_breakpoint_analysis_slide(content)
        else:
            # Check if content has chart_data even if type isn't explicitly "chart"
            if content.get("chart_data"):
                chart_data = content.get("chart_data")
                # Check for special chart types
                if chart_data.get("type") == "sankey":
                    self._add_sankey_slide(content, chart_data)
                else:
                    self._add_chart_slide(content, chart_data)
            elif content.get("devices"):
                # Handle slides with devices (like side-by-side Sankey)
                self._add_devices_slide(content)
            else:
                self._add_content_slide(content)
    
    def _add_title_slide(self, content: Dict[str, Any]):
        """Add title slide with company logos"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = content.get("title", "Presentation")
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = content.get("subtitle", "") + "\n" + content.get("date", "")
        
        # Try to add company logos if available
        companies = content.get("companies", [])
        if companies:
            self._add_company_logos(slide, companies)
    
    def _add_summary_slide(self, content: Dict[str, Any]):
        """Add summary slide with bullets"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = content.get("title", "Summary")
        
        # Bullets
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear existing text
        
        bullets = content.get("bullets", [])
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    
    def _add_company_slide(self, content: Dict[str, Any]):
        """Add company profile slide"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Company Profile")
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Metrics table
        metrics = content.get("metrics", {})
        if metrics:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(7)
            height = Inches(5)
            
            # Create table shape
            rows = len(metrics) + 1
            cols = 2
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Header
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"
    
    def _add_company_logos(self, slide, companies: List[str]):
        """Add company logos to slide using Clearbit or fallback to initials"""
        try:
            logo_size = Inches(1.5)
            spacing = Inches(0.5)
            total_width = len(companies) * (logo_size + spacing) - spacing
            start_left = (self.prs.slide_width - total_width) / 2
            top = Inches(5)
            
            for i, company_name in enumerate(companies[:4]):  # Max 4 logos
                left = start_left + i * (logo_size + spacing)
                
                # Try to get logo from Clearbit
                logo_path = self._fetch_company_logo(company_name)
                
                if logo_path and os.path.exists(logo_path):
                    # Add actual logo image
                    slide.shapes.add_picture(logo_path, left, top, width=logo_size, height=logo_size)
                    # Clean up temp file
                    os.remove(logo_path)
                else:
                    # Create fallback initial badge
                    self._add_initial_badge(slide, company_name, left, top, logo_size)
                
                # Add company name below logo
                text_left = left
                text_top = top + logo_size + Inches(0.2)
                text_width = logo_size
                text_height = Inches(0.5)
                
                text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                text_frame = text_box.text_frame
                text_frame.text = company_name
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.size = Pt(12)
                
        except Exception as e:
            logger.warning(f"Failed to add company logos: {e}")
    
    def _fetch_company_logo(self, company_name: str) -> Optional[str]:
        """Fetch company logo from Clearbit API"""
        try:
            # Clean company name for domain search
            clean_name = company_name.lower().replace(' ', '').replace(',', '').replace('.', '')
            
            # Try common domain patterns
            domains = [
                f"{clean_name}.com",
                f"{clean_name}.ai",
                f"{clean_name}.io",
                f"get{clean_name}.com",
                f"use{clean_name}.com"
            ]
            
            for domain in domains:
                logo_url = f"https://logo.clearbit.com/{domain}"
                
                # Try to fetch the logo
                response = requests.get(logo_url, timeout=2)
                if response.status_code == 200:
                    # Save to temp file
                    temp_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
                    with open(temp_path, 'wb') as f:
                        f.write(response.content)
                    return temp_path
            
            return None
            
        except Exception as e:
            logger.debug(f"Could not fetch logo for {company_name}: {e}")
            return None
    
    def _add_initial_badge(self, slide, company_name: str, left, top, size):
        """Add a colored badge with company initials as fallback"""
        try:
            # Get initials
            words = company_name.split()
            if len(words) >= 2:
                initials = words[0][0].upper() + words[1][0].upper()
            else:
                initials = company_name[:2].upper()
            
            # Generate consistent neo-noir color from company name
            color_hash = hashlib.md5(company_name.encode()).hexdigest()
            # Use neo-noir palette variations
            base_colors = [
                (0, 0, 0),           # Pure black
                (45, 45, 45),        # Dark charcoal
                (74, 74, 74),        # Medium charcoal
                (107, 107, 107),     # Steel gray
                (138, 138, 138),     # Light steel
                (176, 176, 176),     # Silver
            ]
            color_index = int(color_hash[:2], 16) % len(base_colors)
            r, g, b = base_colors[color_index]
            
            # Create rectangle shape as background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,  # Rectangle
                left, top, size, size
            )
            
            # Set fill color
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()
            
            # Add text
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = initials
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Center vertically
            text_frame.margin_top = Inches(0.4)
            
        except Exception as e:
            logger.warning(f"Failed to create initial badge: {e}")
            
            # Data
            for i, (key, value) in enumerate(metrics.items(), 1):
                table.cell(i, 0).text = key
                table.cell(i, 1).text = str(value)
        
        # Investment thesis if present
        thesis = content.get("investment_thesis", "")
        if thesis:
            left = Inches(8)
            top = Inches(2)
            width = Inches(7)
            height = Inches(5)
            thesis_box = slide.shapes.add_textbox(left, top, width, height)
            thesis_frame = thesis_box.text_frame
            thesis_frame.text = thesis
            thesis_frame.word_wrap = True
    
    def _add_chart_slide(self, content: Dict[str, Any], chart_data: Optional[Dict[str, Any]]):
        """Add slide with chart"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(15)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Chart")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Check if chart is pre-rendered image
        if chart_data and chart_data.get("type") == "image":
            try:
                # Decode base64 image
                img_data = base64.b64decode(chart_data['src'].split(',')[1])
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    temp_file.write(img_data)
                    temp_file_path = temp_file.name
                
                # Add image to slide
                left = Inches(1)
                top = Inches(2)
                width = Inches(14)
                height = Inches(5.5)
                
                slide.shapes.add_picture(temp_file_path, left, top, width, height)
                
                # Clean up temp file
                os.unlink(temp_file_path)
                
                logger.info("Successfully added pre-rendered chart image to PowerPoint")
                
            except Exception as e:
                logger.error(f"Failed to add pre-rendered chart image: {e}")
                # Fallback to placeholder
                self._add_chart_placeholder(slide, chart_data.get("alt", "Chart"))
        
        # Try to add actual chart using python-pptx chart capabilities
        elif chart_data and chart_data.get("data"):
            data = chart_data.get("data", {})
            labels = data.get("labels", [])
            datasets = data.get("datasets", [{}])
            chart_type = chart_data.get("type", "bar").lower()
            
            if labels and datasets and datasets[0].get("data"):
                try:
                    from pptx.chart.data import ChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    
                    # Create chart data
                    chart_data_obj = ChartData()
                    chart_data_obj.categories = labels
                    
                    # Add series for each dataset
                    for dataset in datasets:
                        series_name = dataset.get("label", "Series")
                        series_values = dataset.get("data", [])
                        if series_values:
                            chart_data_obj.add_series(series_name, series_values)
                    
                    # Determine chart type
                    if chart_type == "bar":
                        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
                    elif chart_type == "line":
                        chart_type_enum = XL_CHART_TYPE.LINE
                    elif chart_type == "pie":
                        chart_type_enum = XL_CHART_TYPE.PIE
                    else:
                        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
                    
                    # Add chart to slide
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(14)
                    height = Inches(5.5)
                    
                    chart = slide.shapes.add_chart(
                        chart_type_enum, left, top, width, height, chart_data_obj
                    ).chart
                    
                    # Style the chart
                    if hasattr(chart, "has_legend"):
                        chart.has_legend = True
                        if chart.legend:
                            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                            chart.legend.include_in_layout = False
                    
                    # Add chart title if available
                    if hasattr(chart, "chart_title") and chart_data.get("title"):
                        chart.has_title = True
                        chart.chart_title.text_frame.text = chart_data.get("title")
                    
                except Exception as e:
                    logger.error(f"Chart creation failed: {e}")
                    import traceback
                    logger.error(f"Full traceback: {traceback.format_exc()}")
                    # Use the existing fallback method with correct slide parameter
                    self._add_chart_as_table(slide, chart_data, Inches(1), Inches(2), Inches(14), Inches(4))
            else:
                # No valid data for chart
                logger.warning(f"No valid data for chart in slide '{content.get('title', 'Unknown')}'")
                self._add_chart_as_table(slide, chart_data, Inches(1), Inches(2), Inches(14), Inches(4))
        else:
            # No data - add placeholder
            left = Inches(2)
            top = Inches(2)
            width = Inches(12)
            height = Inches(5)
            chart_box = slide.shapes.add_textbox(left, top, width, height)
            chart_frame = chart_box.text_frame
            chart_frame.text = "[Chart Placeholder - No Data]"
    
    def _add_chart_placeholder(self, slide, alt_text: str):
        """Add chart placeholder when image loading fails"""
        left = Inches(2)
        top = Inches(2)
        width = Inches(12)
        height = Inches(5)
        chart_box = slide.shapes.add_textbox(left, top, width, height)
        chart_frame = chart_box.text_frame
        chart_frame.text = f"[Chart: {alt_text}]"
    
    def _add_chart_as_table(self, slide, data, labels, datasets):
        """Add chart data as a formatted table when native chart fails"""
        left = Inches(2)
        top = Inches(2)
        width = Inches(12)
        height = Inches(5)
        
        if labels and datasets and datasets[0].get("data"):
            # Create table
            rows = len(labels) + 1
            cols = len(datasets) + 1
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Headers
            table.cell(0, 0).text = "Category"
            for col_idx, dataset in enumerate(datasets, 1):
                table.cell(0, col_idx).text = dataset.get("label", f"Series {col_idx}")
            
            # Data
            for row_idx, label in enumerate(labels, 1):
                table.cell(row_idx, 0).text = str(label)
                for col_idx, dataset in enumerate(datasets, 1):
                    values = dataset.get("data", [])
                    if row_idx - 1 < len(values):
                        value = values[row_idx - 1]
                        # Format numbers nicely
                        if isinstance(value, (int, float)):
                            table.cell(row_idx, col_idx).text = f"{value:,.0f}"
                        else:
                            table.cell(row_idx, col_idx).text = str(value)
    
    def _add_comparison_slide(self, content: Dict[str, Any]):
        """Add comparison slide"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(15)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Comparison")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Companies table
        companies = content.get("companies", [])
        if companies:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(15)
            height = Inches(5)
            
            # Create table
            rows = len(companies) + 1
            cols = 4  # Name, Valuation, Revenue, Stage
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Headers
            table.cell(0, 0).text = "Company"
            table.cell(0, 1).text = "Valuation"
            table.cell(0, 2).text = "Revenue"
            table.cell(0, 3).text = "Stage"
            
            # Data
            for i, company in enumerate(companies, 1):
                table.cell(i, 0).text = company.get("name", "")
                table.cell(i, 1).text = f"${company.get('valuation', 0):,.0f}"
                table.cell(i, 2).text = f"${company.get('revenue', 0):,.0f}"
                table.cell(i, 3).text = company.get("stage", "")
    
    def _add_thesis_slide(self, content: Dict[str, Any]):
        """Add investment thesis slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = content.get("title", "Investment Thesis")
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        # Add thesis points
        if "thesis_points" in content:
            for point in content["thesis_points"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        
        # Add metrics
        if "key_metrics" in content:
            p = tf.add_paragraph()
            p.text = "\nKey Metrics:"
            p.level = 0
            p.font.bold = True
            
            for key, value in content["key_metrics"].items():
                p = tf.add_paragraph()
                p.text = f"{key}: {value}"
                p.level = 1
    
    def _add_content_slide(self, content: Dict[str, Any]):
        """Add generic content slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = content.get("title", "Content")
        
        # Body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        # Add bullets if present
        bullets = content.get("bullets", [])
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        
        # Add body text if present
        body_text = content.get("body", "")
        if body_text and not bullets:
            p = tf.add_paragraph()
            p.text = body_text
    
    def _add_cap_table_slide(self, content: Dict[str, Any]):
        """Add cap table slide with pie chart showing current ownership"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout for custom positioning
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(15)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Cap Table")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Check for pie chart data first (preferred format)
        chart_data = content.get("chart_data")
        if chart_data and chart_data.get("type", "").lower() == "pie":
            # Render pie chart
            try:
                from pptx.chart.data import ChartData
                from pptx.enum.chart import XL_CHART_TYPE
                
                data = chart_data.get("data", {})
                labels = data.get("labels", [])
                datasets = data.get("datasets", [{}])
                
                if labels and datasets and datasets[0].get("data"):
                    # Create chart data
                    chart_data_obj = ChartData()
                    chart_data_obj.categories = labels
                    
                    # Add series for pie chart (single series)
                    series_values = datasets[0].get("data", [])
                    if series_values:
                        chart_data_obj.add_series("Ownership", series_values)
                    
                    # Add pie chart to slide
                    chart_left = Inches(2)
                    chart_top = Inches(1.5)
                    chart_width = Inches(8)
                    chart_height = Inches(6)
                    
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.PIE, chart_left, chart_top, chart_width, chart_height, chart_data_obj
                    ).chart
                    
                    # Style the chart
                    if hasattr(chart, "has_legend"):
                        chart.has_legend = True
                        if chart.legend:
                            chart.legend.position = XL_LEGEND_POSITION.RIGHT
                            chart.legend.include_in_layout = True
                    
                    # Add chart title if available
                    if hasattr(chart, "chart_title") and chart_data.get("title"):
                        chart.has_title = True
                        chart.chart_title.text_frame.text = chart_data.get("title")
                    
                    logger.info(f"[CAP_TABLE] Successfully rendered pie chart for cap table slide")
                    
                    # Add subtitle if present
                    subtitle = content.get("subtitle", "")
                    if subtitle:
                        subtitle_box = slide.shapes.add_textbox(Inches(10.5), Inches(1.5), Inches(5), Inches(1))
                        subtitle_frame = subtitle_box.text_frame
                        subtitle_frame.text = subtitle
                        subtitle_frame.paragraphs[0].font.size = Pt(14)
                        subtitle_frame.paragraphs[0].font.italic = True
                    
                    # Add bullets if present
                    bullets = content.get("bullets", [])
                    if bullets:
                        bullets_top = Inches(2.5)
                        bullets_left = Inches(10.5)
                        bullets_width = Inches(5)
                        bullets_box = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, Inches(4))
                        bullets_frame = bullets_box.text_frame
                        bullets_frame.clear()
                        for bullet in bullets[:6]:  # Limit to 6 bullets
                            p = bullets_frame.add_paragraph()
                            p.text = f"• {bullet}"
                            p.font.size = Pt(10)
                            p.level = 0
                    
                    # Add metrics if present
                    metrics = content.get("metrics", {})
                    if metrics:
                        metrics_top = Inches(6.5)
                        metrics_left = Inches(10.5)
                        metrics_width = Inches(5)
                        metrics_box = slide.shapes.add_textbox(metrics_left, metrics_top, metrics_width, Inches(1))
                        metrics_frame = metrics_box.text_frame
                        metrics_frame.clear()
                        for key, value in list(metrics.items())[:4]:  # Limit to 4 metrics
                            p = metrics_frame.add_paragraph()
                            p.text = f"{key}: {value}"
                            p.font.size = Pt(10)
                            p.level = 0
                    
                    return  # Successfully rendered pie chart
                    
            except Exception as e:
                logger.error(f"[CAP_TABLE] Failed to render pie chart: {e}")
                import traceback
                logger.error(f"[CAP_TABLE] Stack trace: {traceback.format_exc()}")
                # Fall through to table rendering
        
        # Fallback: Get company data for table rendering
        companies = content.get("companies", [])
        
        # Also check if investor_details are in content directly
        investor_details = content.get("investor_details", [])
        current_cap_table = content.get("current_cap_table", {})
        
        if len(companies) >= 2:
            # Company 1 - Left side
            # Add investor details to company data if available
            if investor_details and len(companies) > 0:
                companies[0]["investor_details"] = [inv for inv in investor_details if inv.get("round") in [r.get("round", "") for r in companies[0].get("funding_rounds", [])]]
            if current_cap_table and len(companies) > 0:
                companies[0]["current_cap_table"] = current_cap_table
            self._add_cap_table_section(slide, companies[0], Inches(0.5), Inches(1.5), Inches(7.5))
            
            # Company 2 - Right side
            if investor_details and len(companies) > 1:
                companies[1]["investor_details"] = [inv for inv in investor_details if inv.get("round") in [r.get("round", "") for r in companies[1].get("funding_rounds", [])]]
            if current_cap_table and len(companies) > 1:
                companies[1]["current_cap_table"] = current_cap_table
            self._add_cap_table_section(slide, companies[1], Inches(8.5), Inches(1.5), Inches(7.5))
        elif len(companies) == 1:
            # Single company - centered
            if investor_details:
                companies[0]["investor_details"] = investor_details
            if current_cap_table:
                companies[0]["current_cap_table"] = current_cap_table
            self._add_cap_table_section(slide, companies[0], Inches(4.5), Inches(1.5), Inches(7.5))
        elif investor_details or current_cap_table:
            # No companies but we have investor details - show them directly
            company_data = {
                "name": content.get("title", "Cap Table").replace("Cap Table - ", ""),
                "investor_details": investor_details,
                "current_cap_table": current_cap_table
            }
            self._add_cap_table_section(slide, company_data, Inches(2), Inches(1.5), Inches(12))
        else:
            # No companies and no chart data - add placeholder
            placeholder_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(3))
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.text = "[Cap Table Data Not Available]"
            placeholder_frame.paragraphs[0].font.size = Pt(16)
            placeholder_frame.paragraphs[0].font.italic = True
    
    def _add_cap_table_section(self, slide, company_data: Dict[str, Any], left: float, top: float, width: float):
        """Add individual cap table section for a company"""
        # Company name
        name_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.text = company_data.get("name", "Company")
        name_frame.paragraphs[0].font.size = Pt(20)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Try to get investor details from slide content first
        investor_details = company_data.get("investor_details", [])
        current_cap_table = company_data.get("current_cap_table", {})
        
        # If we have investor details, show them in a table
        if investor_details:
            # Create table with investor names and ownership
            rows = len(investor_details) + 1
            cols = 2  # Investor Name, Ownership %
            
            table_top = top + Inches(0.7)
            table = slide.shapes.add_table(rows, cols, left, table_top, width, Inches(min(4, len(investor_details) * 0.3))).table
            
            # Headers
            table.cell(0, 0).text = "Investor"
            table.cell(0, 1).text = "Ownership"
            
            # Data - show investor names and ownership
            for i, investor in enumerate(investor_details[:15], 1):  # Limit to 15 investors
                investor_name = investor.get("name", "Unknown")
                ownership = investor.get("ownership", 0)
                table.cell(i, 0).text = investor_name[:40]  # Truncate long names
                table.cell(i, 1).text = f"{ownership:.1f}%"
            
            # Style the table
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    if row == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
            return
        
        # Fallback: Use current_cap_table if available
        if current_cap_table:
            # Extract investor names from cap table
            investor_rows = []
            for owner, pct in current_cap_table.items():
                owner_str = str(owner)
                pct_float = float(pct) if isinstance(pct, (int, float)) else 0
                # Check if this is an investor (not founder, not employee)
                is_investor = (
                    ('Investor' in owner_str or 'Series' in owner_str or 'Seed' in owner_str or 'Round' in owner_str) 
                    and 'Founder' not in owner_str 
                    and 'Employee' not in owner_str 
                    and 'Option' not in owner_str 
                    and 'ESOP' not in owner_str
                    and pct_float > 0.1
                )
                if is_investor:
                    clean_name = owner_str.replace(' (Lead)', '').replace(' (SAFE)', '').strip()
                    investor_rows.append((clean_name, pct_float))
            
            if investor_rows:
                # Sort by ownership descending
                investor_rows.sort(key=lambda x: x[1], reverse=True)
                
                # Create table
                rows = len(investor_rows[:15]) + 1  # Limit to 15 investors
                cols = 2
                table_top = top + Inches(0.7)
                table = slide.shapes.add_table(rows, cols, left, table_top, width, Inches(min(4, len(investor_rows) * 0.3))).table
                
                # Headers
                table.cell(0, 0).text = "Investor"
                table.cell(0, 1).text = "Ownership"
                
                # Data
                for i, (name, pct) in enumerate(investor_rows[:15], 1):
                    table.cell(i, 0).text = name[:40]
                    table.cell(i, 1).text = f"{pct:.1f}%"
                
                # Style
                for row in range(rows):
                    for col in range(cols):
                        cell = table.cell(row, col)
                        cell.text_frame.paragraphs[0].font.size = Pt(9)
                        if row == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
                return
        
        # Cap table data as table (legacy format)
        cap_table_data = company_data.get("cap_table", {})
        rounds = cap_table_data.get("rounds", [])
        
        if rounds:
            # Create ownership evolution table
            rows = len(rounds) + 1
            cols = 4  # Round, Founders %, Investors %, Employees %
            
            table_top = top + Inches(0.7)
            table = slide.shapes.add_table(rows, cols, left, table_top, width, Inches(3)).table
            
            # Headers
            table.cell(0, 0).text = "Round"
            table.cell(0, 1).text = "Founders"
            table.cell(0, 2).text = "Investors"
            table.cell(0, 3).text = "Employees"
            
            # Data
            for i, round_data in enumerate(rounds, 1):
                table.cell(i, 0).text = round_data.get("round", "")
                table.cell(i, 1).text = f"{round_data.get('founders', 0):.1f}%"
                table.cell(i, 2).text = f"{round_data.get('investors', 0):.1f}%"
                table.cell(i, 3).text = f"{round_data.get('employees', 0):.1f}%"
            
            # Add styling
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if row == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
        
        # Add investment highlights if present
        highlights = company_data.get("highlights", [])
        if highlights:
            highlight_top = top + Inches(4)
            highlight_box = slide.shapes.add_textbox(left, highlight_top, width, Inches(1.5))
            tf = highlight_box.text_frame
            tf.text = "Key Points:"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            
            for highlight in highlights[:3]:  # Limit to 3 highlights
                p = tf.add_paragraph()
                p.text = f"• {highlight}"
                p.font.size = Pt(10)
                p.level = 0
    
    def _add_sankey_slide(self, content: Dict[str, Any], chart_data: Dict[str, Any]):
        """Add slide with Sankey diagram visualization"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(15)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Ownership Flow")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Since PowerPoint doesn't natively support Sankey diagrams,
        # we'll create a table representation of the flow
        sankey_data = chart_data.get("data", {})
        nodes = sankey_data.get("nodes", [])
        links = sankey_data.get("links", [])
        
        if nodes and links:
            # Create a flow table
            left = Inches(2)
            top = Inches(2)
            width = Inches(12)
            height = Inches(5)
            
            # Group links by source for better visualization
            flow_map = {}
            for link in links:
                source_idx = link.get("source", 0)
                target_idx = link.get("target", 0)
                value = link.get("value", 0)
                
                if source_idx < len(nodes) and target_idx < len(nodes):
                    source_name = nodes[source_idx].get("name", "Source")
                    target_name = nodes[target_idx].get("name", "Target")
                    
                    if source_name not in flow_map:
                        flow_map[source_name] = []
                    flow_map[source_name].append((target_name, value))
            
            # Create table
            rows = sum(len(targets) for targets in flow_map.values()) + 1
            cols = 3
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Headers
            table.cell(0, 0).text = "From"
            table.cell(0, 1).text = "To"
            table.cell(0, 2).text = "Ownership %"
            
            # Data
            row_idx = 1
            for source, targets in flow_map.items():
                for target, value in targets:
                    table.cell(row_idx, 0).text = source
                    table.cell(row_idx, 1).text = target
                    table.cell(row_idx, 2).text = f"{value:.1f}%"
                    row_idx += 1
            
            # Style the table
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(11)
                    if row == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
    
    def _add_side_by_side_slide(self, content: Dict[str, Any]):
        """Add side-by-side comparison slide"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(15)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Side-by-Side Comparison")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Get left and right content
        left_content = content.get("left", {})
        right_content = content.get("right", {})
        
        # Left side
        if left_content:
            self._add_side_content(slide, left_content, Inches(0.5), Inches(1.5), Inches(7.5))
        
        # Right side
        if right_content:
            self._add_side_content(slide, right_content, Inches(8.5), Inches(1.5), Inches(7.5))
    
    def _add_side_content(self, slide, content: Dict[str, Any], left: float, top: float, width: float):
        """Add content for one side of a side-by-side slide"""
        # Subtitle
        subtitle = content.get("subtitle", "")
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(18)
            subtitle_frame.paragraphs[0].font.bold = True
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            top += Inches(0.7)
        
        # Check for chart data
        chart_data = content.get("chart_data")
        if chart_data:
            # Add placeholder for chart
            chart_box = slide.shapes.add_textbox(left, top, width, Inches(3))
            tf = chart_box.text_frame
            tf.text = f"[{chart_data.get('type', 'Chart').upper()} VISUALIZATION]\n"
            
            # Add data points as text
            data = chart_data.get("data", {})
            labels = data.get("labels", [])
            datasets = data.get("datasets", [])
            
            if labels and datasets:
                for dataset in datasets[:1]:  # Show first dataset
                    values = dataset.get("data", [])
                    for i, (label, value) in enumerate(zip(labels[:5], values[:5])):
                        p = tf.add_paragraph()
                        p.text = f"{label}: {value:,.0f}"
                        p.font.size = Pt(10)
        
        # Metrics
        metrics = content.get("metrics", {})
        if metrics:
            metrics_top = top + (Inches(3.5) if chart_data else Inches(0))
            for i, (key, value) in enumerate(list(metrics.items())[:5]):
                metric_box = slide.shapes.add_textbox(left, metrics_top + (i * Inches(0.5)), width, Inches(0.5))
                tf = metric_box.text_frame
                tf.text = f"{key}: {value}"
                tf.paragraphs[0].font.size = Pt(11)
        
        # Bullets
        bullets = content.get("bullets", [])
        if bullets:
            bullets_top = top + (Inches(3) if chart_data else Inches(0)) + (Inches(2.5) if metrics else Inches(0))
            bullet_box = slide.shapes.add_textbox(left, bullets_top, width, Inches(2))
            tf = bullet_box.text_frame
            
            for bullet in bullets[:4]:  # Limit bullets
                p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(10)
                p.level = 0
    
    def _add_investment_comparison_slide(self, content: Dict[str, Any]):
        """Add investment comparison slide with radar chart"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Investment Comparison"))
        
        # Add companies side by side
        companies = content.get("companies", [])
        if companies:
            # Add company details and chart
            chart_data = content.get("chart_data")
            if chart_data:
                self._add_chart_with_fallback(slide, chart_data, Inches(1), Inches(2), Inches(14), Inches(5))
    
    def _add_cap_table_comparison_slide(self, content: Dict[str, Any]):
        """Add cap table comparison slide with Sankey diagrams"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Cap Table Comparison"))
        
        # Check for devices with Sankey data
        devices = content.get("devices", [])
        if devices:
            for device in devices:
                if device.get("type") == "side_by_side_sankey":
                    # Add Sankey visualization or fallback
                    self._add_sankey_fallback(slide, device)
        
        # Add metrics if present
        metrics = content.get("metrics", {})
        if metrics:
            self._add_metrics_table(slide, metrics, Inches(1), Inches(5))
    
    def _add_scenario_comparison_slide(self, content: Dict[str, Any]):
        """Add scenario comparison slide with charts"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Scenario Analysis"))
        
        # Add chart if present
        chart_data = content.get("chart_data")
        if chart_data:
            self._add_chart_with_fallback(slide, chart_data, Inches(1), Inches(1.5), Inches(14), Inches(5.5))
    
    def _add_business_model_comparison_slide(self, content: Dict[str, Any]):
        """Add business model comparison slide"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Business Model Analysis"))
        
        # Add companies data
        companies = content.get("companies", [])
        if companies:
            # Create comparison table
            comparison_table = content.get("comparison_table")
            if comparison_table and comparison_table.get("data"):
                self._add_comparison_table(slide, comparison_table, Inches(1), Inches(1.5))
    
    def _add_citations_slide(self, content: Dict[str, Any]):
        """Add citations slide"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Sources & References"))
        
        # Add citations
        citations = content.get("citations", [])
        if citations:
            top = Inches(1.5)
            for i, citation in enumerate(citations[:10]):  # Limit to 10 per slide
                text_box = slide.shapes.add_textbox(Inches(0.5), top + (i * Inches(0.5)), Inches(15), Inches(0.5))
                tf = text_box.text_frame
                tf.text = f"{i+1}. {citation}"
                tf.paragraphs[0].font.size = Pt(9)
    
    def _add_devices_slide(self, content: Dict[str, Any]):
        """Add slide with devices (special visualizations)"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        self._add_slide_title(slide, content.get("title", "Analysis"))
        
        # Process devices
        devices = content.get("devices", [])
        for device in devices:
            device_type = device.get("type")
            if device_type == "side_by_side_sankey":
                self._add_sankey_fallback(slide, device)
    
    def _add_slide_title(self, slide, title_text: str):
        """Helper to add title to slide"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_chart_with_fallback(self, slide, chart_data: Dict[str, Any], left, top, width, height):
        """Add chart with fallback to table if chart creation fails"""
        
        # Handle special chart types
        chart_type_str = chart_data.get("type", "bar").lower()
        
        # Probability cloud needs special handling
        if chart_type_str == "probability_cloud":
            self._add_probability_cloud_fallback(slide, chart_data, left, top, width, height)
            return
        
        try:
            # Try to create actual chart
            if chart_data and chart_data.get("data"):
                data = chart_data.get("data", {})
                labels = data.get("labels", [])
                datasets = data.get("datasets", [{}])
                
                if labels and datasets and datasets[0].get("data"):
                    # Create chart data
                    chart_data_obj = ChartData()
                    chart_data_obj.categories = labels
                    
                    # Add series for each dataset
                    for dataset in datasets:
                        series_name = dataset.get("label", "Series")
                        series_values = dataset.get("data", [])
                        if series_values:
                            chart_data_obj.add_series(series_name, series_values)
                    
                    # Determine chart type
                    if chart_type_str in ["bar", "column"]:
                        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
                    elif chart_type_str == "line":
                        chart_type_enum = XL_CHART_TYPE.LINE
                    elif chart_type_str == "pie":
                        chart_type_enum = XL_CHART_TYPE.PIE
                    elif chart_type_str == "area":
                        chart_type_enum = XL_CHART_TYPE.AREA
                    else:
                        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
                    
                    chart = slide.shapes.add_chart(
                        chart_type_enum, left, top, width, height, chart_data_obj
                    ).chart
                    
                    # Add chart title if available
                    if hasattr(chart, "chart_title") and chart_data.get("title"):
                        chart.has_title = True
                        chart.chart_title.text_frame.text = chart_data.get("title")
                    
                    return  # Success
        except Exception as e:
            logger.warning(f"Chart creation failed, using fallback: {e}")
        
        # Fallback to table representation
        self._add_chart_as_table(slide, chart_data, left, top, width, height)
    
    def _add_chart_as_table(self, slide, chart_data: Dict[str, Any], left, top, width, height):
        """Add chart data as a table when chart rendering fails"""
        data = chart_data.get("data", {})
        labels = data.get("labels", [])
        datasets = data.get("datasets", [])
        
        if labels and datasets:
            rows = len(labels) + 1
            cols = len(datasets) + 1
            
            table = slide.shapes.add_table(rows, cols, left, top, width, Inches(3)).table
            
            # Headers
            table.cell(0, 0).text = "Category"
            for i, dataset in enumerate(datasets):
                table.cell(0, i + 1).text = dataset.get("label", f"Series {i+1}")
            
            # Data
            for i, label in enumerate(labels):
                table.cell(i + 1, 0).text = str(label)
                for j, dataset in enumerate(datasets):
                    values = dataset.get("data", [])
                    if i < len(values):
                        table.cell(i + 1, j + 1).text = f"{values[i]:,.0f}" if isinstance(values[i], (int, float)) else str(values[i])
            
            # Style the table
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if row == 0 or col == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
    
    def _add_probability_cloud_fallback(self, slide, chart_data: Dict[str, Any], left, top, width, height):
        """Add probability cloud visualization for PowerPoint - try to render as line chart first"""
        data = chart_data.get("data", {})
        scenario_curves = data.get("scenario_curves", [])
        breakpoint_clouds = data.get("breakpoint_clouds", [])
        insights = data.get("insights", {})
        
        try:
            # Try to create actual line chart with multiple series
            if scenario_curves and len(scenario_curves) > 0:
                # Get common x-axis values (exit values) from first scenario
                first_curve = scenario_curves[0].get("return_curve", {})
                exit_values = first_curve.get("exit_values", [])
                
                if exit_values:
                    # Create chart data
                    chart_data_obj = ChartData()
                    
                    # Set categories (exit values in millions)
                    categories = [f"${val/1e6:.0f}M" if val < 1e9 else f"${val/1e9:.1f}B" 
                                  for val in exit_values[::3]]  # Sample every 3rd point for readability
                    chart_data_obj.categories = categories
                    
                    # Add top 5-8 most probable scenarios as series
                    sorted_scenarios = sorted(scenario_curves, 
                                            key=lambda s: s.get("probability", 0), 
                                            reverse=True)[:8]
                    
                    for scenario in sorted_scenarios:
                        return_curve = scenario.get("return_curve", {})
                        return_multiples = return_curve.get("return_multiples", [])
                        if return_multiples:
                            # Sample every 3rd point to match categories
                            sampled_returns = return_multiples[::3]
                            scenario_name = scenario.get("name", "Scenario")[:20]  # Limit name length
                            chart_data_obj.add_series(scenario_name, sampled_returns)
                    
                    # Create line chart
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.LINE, left, top, width, Inches(4), chart_data_obj
                    ).chart
                    
                    # Configure chart
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_data.get("title", "Return Scenarios")
                    
                    # Make lines smoother with markers
                    for series in chart.series:
                        series.smooth = True
                        series.marker.style = XL_MARKER_STYLE.CIRCLE
                        series.marker.size = 3
                    
                    top += Inches(4.5)
                    
                    # Add breakpoint summary below chart
                    if breakpoint_clouds:
                        self._add_breakpoint_summary(slide, breakpoint_clouds, left, top)
                        top += Inches(1.5)
                    
                    # Add insights
                    if insights:
                        self._add_insights_box(slide, insights, left, top, width)
                    
                    return  # Success - no need for fallback
                    
        except Exception as e:
            logger.warning(f"Failed to create line chart for probability cloud: {e}")
        
        # Fallback to table representation if chart creation fails
        self._add_probability_cloud_table_fallback(slide, chart_data, left, top, width, height)
    
    def _add_probability_cloud_table_fallback(self, slide, chart_data: Dict[str, Any], left, top, width, height):
        """Fallback table representation for probability cloud"""
        data = chart_data.get("data", {})
        scenario_curves = data.get("scenario_curves", [])
        breakpoint_clouds = data.get("breakpoint_clouds", [])
        insights = data.get("insights", {})
        
        # Add title
        if chart_data.get("title"):
            title_shape = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            title_frame = title_shape.text_frame
            title_frame.text = chart_data["title"]
            title_frame.paragraphs[0].font.size = Pt(14)
            title_frame.paragraphs[0].font.bold = True
            top += Inches(0.6)
        
        # Create scenarios summary table - NO LIMIT, show all scenarios
        if scenario_curves:
            num_scenarios = len(scenario_curves)
            
            table = slide.shapes.add_table(
                min(num_scenarios + 1, 15),  # Cap at 15 rows for space
                5,  # columns
                left, top, width, Inches(3)
            ).table
            
            # Headers
            headers = ["Scenario", "Probability", "Ownership", "Final LP", "Max Return"]
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Add scenario data - show ALL scenarios dynamically
            display_scenarios = scenario_curves[:min(num_scenarios, 14)]  # Leave room for header
            for i, scenario in enumerate(display_scenarios):
                table.cell(i + 1, 0).text = (scenario.get("name", "")[:25])
                table.cell(i + 1, 1).text = f"{safe_multiply(scenario.get('probability', 0), 100, 0):.0f}%"
                
                # Show final ownership percentage
                final_ownership = scenario.get("final_ownership", 0)
                table.cell(i + 1, 2).text = f"{final_ownership * 100:.1f}%"
                
                # Show final liquidation preference
                final_lp = scenario.get("final_liq_pref", 0)
                if final_lp >= 1e9:
                    table.cell(i + 1, 3).text = f"${final_lp/1e9:.1f}B"
                else:
                    table.cell(i + 1, 3).text = f"${final_lp/1e6:.0f}M"
                
                # Show max return from return curve
                return_curve = scenario.get("return_curve", {})
                return_multiples = return_curve.get("return_multiples", [])
                # Filter out non-numeric values and ensure we have valid data
                numeric_multiples = [m for m in return_multiples if isinstance(m, (int, float)) and not isinstance(m, bool)]
                max_return = max(numeric_multiples) if numeric_multiples else 0
                table.cell(i + 1, 4).text = f"{max_return:.1f}x"
                
                # Style cells
                for j in range(5):
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
            
            top += Inches(2.2)
        
        # Add breakpoint distributions summary  
        if breakpoint_clouds:
            self._add_breakpoint_summary(slide, breakpoint_clouds, left, top)
            top += Inches(1.7)
        
        # Add key insights
        if insights:
            self._add_insights_box(slide, insights, left, top, width)
    
    def _add_breakpoint_summary(self, slide, breakpoint_clouds: list, left, top):
        """Add breakpoint summary table to slide"""
        if not breakpoint_clouds:
            return
            
        bp_table = slide.shapes.add_table(
            len(breakpoint_clouds) + 1, 4,
            left, top, Inches(6), Inches(1.5)
        ).table
        
        # Headers
        bp_headers = ["Breakpoint", "Median", "P25-P75 Range", "Confidence"]
        for i, header in enumerate(bp_headers):
            cell = bp_table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Add breakpoint data
        for i, cloud in enumerate(breakpoint_clouds):
            bp_table.cell(i + 1, 0).text = cloud.get("label", "")
            
            median = cloud.get("median", 0)
            if median >= 1e9:
                bp_table.cell(i + 1, 1).text = f"${median/1e9:.1f}B"
            else:
                bp_table.cell(i + 1, 1).text = f"${median/1e6:.0f}M"
            
            p25_p75 = cloud.get("p25_p75", [0, 0])
            p25 = p25_p75[0] if len(p25_p75) > 0 else 0
            p75 = p25_p75[1] if len(p25_p75) > 1 else 0
            
            if p75 >= 1e9:
                bp_table.cell(i + 1, 2).text = f"${p25/1e9:.1f}B-${p75/1e9:.1f}B"
            else:
                bp_table.cell(i + 1, 2).text = f"${p25/1e6:.0f}M-${p75/1e6:.0f}M"
                
            # Confidence based on range tightness
            range_ratio = (p75 - p25) / median if median > 0 else float('inf')
            if range_ratio < 0.5:
                confidence = "High"
            elif range_ratio < 1.0:
                confidence = "Medium"
            else:
                confidence = "Low"
            bp_table.cell(i + 1, 3).text = confidence
            
            # Style cells
            for j in range(4):
                bp_table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
    
    def _add_insights_box(self, slide, insights: Dict[str, Any], left, top, width):
        """Add insights box to slide"""
        insights_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        insights_frame = insights_box.text_frame
        insights_frame.text = "Key Insights:"
        insights_frame.paragraphs[0].font.bold = True
        insights_frame.paragraphs[0].font.size = Pt(11)
        
        p = insights_frame.add_paragraph()
        insights_text = []
        if insights.get("probability_of_3x"):
            insights_text.append(f"• {insights['probability_of_3x']*100:.0f}% probability of 3x return")
        if insights.get("expected_breakeven"):
            val = insights['expected_breakeven']
            if val >= 1e9:
                insights_text.append(f"• Expected breakeven at ${val/1e9:.1f}B exit")
            else:
                insights_text.append(f"• Expected breakeven at ${val/1e6:.0f}M exit")
        if insights.get("median_liq_pref"):
            val = insights['median_liq_pref']
            if val >= 1e9:
                insights_text.append(f"• Median liquidation preference: ${val/1e9:.1f}B")
            else:
                insights_text.append(f"• Median liquidation preference: ${val/1e6:.0f}M")
        
        p.text = "\n".join(insights_text)
        p.font.size = Pt(10)

    def _add_sankey_fallback(self, slide, device: Dict[str, Any]):
        """Add Sankey diagram fallback visualization"""
        # Since PowerPoint doesn't support Sankey natively, create a flow table
        data = device.get("data", {})
        company1_data = data.get("company1_data", {})
        company2_data = data.get("company2_data", {})
        company1_name = data.get("company1_name", "Company 1")
        company2_name = data.get("company2_name", "Company 2")
        
        # Add company 1 flow
        if company1_data:
            self._add_flow_table(slide, company1_data, company1_name, Inches(1), Inches(2))
        
        # Add company 2 flow
        if company2_data:
            self._add_flow_table(slide, company2_data, company2_name, Inches(8.5), Inches(2))
    
    def _add_flow_table(self, slide, flow_data: Dict[str, Any], company_name: str, left, top):
        """Add flow table representation of Sankey data"""
        # Add company name
        name_box = slide.shapes.add_textbox(left, top, Inches(6), Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.text = company_name
        name_frame.paragraphs[0].font.size = Pt(16)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Create flow table from links
        links = flow_data.get("links", [])
        nodes = flow_data.get("nodes", [])
        
        if links:
            # Create a simple flow representation
            rows = min(len(links) + 1, 10)  # Limit rows
            table = slide.shapes.add_table(rows, 3, left, top + Inches(0.7), Inches(6), Inches(3)).table
            
            # Headers
            table.cell(0, 0).text = "From"
            table.cell(0, 1).text = "To"
            table.cell(0, 2).text = "Stake %"
            
            # Add flow data
            for i, link in enumerate(links[:rows-1]):
                source_idx = link.get("source", 0)
                target_idx = link.get("target", 0)
                value = link.get("value", 0)
                
                # Get node names
                source_name = nodes[source_idx].get("name", "Unknown") if source_idx < len(nodes) else "Unknown"
                target_name = nodes[target_idx].get("name", "Unknown") if target_idx < len(nodes) else "Unknown"
                
                table.cell(i + 1, 0).text = source_name
                table.cell(i + 1, 1).text = target_name
                table.cell(i + 1, 2).text = f"{value:.1f}%"
            
            # Style the table
            for row in range(rows):
                for col in range(3):
                    cell = table.cell(row, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    if row == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
    
    def _add_metrics_table(self, slide, metrics: Dict[str, Any], left, top):
        """Add metrics comparison table"""
        if metrics:
            # Count companies
            companies = list(metrics.keys())
            if companies:
                rows = 4  # Typical metrics rows
                cols = len(companies) + 1
                
                table = slide.shapes.add_table(rows, cols, left, top, Inches(10), Inches(2)).table
                
                # Headers
                table.cell(0, 0).text = "Metric"
                for i, company in enumerate(companies):
                    table.cell(0, i + 1).text = company
                
                # Add metrics (if they exist in the data)
                metric_rows = [
                    ("Founder %", "Founder %"),
                    ("Total Raised", "Total Raised"),
                    ("Our Stake", "Our Stake")
                ]
                
                for row_idx, (label, key) in enumerate(metric_rows, 1):
                    table.cell(row_idx, 0).text = label
                    for col_idx, company in enumerate(companies):
                        company_metrics = metrics.get(company, {})
                        value = company_metrics.get(key, "N/A")
                        table.cell(row_idx, col_idx + 1).text = str(value)
                
                # Style
                for row in range(rows):
                    for col in range(cols):
                        cell = table.cell(row, col)
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        if row == 0 or col == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
    
    def _add_comparison_table(self, slide, table_data: Dict[str, Any], left, top):
        """Add comparison table to slide"""
        data = table_data.get("data", [])
        if data:
            rows = len(data)
            cols = len(data[0]) if data else 0
            
            if rows > 0 and cols > 0:
                table = slide.shapes.add_table(rows, cols, left, top, Inches(14), Inches(5)).table
                
                for row_idx, row_data in enumerate(data):
                    for col_idx, cell_value in enumerate(row_data):
                        table.cell(row_idx, col_idx).text = str(cell_value)
                        # Bold first row and column
                        if row_idx == 0 or col_idx == 0:
                            table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.bold = True
                        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
    
    def _add_exit_scenarios_comprehensive_slide(self, content: Dict[str, Any]):
        """Add comprehensive exit scenarios slide with charts to PowerPoint"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Exit Scenarios & Ownership Analysis")
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content.get("subtitle")
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        
        # Companies data in tables
        companies = content.get("companies", {})
        if companies:
            # Create two columns for company data
            company_names = list(companies.keys())[:2]
            for idx, company_name in enumerate(company_names):
                x_pos = Inches(0.5 + idx * 7.5)
                y_pos = Inches(1.8)
                
                company_data = companies[company_name]
                
                # Company name
                name_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(7), Inches(0.5))
                name_frame = name_box.text_frame
                name_frame.text = company_name
                name_frame.paragraphs[0].font.size = Pt(18)
                name_frame.paragraphs[0].font.bold = True
                
                # Add breakpoints table if available
                if company_data.get("breakpoints"):
                    breakpoints = company_data["breakpoints"]
                    table = slide.shapes.add_table(4, 2, x_pos, y_pos + Inches(0.6), Inches(7), Inches(1.5)).table
                    
                    # Headers
                    table.cell(0, 0).text = "Breakpoint"
                    table.cell(0, 1).text = "Value"
                    
                    # Data
                    row = 1
                    if breakpoints.get("liquidation_preference"):
                        table.cell(row, 0).text = "Liquidation Preference"
                        table.cell(row, 1).text = f"${breakpoints['liquidation_preference']/1e6:.1f}M"
                        row += 1
                    if breakpoints.get("conversion_point"):
                        table.cell(row, 0).text = "Conversion Point"
                        table.cell(row, 1).text = f"${breakpoints['conversion_point']/1e6:.1f}M"
                        row += 1
                    if breakpoints.get("target_3x_exit"):
                        table.cell(row, 0).text = "3x Target"
                        table.cell(row, 1).text = f"${breakpoints['target_3x_exit']/1e6:.1f}M"
                    
                    # Format table
                    for row in range(4):
                        for col in range(2):
                            cell = table.cell(row, col)
                            cell.text_frame.paragraphs[0].font.size = Pt(10)
                            if row == 0:
                                cell.text_frame.paragraphs[0].font.bold = True
                
                # Add scenarios summary
                if company_data.get("scenarios"):
                    y_pos += Inches(2.2)
                    scenarios_text = "Top Scenarios:\n"
                    for scenario in company_data["scenarios"][:3]:
                        scenarios_text += f"• {scenario['scenario']}: ${scenario['exit_value']/1e6:.0f}M ({scenario['probability']*100:.0f}%)\n"
                    
                    scenario_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(7), Inches(1.5))
                    scenario_frame = scenario_box.text_frame
                    scenario_frame.text = scenarios_text
                    scenario_frame.paragraphs[0].font.size = Pt(10)
    
    def _add_probability_cloud_slide(self, content: Dict[str, Any]):
        """Add probability cloud visualization slide to PowerPoint"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Exit Probability Cloud")
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content.get("subtitle")
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        
        # Add chart placeholder note
        note_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
        note_frame = note_box.text_frame
        note_frame.text = "Probability Cloud Visualization\n\n"
        note_frame.text += "• Shows return distribution across exit scenarios\n"
        note_frame.text += "• Includes defensive breakpoints\n"
        note_frame.text += "• Decision zones for investment thresholds\n"
        note_frame.text += "\n(View in PDF for full interactive chart)"
        note_frame.paragraphs[0].font.size = Pt(14)
    
    def _add_breakpoint_analysis_slide(self, content: Dict[str, Any]):
        """Add breakpoint analysis slide to PowerPoint"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Breakpoint Analysis")
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content.get("subtitle")
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        
        # Reality check table if available
        if content.get("reality_check"):
            reality_data = content["reality_check"]
            if reality_data:
                # Create table
                rows = min(len(reality_data) + 1, 10)  # Header + data rows (max 10)
                cols = 5
                table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(15), Inches(4)).table
                
                # Headers
                headers = ["Company", "Exit Value", "Our Return", "MOIC", "DPI Impact"]
                for col, header in enumerate(headers):
                    table.cell(0, col).text = header
                    table.cell(0, col).text_frame.paragraphs[0].font.bold = True
                    table.cell(0, col).text_frame.paragraphs[0].font.size = Pt(11)
                
                # Data
                for row_idx, row_data in enumerate(reality_data[:rows-1], 1):
                    table.cell(row_idx, 0).text = row_data.get('company', '')
                    table.cell(row_idx, 1).text = row_data.get('exit_value', '')
                    table.cell(row_idx, 2).text = row_data.get('our_proceeds', '')
                    table.cell(row_idx, 3).text = row_data.get('moic', '')
                    table.cell(row_idx, 4).text = row_data.get('dpi_contribution', '')
                    
                    # Format cells
                    for col in range(5):
                        table.cell(row_idx, col).text_frame.paragraphs[0].font.size = Pt(10)
        
        # Insights if available
        if content.get("insights"):
            insights = content["insights"]
            if isinstance(insights, list):
                insight_text = "Key Insights:\n" + "\n".join(f"• {insight}" for insight in insights[:5])
            else:
                insight_text = str(insights)
            
            insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(15), Inches(2))
            insights_frame = insights_box.text_frame
            insights_frame.text = insight_text
            insights_frame.paragraphs[0].font.size = Pt(11)
    
    def _generate_html_deck(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]], chart_images: Dict[str, str] = None) -> str:
        """Generate professional HTML deck with Tailwind CSS and Chart.js"""
        slides_html = []
        
        # Handle both dict with slides key and direct list of slides
        if isinstance(deck_data, list):
            slides = deck_data
        else:
            slides = deck_data.get("slides", deck_data.get("deck_slides", []))
        
        logger.info(f"[HTML_DECK] Generating HTML for {len(slides)} slides")
        
        for slide_idx, slide_data in enumerate(slides):
            slide_type = slide_data.get("type", "unknown")
            logger.info(f"[HTML_DECK] Processing slide {slide_idx + 1}/{len(slides)}: type={slide_type}")
            try:
                slide_html = self._generate_html_slide(slide_data, slide_idx, chart_images)
                slides_html.append(slide_html)
            except Exception as e:
                logger.error(f"[HTML_DECK] Error generating slide {slide_idx + 1} (type={slide_type}): {e}")
                # Continue processing other slides
        
        # Create complete HTML document
        html = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Investment Analysis Report</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
    <script src="https://d3js.org/d3.v7.min.js"></script>
    <script src="https://unpkg.com/d3-sankey@0.12.3/dist/d3-sankey.min.js"></script>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
    <style>
        /* Monochrome Design System - Landing Page Aesthetic */
        :root {{
            --background: 220 20% 98%;
            --foreground: 224 28% 12%;
            --muted-foreground: 224 12% 38%;
            --primary: 220 22% 20%;
            --secondary: 220 15% 92%;
            --border: 220 13% 88%;
            --brand-gradient-start: #111827;
            --brand-gradient-end: #1f2937;
        }}
        
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        html {{
            width: 1024px;
            max-width: 1024px;
            margin: 0 auto;
            overflow-x: hidden;
        }}
        
        body {{
            font-family: 'Inter', system-ui, -apple-system, 'Segoe UI', Roboto, sans-serif;
            -webkit-font-smoothing: antialiased;
            -moz-osx-font-smoothing: grayscale;
            background: hsl(var(--background));
            color: hsl(var(--foreground));
            line-height: 1.6;
            width: 1024px;
            max-width: 1024px;
            margin: 0 auto;
            overflow-x: hidden;
        }}
        
        /* Ensure all containers respect max-width */
        .container, .chart-container, table, .grid {{
            max-width: 100%;
            overflow-x: auto;
        }}
        
        .slide {{
            width: 1024px;
            min-height: 768px;
            padding: 3rem;
            page-break-after: always;
            page-break-inside: avoid;
            background: #ffffff;
            box-shadow: 0 20px 40px -12px rgba(0, 0, 0, 0.08);
            display: flex;
            flex-direction: column;
            position: relative;
            overflow: hidden;
        }}
        
        @media print {{
            .slide {{
                page-break-after: always;
                page-break-inside: avoid;
                border: none;
            }}
        }}
        
        /* Typography - Professional Hierarchy */
        h1, h2, h3, h4, h5, h6 {{
            font-family: 'Inter', sans-serif;
            color: hsl(220, 22%, 20%);
            font-weight: 700;
        }}
        
        .slide-title {{
            font-size: 2.5rem;
            font-weight: 600;
            line-height: 1.2;
            letter-spacing: -0.02em;
            margin-bottom: 2rem;
            color: hsl(224, 28%, 12%);
        }}
        
        .slide-subtitle {{
            font-size: 1.25rem;
            font-weight: 500;
            line-height: 1.4;
            color: hsl(224, 12%, 38%);
            margin-bottom: 1.5rem;
        }}
        
        .slide-body {{
            font-size: 1rem;
            font-weight: 400;
            line-height: 1.6;
            color: hsl(224, 28%, 12%);
        }}
        
        .label-text {{
            font-size: 0.75rem;
            font-weight: 500;
            text-transform: uppercase;
            letter-spacing: 0.3em;
            color: hsl(224, 12%, 38%);
        }}
        
        /* Chart Containers - Monochrome Style */
        .chart-container {{
            position: relative;
            width: 100%;
            height: 400px;
            padding: 1.5rem;
            background: linear-gradient(135deg, rgba(15, 23, 42, 0.02) 0%, rgba(30, 41, 59, 0.03) 100%);
            border-radius: 16px;
            border: 1px solid hsl(var(--border));
            box-shadow: 0 10px 30px -12px rgba(15, 23, 42, 0.15);
        }}
        
        /* Marketing Cards - Landing Page Style */
        .metric-card, .marketing-card {{
            background: linear-gradient(145deg, rgba(255, 255, 255, 0.75), rgba(248, 250, 252, 0.95));
            border-radius: 24px;
            border: 1px solid rgba(15, 23, 42, 0.08);
            box-shadow: 0 18px 55px -25px rgba(15, 23, 42, 0.25);
            backdrop-filter: blur(14px);
            padding: 1.5rem;
            text-align: center;
        }}
        
        .metric-value {{
            font-family: 'Inter', sans-serif;
            font-size: 2rem;
            font-weight: 700;
            color: hsl(220, 22%, 20%);
            margin-bottom: 0.5rem;
            letter-spacing: -0.02em;
        }}
        
        .metric-label {{
            font-family: 'Inter', sans-serif;
            font-size: 0.75rem;
            font-weight: 500;
            color: hsl(224, 12%, 38%);
            text-transform: uppercase;
            letter-spacing: 0.3em;
        }}
        
        /* Tables - Clean Professional Style */
        .professional-table {{
            width: 100%;
            border-collapse: collapse;
            font-size: 0.875rem;
            margin: 1.5rem 0;
        }}
        
        .professional-table th {{
            background: hsl(220, 15%, 92%);
            border: 1px solid hsl(220, 13%, 88%);
            padding: 0.75rem;
            text-align: left;
            font-weight: 600;
            color: hsl(224, 28%, 12%);
            font-size: 0.75rem;
            text-transform: uppercase;
            letter-spacing: 0.05em;
        }}
        
        .professional-table td {{
            border: 1px solid hsl(220, 13%, 88%);
            padding: 0.75rem;
            color: hsl(224, 28%, 12%);
        }}
        
        .professional-table tr:nth-child(even) {{
            background: hsl(220, 20%, 98%);
        }}
        
        /* Title Slide - Brand Gradient */
        .title-slide {{
            background: linear-gradient(120deg, var(--brand-gradient-start), var(--brand-gradient-end));
            color: #f9fafb !important;
            position: relative;
        }}
        
        .company-logo {{
            position: absolute;
            top: 2rem;
            right: 2rem;
            width: 60px;
            height: 60px;
            background: white;
            border-radius: 12px;
            padding: 8px;
            box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
        }}
        
        .title-slide * {{
            color: #ffffff !important;
        }}
        
        .title-slide h1, .title-slide h2, .title-slide h3,
        .title-slide p, .title-slide .slide-title {{
            color: #ffffff !important;
        }}
        
        /* Content Slides */
        .content-slide {{
            background: #ffffff;
        }}
        
        .section-slide {{
            background: linear-gradient(145deg, #ffffff, hsl(220, 20%, 98%));
        }}
        
        /* Professional Bullets */
        .bullet-list {{
            list-style: none;
            padding: 0;
            margin: 1rem 0;
        }}
        
        .bullet-list li {{
            padding-left: 1.5rem;
            margin-bottom: 0.75rem;
            position: relative;
            color: hsl(224, 28%, 12%);
        }}
        
        .bullet-list li:before {{
            content: "•";
            position: absolute;
            left: 0;
            color: hsl(220, 22%, 20%);
            font-weight: bold;
        }}
    </style>
</head>
<body class="bg-gray-50">
    {''.join(slides_html)}
    <script>
        // Initialize charts - wait for Chart.js and D3 to load
        function initializeCharts() {{
            // Check if Chart.js is loaded
            if (typeof Chart === 'undefined') {{
                console.log('Chart.js not loaded yet, waiting...');
                setTimeout(initializeCharts, 100);
                return;
            }}
            
            // Check if D3 is loaded (for Sankey and other D3 charts)
            if (typeof d3 === 'undefined') {{
                console.log('D3.js not loaded yet, waiting...');
                setTimeout(initializeCharts, 100);
                return;
            }}
            
            console.log('Chart.js and D3.js loaded, initializing charts...');
            try {{
                {self._generate_chart_scripts(deck_data)}
                console.log('All charts initialized successfully');
            }} catch(e) {{
                console.error('Error initializing charts:', e);
            }}
        }}
        
        // Start initialization when DOM is ready
        if (document.readyState === 'loading') {{
            document.addEventListener('DOMContentLoaded', initializeCharts);
        }} else {{
            initializeCharts();
        }}
    </script>
</body>
</html>
        """
        return html

    def _generate_html_slide(self, slide_data: Dict[str, Any], slide_idx: int, chart_images: Dict[str, str] = None) -> str:
        """Generate individual slide HTML"""
        slide_type = slide_data.get("type", "content")
        content = slide_data.get("content", {})
        
        # Generate base slide HTML
        slide_html = ""
        if slide_type == "title":
            slide_html = self._html_title_slide(content, slide_idx)
        elif slide_type == "summary":
            slide_html = self._html_summary_slide(content, slide_idx)
        elif slide_type == "company":
            slide_html = self._html_company_slide(content, slide_idx)
        elif slide_type == "chart" or content.get("chart_data"):
            chart_data = content.get("chart_data", {})
            if chart_data.get("type") == "sankey":
                slide_html = self._html_sankey_slide(content, slide_idx, chart_images)
            else:
                slide_html = self._html_chart_slide(content, slide_idx, chart_images)
        elif slide_type == "comparison":
            slide_html = self._html_comparison_slide(content, slide_idx)
        elif slide_type == "investment_thesis":
            slide_html = self._html_thesis_slide(content, slide_idx)
        elif slide_type == "investment_recommendations":
            slide_html = self._html_investment_recommendations_slide(content, slide_idx)
        elif slide_type == "company_comparison":
            slide_html = self._html_company_comparison_slide(content, slide_idx)
        elif slide_type == "path_to_100m_comparison":
            slide_html = self._html_path_to_100m_slide(content, slide_idx)
        elif slide_type == "business_analysis_comparison":
            slide_html = self._html_business_analysis_slide(content, slide_idx)
        elif slide_type == "tam_pincer":
            slide_html = self._html_tam_pincer_slide(content, slide_idx)
        elif slide_type == "founder_team_analysis":
            slide_html = self._html_founder_team_slide(content, slide_idx)
        elif slide_type == "cap_table_comparison":
            slide_html = self._html_cap_table_comparison_slide(content, slide_idx)
        elif slide_type == "citations":
            slide_html = self._html_citations_slide(content, slide_idx)
        elif slide_type == "cap_table" or slide_type == "sankey":
            slide_html = self._html_cap_table_slide(content, slide_idx)
        elif slide_type == "side_by_side":
            slide_html = self._html_side_by_side_slide(content, slide_idx)
        elif slide_type == "exit_scenarios_comprehensive":
            slide_html = self._html_exit_scenarios_comprehensive_slide(content, slide_idx)
        elif slide_type == "followon_strategy_table":
            slide_html = self._html_followon_strategy_slide(content, slide_idx)
        elif slide_type == "fund_return_impact_enhanced":
            slide_html = self._html_fund_impact_slide(content, slide_idx)
        elif slide_type == "risk_analysis":
            slide_html = self._html_risk_analysis_slide(content, slide_idx)
        elif slide_type == "probability_cloud":
            slide_html = self._html_probability_cloud_slide(content, slide_idx)
        elif slide_type == "breakpoint_analysis":
            slide_html = self._html_breakpoint_analysis_slide(content, slide_idx)
        else:
            slide_html = self._html_content_slide(content, slide_idx)
        
        # Add citations to all slides (except citations slides themselves)
        if slide_type != "citations" and content.get("citations"):
            citations = content.get("citations", [])
            slide_html = self._add_citations_to_slide(slide_html, citations)
        
        return slide_html
    
    def _html_title_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate title slide HTML with company logos"""
        # Extract company logos if available
        company_logos = content.get('company_logos', [])
        logo_html = ''
        if company_logos:
            logo_html = '<div class="flex justify-center gap-4 mt-8">'
            for logo in company_logos[:2]:  # Show max 2 logos
                if logo:
                    logo_html += f'<div class="company-logo"><img src="{logo}" alt="Company Logo" style="width: 100%; height: 100%; object-fit: contain;"/></div>'
            logo_html += '</div>'
        
        # Ensure title doesn't overflow - use responsive sizing
        title_text = content.get('title', 'Investment Analysis')
        subtitle_text = content.get('subtitle', '')
        date_text = content.get('date', datetime.now().strftime('%B %Y'))
        
        return f"""
<div class="slide title-slide flex items-center justify-center">
    <div class="text-center" style="width: 100%; padding: 2rem 4rem; max-width: 100%; box-sizing: border-box;">
        <h1 style="color: #f9fafb; font-size: 3rem; font-weight: 700; margin-bottom: 1.5rem; line-height: 1.2; word-wrap: break-word; overflow-wrap: break-word;">{title_text}</h1>
        {f'<p style="color: #f9fafb; opacity: 0.9; font-size: 1.5rem; margin-bottom: 1.5rem; line-height: 1.4; word-wrap: break-word;">{subtitle_text}</p>' if subtitle_text else ''}
        {f'<p style="color: #f9fafb; opacity: 0.8; font-size: 1.125rem; margin-bottom: 2rem;">{date_text}</p>' if date_text else ''}
        {logo_html}
    </div>
</div>
        """
    
    def _html_summary_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate summary slide HTML with clean theme"""
        bullets = content.get('bullets', [])
        bullets_html = ''.join([
            f'<li class="mb-3" style="color: hsl(224, 20%, 25%);">{bullet}</li>' 
            for bullet in bullets
        ])
        
        return f"""
<div class="slide">
    <h2 class="slide-title">{content.get('title', 'Executive Summary')}</h2>
    {f'<p class="slide-subtitle">{content.get("subtitle", "")}</p>' if content.get('subtitle') else ''}
    {f'<ul class="list-disc list-inside text-lg mt-6">{bullets_html}</ul>' if bullets else ''}
    {f'<p class="slide-body mt-4">{content.get("body", "")}</p>' if content.get('body') else ''}
</div>
        """
    
    def _html_company_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate company profile slide HTML with dark theme"""
        metrics = content.get('metrics', {})
        metrics_html = ''
        
        # Determine metric type for proper formatting
        metric_types = {
            'revenue': 'currency', 'arr': 'currency', 'valuation': 'currency',
            'funding': 'currency', 'runway': 'text', 'growth_rate': 'percentage',
            'team_size': 'number', 'customers': 'number', 'stage': 'text',
            'founded': 'text', 'location': 'text', 'business_model': 'text'
        }
        
        for key, value in metrics.items():
            # Format value based on metric type
            metric_key = key.lower().replace(' ', '_')
            metric_type = metric_types.get(metric_key, 'text')
            formatted_value = self._format_metric_value(value, metric_type)
            
            metrics_html += f"""
            <div class="metric-card">
                <p class="label-text">{key}</p>
                <p class="value-text">{formatted_value}</p>
            </div>
            """
        
        return f"""
<div class="slide">
    <h2 class="slide-title">{content.get('title', 'Company Profile')}</h2>
    <div class="grid grid-cols-3 gap-6">
        {metrics_html}
    </div>
    {f'<div class="mt-8 p-6 rounded-xl" style="background: rgba(31, 41, 55, 0.3); border: 1px solid rgba(255, 255, 255, 0.1);"><p class="text-gray-300">{content.get("investment_thesis", "")}</p></div>' if content.get('investment_thesis') else ''}
</div>
        """
    
    def _html_chart_slide(self, content: Dict[str, Any], slide_idx: int, chart_images: Dict[str, str] = None) -> str:
        """Generate chart slide HTML with dark theme"""
        chart_data = content.get('chart_data', {})
        subtitle = content.get('subtitle', '')
        
        # Check if we have a pre-rendered image for this chart
        chart_key = f"slide_{slide_idx}_chart_0"
        chart_html = ""
        
        if chart_images and chart_key in chart_images:
            # Use pre-rendered high-quality PNG
            img_base64 = chart_images[chart_key]
            chart_html = f'<img src="data:image/png;base64,{img_base64}" style="width: 100%; height: 400px; object-fit: contain;" />'
            logger.info(f"[HTML_CHART] Using pre-rendered image for slide {slide_idx}")
        else:
            # Fall back to canvas for simple charts
            chart_html = f'<canvas id="chart-{slide_idx}"></canvas>'
            logger.info(f"[HTML_CHART] Using canvas for slide {slide_idx}")
        
        return f"""
<div class="slide">
    <h2 class="slide-title">{content.get('title', 'Chart')}</h2>
    {f'<p class="slide-subtitle">{subtitle}</p>' if subtitle else ''}
    <div class="chart-container">
        {chart_html}
    </div>
</div>
        """
    
    def _html_comparison_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate comparison slide HTML with dark theme"""
        companies = content.get('companies', [])
        comparison_html = ''
        
        for company in companies[:2]:  # Show 2 companies side by side
            # Format metrics properly
            valuation = self._validate_and_format_data(company.get('valuation', 0))
            revenue = self._validate_and_format_data(company.get('revenue', 0))
            stage = company.get('stage', 'N/A')
            
            comparison_html += f"""
            <div class="flex-1 rounded-xl p-6" style="background: hsl(220, 20%, 98%); border: 1px solid hsl(220, 13%, 88%);">
                <h3 class="text-2xl font-bold mb-4" style="color: hsl(220, 22%, 20%);">{company.get('name', 'Company')}</h3>
                <div class="space-y-3">
                    <div><span style="color: hsl(224, 12%, 45%);">Valuation:</span> <span class="font-semibold" style="color: hsl(224, 28%, 12%);">{valuation}</span></div>
                    <div><span style="color: hsl(224, 12%, 45%);">Revenue:</span> <span class="font-semibold" style="color: hsl(224, 28%, 12%);">{revenue}</span></div>
                    <div><span style="color: hsl(224, 12%, 45%);">Stage:</span> <span class="font-semibold" style="color: hsl(224, 28%, 12%);">{stage}</span></div>
                </div>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-8">{content.get('title', 'Company Comparison')}</h2>
    <div class="flex gap-8">
        {comparison_html}
    </div>
</div>
        """
    
    def _html_thesis_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate investment thesis slide HTML"""
        thesis_points = content.get('thesis_points', [])
        points_html = ''.join([f'<li class="mb-3">{point}</li>' for point in thesis_points])
        
        metrics = content.get('key_metrics', {})
        metrics_html = ''
        for key, value in metrics.items():
            metrics_html += f"""
            <div class="bg-gray-50 rounded-lg p-4 border-l-3 border-gray-800">
                <p class="text-sm text-gray-600 mb-1">{key}</p>
                <p class="text-xl font-semibold text-gray-900">{value}</p>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-6">{content.get('title', 'Investment Thesis')}</h2>
    {f'<ul class="list-disc list-inside text-lg text-gray-700 mb-8">{points_html}</ul>' if thesis_points else ''}
    {f'<div class="grid grid-cols-3 gap-4">{metrics_html}</div>' if metrics else ''}
</div>
        """
    
    def _html_content_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate generic content slide HTML"""
        bullets = content.get('bullets', [])
        bullets_html = ''.join([f'<li class="mb-3 text-gray-700">{bullet}</li>' for bullet in bullets])
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-6">{content.get('title', 'Content')}</h2>
    {f'<ul class="list-disc list-inside text-lg">{bullets_html}</ul>' if bullets else ''}
    {f'<p class="text-lg text-gray-700">{content.get("body", "")}</p>' if content.get('body') else ''}
</div>
        """
    
    def _html_investment_recommendations_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate investment recommendations slide HTML"""
        recommendations = content.get('recommendations', [])
        recs_html = ''
        
        for rec in recommendations:
            # Determine color scheme based on recommendation
            recommendation_text = rec.get('recommendation', '')
            if 'INVEST' in recommendation_text:
                bg_color = 'bg-green-50 border-green-500'
                text_color = 'text-green-800'
            elif 'CONSIDER' in recommendation_text:
                bg_color = 'bg-yellow-50 border-yellow-500'
                text_color = 'text-yellow-800'
            else:
                bg_color = 'bg-red-50 border-red-500'
                text_color = 'text-red-800'
            
            recs_html += f"""
            <div class="rounded-lg p-6 border-2 {bg_color} mb-4">
                <h3 class="text-xl font-bold mb-2">{rec.get('company', 'Unknown')}</h3>
                <div class="{text_color} text-lg font-semibold mb-3">{recommendation_text}</div>
                <div class="text-gray-700 mb-2">{rec.get('reasoning', '')}</div>
                <div class="grid grid-cols-3 gap-4 mt-4 pt-4 border-t border-gray-300">
                    <div>
                        <span class="text-sm text-gray-600">Ownership</span>
                        <div class="font-semibold">{rec.get('ownership_details', 'N/A')}</div>
                    </div>
                    <div>
                        <span class="text-sm text-gray-600">Expected Proceeds</span>
                        <div class="font-semibold">{rec.get('expected_proceeds', 'TBD')}</div>
                    </div>
                    <div>
                        <span class="text-sm text-gray-600">Expected IRR</span>
                        <div class="font-semibold">{rec.get('expected_irr', 'TBD')}</div>
                    </div>
                </div>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Investment Recommendations')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    <div>
        {recs_html}
    </div>
</div>
        """
    
    def _html_company_comparison_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate company comparison slide HTML WITH ANALYSIS AND RECOMMENDATIONS"""
        companies = content.get('companies', [])
        companies_html = ''
        
        for company in companies[:2]:
            metrics = company.get('metrics', {})
            recommendation = company.get('recommendation', {})
            
            # Format metrics as readable text
            metrics_html = ''
            for key, value in metrics.items():
                if value and str(value).strip() and str(value) != 'N/A':
                    metrics_html += f"""
                    <div class="flex justify-between py-2 border-b border-gray-100">
                        <span class="text-gray-600 text-sm">{key}:</span>
                        <span class="font-semibold text-sm">{value}</span>
                    </div>
                    """
            
            # Add investment recommendation section
            recommendation_html = ""
            if recommendation:
                decision = recommendation.get('decision', 'WATCH')
                action = recommendation.get('action', 'Monitor')
                reasoning = recommendation.get('reasoning', 'Analysis in progress')
                score = recommendation.get('score', 'N/A')
                color = recommendation.get('color', 'gray')
                
                # Color mapping for recommendation badges
                color_classes = {
                    'green': 'bg-green-100 text-green-800 border-green-300',
                    'blue': 'bg-blue-100 text-blue-800 border-blue-300',
                    'yellow': 'bg-yellow-100 text-yellow-800 border-yellow-300',
                    'red': 'bg-red-100 text-red-800 border-red-300',
                    'gray': 'bg-gray-100 text-gray-800 border-gray-300'
                }
                badge_class = color_classes.get(color, color_classes['gray'])
                
                recommendation_html = f"""
                <div class="mt-4 pt-4 border-t-2 border-gray-300">
                    <div class="flex items-center justify-between mb-3">
                        <h4 class="text-sm font-bold text-gray-900 uppercase tracking-wide">Investment Recommendation</h4>
                        <span class="px-3 py-1 rounded-full text-xs font-bold border {badge_class}">{decision}</span>
                    </div>
                    <div class="space-y-2">
                        <p class="text-sm font-semibold text-gray-800">{action}</p>
                        <p class="text-xs text-gray-700 leading-relaxed">{reasoning}</p>
                        {f'<p class="text-xs text-gray-600">Score: {score}</p>' if score != 'N/A' else ''}
                    </div>
                </div>
                """
            
            companies_html += f"""
            <div class="flex-1 bg-white rounded-xl p-6 border-l-4 border-gray-800 shadow-sm">
                <h3 class="text-xl font-semibold text-gray-900 mb-2">{company.get('name', 'Company')}</h3>
                <div class="text-sm text-gray-600 mb-4">{company.get('business_model', 'N/A')}</div>
                <div class="space-y-1 mb-2">
                    {metrics_html}
                </div>
                {recommendation_html}
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-semibold text-gray-900 mb-2">{content.get('title', 'Company Overview & Financials')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', 'Pre/post money valuations and investment recommendations')}</p>
    <div class="flex gap-6">
        {companies_html}
    </div>
</div>
        """
    
    def _html_path_to_100m_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate Path to $100M ARR slide HTML"""
        companies = content.get('companies', {})
        metrics = content.get('metrics', {})
        
        # Create metrics tables
        metrics_html = ''
        for company_name, company_metrics in metrics.items():
            metrics_rows = ''
            for key, value in company_metrics.items():
                metrics_rows += f"""
                <tr>
                    <td class="px-3 py-2 text-sm text-gray-600">{key}</td>
                    <td class="px-3 py-2 text-sm font-semibold text-right">{value}</td>
                </tr>
                """
            
            metrics_html += f"""
            <div class="flex-1">
                <h4 class="font-bold text-lg mb-3">{company_name}</h4>
                <table class="w-full bg-gray-50 rounded-lg">
                    <tbody>{metrics_rows}</tbody>
                </table>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Path to $100M ARR')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    
    <div class="mb-8 bg-gray-100 rounded-xl p-8">
        <canvas id="chart-{slide_idx}" style="max-height: 300px;"></canvas>
    </div>
    
    <div class="flex gap-8">
        {metrics_html}
    </div>
</div>
        """
    
    def _html_business_analysis_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate business analysis comparison slide HTML"""
        companies = content.get('companies', {})
        companies_html = ''
        
        for company_name, data in companies.items():
            companies_html += f"""
            <div class="flex-1 bg-gray-50 rounded-xl p-6">
                <h3 class="text-xl font-bold text-gray-900 mb-4">{company_name}</h3>
                <div class="space-y-4">
                    <div>
                        <span class="text-sm text-gray-500 uppercase tracking-wider">What They Do</span>
                        <p class="mt-1 text-gray-800">{data.get('what_they_do', 'N/A')}</p>
                    </div>
                    <div>
                        <span class="text-sm text-gray-500 uppercase tracking-wider">What They Sell</span>
                        <p class="mt-1 text-gray-800">{data.get('what_they_sell', 'N/A')}</p>
                    </div>
                    <div>
                        <span class="text-sm text-gray-500 uppercase tracking-wider">Who They Sell To</span>
                        <p class="mt-1 text-gray-800">{data.get('who_they_sell_to', 'N/A')}</p>
                    </div>
                    <div class="pt-3 border-t border-gray-300">
                        <div class="grid grid-cols-2 gap-2 text-sm">
                            <div><span class="text-gray-600">Sector:</span> <span class="font-semibold">{data.get('sector', 'N/A')}</span></div>
                            <div><span class="text-gray-600">Founded:</span> <span class="font-semibold">{data.get('founded', 'N/A')}</span></div>
                            <div><span class="text-gray-600">Pricing:</span> <span class="font-semibold">{data.get('pricing_model', 'N/A')}</span></div>
                            <div><span class="text-gray-600">Team Size:</span> <span class="font-semibold">{data.get('team_size', 'N/A')}</span></div>
                        </div>
                    </div>
                </div>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-8">{content.get('title', 'Business Analysis Comparison')}</h2>
    <div class="flex gap-8">
        {companies_html if companies_html else '<p class="text-gray-600">No business data available</p>'}
    </div>
</div>
        """
    
    def _html_tam_pincer_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate TAM Pincer Analysis slide HTML"""
        companies = content.get('companies', {})
        tam_html = ''
        
        for company_name, data in companies.items():
            # Format TAM numbers with safety checks for None/0 values
            trad_tam_raw = data.get('traditional_tam', 0)
            labor_tam_raw = data.get('labor_tam', 0) 
            selected_tam_raw = data.get('selected_tam', 0)
            sam_raw = data.get('sam', 0)
            som_raw = data.get('som', 0)
            
            # Safe division with fallback to 0 if None
            trad_tam = (trad_tam_raw / 1e9) if trad_tam_raw and trad_tam_raw > 0 else 0
            labor_tam = (labor_tam_raw / 1e9) if labor_tam_raw and labor_tam_raw > 0 else 0
            selected_tam = (selected_tam_raw / 1e9) if selected_tam_raw and selected_tam_raw > 0 else 0
            sam = (sam_raw / 1e9) if sam_raw and sam_raw > 0 else 0
            som = (som_raw / 1e9) if som_raw and som_raw > 0 else 0
            penetration = data.get('penetration', 0) if data.get('penetration') is not None else 0
            
            tam_html += f"""
            <div class="flex-1 bg-gray-50 rounded-xl p-6">
                <h3 class="text-xl font-bold text-gray-900 mb-4">{company_name}</h3>
                <div class="space-y-4">
                    <div class="grid grid-cols-2 gap-4">
                        <div class="bg-blue-100 rounded-lg p-3">
                            <div class="text-sm text-blue-600 font-semibold">Software TAM</div>
                            <div class="text-2xl font-bold text-blue-900">${trad_tam:.1f}B</div>
                            <div class="text-xs text-blue-700 mt-1">{data.get('software_citation', 'Market research')}</div>
                        </div>
                        <div class="bg-gray-100 rounded-lg p-3 border-l-3 border-gray-700">
                            <div class="text-sm text-gray-600 font-semibold">Labor TAM</div>
                            <div class="text-2xl font-semibold text-gray-900">${labor_tam:.1f}B</div>
                            <div class="text-xs text-gray-600 mt-1">{data.get('labor_citation', 'Labor analysis')}</div>
                        </div>
                    </div>
                    <div class="border-t pt-4">
                        <div class="space-y-2">
                            <div class="flex justify-between">
                                <span class="text-gray-600">Selected TAM:</span>
                                <span class="font-bold text-lg">${selected_tam:.1f}B</span>
                            </div>
                            <div class="flex justify-between">
                                <span class="text-gray-600">SAM:</span>
                                <span class="font-semibold">${sam:.2f}B</span>
                            </div>
                            <div class="flex justify-between">
                                <span class="text-gray-600">SOM:</span>
                                <span class="font-semibold">${som:.2f}B</span>
                            </div>
                            <div class="flex justify-between">
                                <span class="text-gray-600">Current Penetration:</span>
                                <span class="font-semibold">{penetration:.2f}%</span>
                            </div>
                        </div>
                    </div>
                </div>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'TAM Pincer Analysis')}</h2>
    <p class="text-gray-600 mb-8">{content.get('subtitle', 'Market Opportunity & Penetration')}</p>
    <div class="flex gap-8">
        {tam_html if tam_html else '<p class="text-gray-600">No TAM data available</p>'}
    </div>
</div>
        """
    
    def _html_founder_team_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate founder & team analysis slide HTML"""
        companies = content.get('companies', {})
        team_html = ''
        
        for company_name, data in companies.items():
            founders = data.get('founders', [])
            founders_html = ''
            
            for founder in founders[:2]:  # Show max 2 founders
                founders_html += f"""
                <div class="mb-3">
                    <div class="font-semibold">{founder.get('name', 'Unknown')}</div>
                    <div class="text-sm text-gray-600">{founder.get('role', '')}</div>
                    <div class="text-sm text-gray-700 mt-1">{founder.get('background', 'N/A')[:100]}...</div>
                </div>
                """
            
            quality_html = ''.join([f'<li class="text-sm text-green-700">• {signal}</li>' for signal in data.get('quality_signals', [])])
            risk_html = ''.join([f'<li class="text-sm text-red-700">• {risk}</li>' for risk in data.get('risk_factors', [])])
            
            team_html += f"""
            <div class="flex-1 bg-gray-50 rounded-xl p-6">
                <h3 class="text-xl font-bold text-gray-900 mb-4">{company_name}</h3>
                <div class="mb-4">
                    <h4 class="font-semibold text-gray-700 mb-2">Founders</h4>
                    {founders_html if founders_html else '<p class="text-sm text-gray-600">Founder information not available</p>'}
                </div>
                <div class="mb-3">
                    <span class="text-sm text-gray-600">Team Size:</span>
                    <span class="font-semibold ml-2">{data.get('team_size', 'Unknown')} employees</span>
                </div>
                {f'<div class="mb-3"><h4 class="font-semibold text-gray-700 mb-1">Quality Signals</h4><ul>{quality_html}</ul></div>' if quality_html else ''}
                {f'<div><h4 class="font-semibold text-gray-700 mb-1">Risk Factors</h4><ul>{risk_html}</ul></div>' if risk_html else ''}
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Founder & Team Analysis')}</h2>
    <p class="text-gray-600 mb-8">{content.get('subtitle', 'Leadership profiles and quality assessment')}</p>
    <div class="flex gap-8">
        {team_html if team_html else '<p class="text-gray-600">No team data available</p>'}
    </div>
</div>
        """
    
    def _html_cap_table_comparison_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate cap table comparison slide HTML"""
        devices = content.get('devices', [])
        
        # Check if this has devices with side-by-side Sankey
        if devices:
            for device in devices:
                if device.get('type') == 'side_by_side_sankey':
                    data = device.get('data', {})
                    company1_name = data.get('company1_name', 'Company 1')
                    company2_name = data.get('company2_name', 'Company 2')
                    
                    # For now, create a simple comparison visualization
                    # We'll use simple bars to show ownership changes
                    # Pre-compute to avoid backslash in f-string (Python 3.11)
                    _insights_items = ''.join([
                        '<p class="text-sm text-gray-700">' + '\u2022 ' + insight + '</p>'
                        for insight in content.get("insights", [])
                    ])
                    _insights_html = f'<div class="mt-6 grid grid-cols-2 gap-4">{_insights_items}</div>' if content.get('insights') else ''
                    return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-4">{content.get('title', 'Cap Table Evolution')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    <div class="grid grid-cols-2 gap-8">
        <div>
            <h3 class="text-xl font-semibold mb-4 text-center">{company1_name}</h3>
            <div class="chart-container" style="height: 350px;">
                <canvas id="chart-{slide_idx}-0"></canvas>
            </div>
        </div>
        <div>
            <h3 class="text-xl font-semibold mb-4 text-center">{company2_name}</h3>
            <div class="chart-container" style="height: 350px;">
                <canvas id="chart-{slide_idx}-1"></canvas>
            </div>
        </div>
    </div>
    {_insights_html}
</div>
                    """
        
        # Fall back to original implementation
        return self._html_cap_table_slide(content, slide_idx)
    
    def _create_ownership_chart_from_sankey(self, sankey_data: Dict[str, Any], company_name: str) -> Dict[str, Any]:
        """Convert Sankey data to a stacked bar chart for PDF export"""
        nodes = sankey_data.get('nodes', [])
        links = sankey_data.get('links', [])
        
        # Extract ownership percentages from nodes
        rounds = []
        ownership_data = {}
        
        # Find unique rounds and stakeholders
        round_set = set()
        stakeholder_set = set()
        
        for node in nodes:
            if 'round' in node:
                round_set.add(node['round'])
            name = node.get('name', '')
            if name and not any(x in name.lower() for x in ['round', 'total', 'exit']):
                stakeholder_set.add(name)
        
        # Sort rounds chronologically
        round_order = ['Seed', 'Series A', 'Series B', 'Series C', 'Series D', 'Series E', 'Exit']
        rounds = sorted(list(round_set), key=lambda x: round_order.index(x) if x in round_order else 999)
        
        if not rounds:
            rounds = ['Current']
        
        # Initialize ownership data for each stakeholder
        for stakeholder in stakeholder_set:
            ownership_data[stakeholder] = []
        
        # Process each round
        for round_name in rounds:
            round_values = {}
            for node in nodes:
                if node.get('round') == round_name:
                    name = node.get('name', '')
                    value = node.get('value', 0) or node.get('ownership', 0)
                    if name and not any(x in name.lower() for x in ['round', 'total', 'exit']):
                        round_values[name] = value
            
            # Add values for this round
            for stakeholder in stakeholder_set:
                ownership_data[stakeholder].append(round_values.get(stakeholder, 0))
        
        # If no valid data, create simple fallback
        if not ownership_data or all(sum(v) == 0 for v in ownership_data.values()):
            rounds = ['Seed', 'Series A', 'Current']
            ownership_data = {
                'Founders': [80, 60, 40],
                'Investors': [15, 30, 45],
                'ESOP': [5, 10, 15]
            }
        
        # Create datasets for stacked bar chart
        datasets = []
        colors = [
            'rgba(17, 24, 39, 0.9)',    # #111827 - Darkest gray
            'rgba(31, 41, 55, 0.85)',   # #1F2937 - Dark gray
            'rgba(55, 65, 81, 0.8)',    # #374151 - Medium gray
            'rgba(75, 85, 99, 0.75)',   # #4B5563 - Gray
            'rgba(107, 114, 128, 0.7)', # #6B7280 - Light gray
            'rgba(156, 163, 175, 0.65)',# #9CA3AF - Lighter gray
        ]
        
        for idx, (stakeholder, values) in enumerate(ownership_data.items()):
            datasets.append({
                'label': stakeholder[:20],  # Limit label length
                'data': values,
                'backgroundColor': colors[idx % len(colors)],
                'borderColor': colors[idx % len(colors)].replace('0.8', '1'),
                'borderWidth': 1
            })
        
        # Fix case-insensitive key lookup
        founders_key = next((k for k in ownership_data.keys() if k.lower() == 'founders'), None)
        investors_key = next((k for k in ownership_data.keys() if k.lower() == 'investors'), None)
        employees_key = next((k for k in ownership_data.keys() if k.lower() in ['employees', 'esop']), None)
        
        chart_datasets = []
        if founders_key:
            chart_datasets.append({
                'label': 'Founders',
                'data': ownership_data[founders_key],
                'backgroundColor': 'rgba(17, 24, 39, 0.9)',
                'borderColor': 'rgba(17, 24, 39, 1)',
                'borderWidth': 1
            })
        if investors_key:
            chart_datasets.append({
                'label': 'Investors',
                'data': ownership_data[investors_key],
                'backgroundColor': 'rgba(55, 65, 81, 0.8)',
                'borderColor': 'rgba(55, 65, 81, 1)',
                'borderWidth': 1
            })
        if employees_key:
            chart_datasets.append({
                'label': 'Employees/ESOP',
                'data': ownership_data[employees_key],
                'backgroundColor': 'rgba(107, 114, 128, 0.7)',
                'borderColor': 'rgba(107, 114, 128, 1)',
                'borderWidth': 1
            })
        
        # If no valid datasets, use the fallback data
        if not chart_datasets:
            chart_datasets = datasets  # Use the dynamically created datasets from above
        
        return {
            'type': 'bar',
            'data': {
                'labels': rounds,
                'datasets': chart_datasets
            },
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'animation': False,  # Disable animations for PDF rendering
                'plugins': {
                    'legend': {
                        'position': 'bottom',
                        'labels': {
                            'padding': 10,
                            'font': {
                                'size': 11
                            }
                        }
                    },
                    'title': {
                        'display': False
                    }
                },
                'scales': {
                    'x': {
                        'stacked': True,
                        'grid': {
                            'display': False
                        }
                    },
                    'y': {
                        'stacked': True,
                        'beginAtZero': True,
                        'max': 100,
                        'grid': {
                            'color': 'rgba(0, 0, 0, 0.05)'
                        }
                    }
                }
            }
        }
    
    def _create_waterfall_config(self, chart_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create Chart.js configuration for waterfall charts"""
        data = chart_data.get('data', {})
        labels = data.get('labels', [])
        waterfall_values = data.get('datasets', [{}])[0].get('data', [])
        
        # Determine if this is ownership (starts at 100) or valuation (starts at initial value)
        # Check first value to determine starting point
        initial_value = waterfall_values[0] if waterfall_values else 100
        is_valuation_chart = initial_value < 100 and initial_value > 0
        
        # Calculate cumulative values for waterfall effect
        if is_valuation_chart:
            # For valuation charts, start at initial value and ensure it never goes below
            cumulative = initial_value
            min_value = initial_value  # Track minimum to prevent going below initial
        else:
            # For ownership charts, start at 100%
            cumulative = 100
            min_value = 0
        
        processed_data = []
        colors = []
        
        for i, value in enumerate(waterfall_values):
            if i == 0:
                # First bar - initial value
                processed_data.append([0, cumulative])
                colors.append('rgba(17, 24, 39, 0.9)')  # Dark gray for initial
            elif i == len(waterfall_values) - 1:
                # Last bar - final value (ensure it's not below initial for valuation charts)
                if is_valuation_chart:
                    cumulative = max(cumulative, min_value)
                processed_data.append([0, cumulative])
                colors.append('rgba(17, 24, 39, 0.9)')  # Dark gray for final
            else:
                # Middle bars - changes (can be positive or negative)
                start = cumulative
                cumulative += value
                # For valuation charts, ensure cumulative never goes below initial
                if is_valuation_chart:
                    cumulative = max(cumulative, min_value)
                processed_data.append([cumulative, start])
                colors.append('rgba(75, 85, 99, 0.75)' if value < 0 else 'rgba(31, 41, 55, 0.85)')
        
        return {
            'type': 'bar',
            'data': {
                'labels': labels,
                'datasets': [{
                    'label': 'Ownership %',
                    'data': processed_data,
                    'backgroundColor': colors,
                    'borderColor': [c.replace('0.6', '1').replace('0.8', '1') for c in colors],
                    'borderWidth': 1
                }]
            },
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'animation': False,  # Disable animations for PDF rendering
                'indexAxis': 'x',
                'plugins': {
                    'legend': {
                        'display': False
                    },
                    'title': {
                        'display': bool(chart_data.get('title')),
                        'text': chart_data.get('title', ''),
                        'font': {
                            'size': 16,
                            'weight': 'bold'
                        }
                    }
                },
                'scales': {
                    'x': {
                        'grid': {
                            'display': False
                        },
                        'ticks': {
                            'autoSkip': True,
                            'maxRotation': 45,
                            'minRotation': 0,
                            'maxTicksLimit': 10,
                            'font': {
                                'size': 10
                            }
                        }
                    },
                    'y': {
                        'beginAtZero': True,
                        'max': 100,
                        'grid': {
                            'color': 'rgba(0, 0, 0, 0.05)'
                        }
                    }
                }
            }
        }

    def _html_citations_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate citations slide HTML"""
        citations = content.get('citations', [])
        citations_html = ''
        
        for i, citation in enumerate(citations[:15], 1):  # Limit to 15 citations
            citations_html += f"""
            <div class="flex gap-2 mb-2">
                <span class="text-gray-500">{i}.</span>
                <span class="text-sm text-gray-700">{citation}</span>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-8">{content.get('title', 'Sources & References')}</h2>
    <div class="columns-2 gap-8">
        {citations_html if citations_html else '<p class="text-gray-600">No citations available</p>'}
    </div>
</div>
        """
    
    def _add_citations_to_slide(self, slide_html: str, citations: List[Dict[str, Any]]) -> str:
        """Add citations section to any slide HTML"""
        if not citations or len(citations) == 0:
            return slide_html
        
        citations_html = ""
        for i, citation in enumerate(citations, 1):
            url = citation.get('url', citation.get('source', ''))
            title = citation.get('title', citation.get('text', citation.get('source', f'Source {i}')))
            
            if url and url.startswith('http'):
                citation_link = f'<a href="{url}" target="_blank" rel="noopener noreferrer" class="text-blue-600 hover:text-blue-800 underline">{title}</a>'
            else:
                citation_link = f'<span>{title}</span>'
            
            citations_html += f"""
            <div class="flex gap-2 mb-1">
                <span class="text-xs text-gray-500 font-medium">[{i}]</span>
                <span class="text-xs text-gray-600">{citation_link}</span>
            </div>
            """
        
        citations_section = f"""
        <div class="mt-6 pt-4 border-t border-gray-200">
            <h4 class="text-sm font-semibold text-gray-700 mb-3">Sources</h4>
            <div class="space-y-1">
                {citations_html}
            </div>
        </div>
        """
        
        # Insert citations before the closing </div> of the slide
        if slide_html.rstrip().endswith('</div>'):
            slide_html = slide_html.rstrip()[:-6] + citations_section + '\n</div>'
        else:
            slide_html += citations_section
        
        return slide_html
    
    def _html_cap_table_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate cap table slide HTML with waterfall chart"""
        chart_data = content.get('chart_data')
        future_chart_data = content.get('future_chart_data')
        future_pie_charts = content.get('future_pie_charts', [])
        
        # Build main chart container
        chart_container = ""
        if chart_data:
            chart_type = chart_data.get('type', 'bar').lower()
            
            # For D3.js-based charts (sankey, side_by_side_sankey, heatmap, sunburst), use SVG container
            # For other charts, use canvas container
            if chart_type in ['sankey', 'side_by_side_sankey', 'heatmap', 'sunburst']:
                chart_container = f"""
    <div class="chart-container" style="height: 400px;" data-chart-type="{chart_type}" id="chart-{slide_idx}">
        <!-- D3.js will render SVG here -->
    </div>
                """
            else:
                chart_container = f"""
    <div class="chart-container" style="height: 400px;" data-chart-type="{chart_type}">
        <canvas id="chart-{slide_idx}"></canvas>
    </div>
                """
        
        # Build future chart container
        future_chart_container = ""
        if future_chart_data:
            future_chart_type = future_chart_data.get('type', 'bar').lower()
            if future_chart_type == 'sankey':
                future_chart_container = f"""
    <div class="mt-8 pt-8 border-t border-gray-300">
        <h3 class="text-xl font-semibold mb-4 text-gray-900">Future Ownership Scenarios</h3>
        <div class="chart-container" style="height: 400px;" data-chart-type="{future_chart_type}" id="future-chart-{slide_idx}">
            <!-- D3.js will render SVG here -->
        </div>
    </div>
                """
            else:
                future_chart_container = f"""
    <div class="mt-8 pt-8 border-t border-gray-300">
        <h3 class="text-xl font-semibold mb-4 text-gray-900">Future Ownership Scenarios</h3>
        <div class="chart-container" style="height: 400px;" data-chart-type="{future_chart_type}">
            <canvas id="future-chart-{slide_idx}"></canvas>
        </div>
    </div>
                """
        
        # Build future pie charts container
        future_pie_charts_html = ""
        if future_pie_charts and len(future_pie_charts) > 0:
            future_pie_charts_html = '<div class="mt-8 pt-8 border-t border-gray-300"><h3 class="text-xl font-semibold mb-4 text-gray-900">Future Ownership Scenarios</h3><div class="grid grid-cols-2 gap-6">'
            for idx, pie_chart in enumerate(future_pie_charts):
                pie_title = pie_chart.get('title', f'Scenario {idx + 1}')
                future_pie_charts_html += f"""
    <div class="p-4 border border-gray-300 rounded-lg">
        <h4 class="text-sm font-semibold mb-2 text-gray-700">{pie_title}</h4>
        <div class="chart-container" style="height: 300px;" data-chart-type="pie">
            <canvas id="future-pie-{slide_idx}-{idx}"></canvas>
        </div>
    </div>
                """
            future_pie_charts_html += '</div></div>'
        
        # Pre-compute to avoid backslash in f-string (Python 3.11)
        _bullets_items = ''.join([
            '<p class="text-sm text-gray-700">' + '\u2022 ' + bullet + '</p>'
            for bullet in content.get("bullets", [])
        ])
        _bullets_html = f'<div class="mt-6 space-y-2">{_bullets_items}</div>' if content.get('bullets') else ''

        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-4">{content.get('title', 'Cap Table Evolution')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    {chart_container}
    {future_chart_container}
    {future_pie_charts_html}
    {_bullets_html}
</div>
            """
        
        # Otherwise handle as a comparison table (legacy format)
        companies = content.get('companies', [])
        companies_html = ''
        
        for company in companies[:2]:  # Show 2 companies side by side
            cap_table = company.get('cap_table', {})
            rounds = cap_table.get('rounds', [])
            
            rounds_html = ''
            for round_data in rounds:
                rounds_html += f"""
                <tr>
                    <td class="px-4 py-2 text-sm">{round_data.get('round', '')}</td>
                    <td class="px-4 py-2 text-sm text-right">{round_data.get('founders', 0):.1f}%</td>
                    <td class="px-4 py-2 text-sm text-right">{round_data.get('investors', 0):.1f}%</td>
                    <td class="px-4 py-2 text-sm text-right">{round_data.get('employees', 0):.1f}%</td>
                </tr>
                """
            
            companies_html += f"""
            <div class="flex-1">
                <h3 class="text-xl font-bold text-gray-900 mb-4 text-center">{company.get('name', 'Company')}</h3>
                <div class="bg-gray-50 rounded-lg p-4">
                    <table class="w-full">
                        <thead>
                            <tr class="border-b border-gray-200">
                                <th class="text-left px-4 py-2 text-sm font-semibold">Round</th>
                                <th class="text-right px-4 py-2 text-sm font-semibold">Founders</th>
                                <th class="text-right px-4 py-2 text-sm font-semibold">Investors</th>
                                <th class="text-right px-4 py-2 text-sm font-semibold">Employees</th>
                            </tr>
                        </thead>
                        <tbody>
                            {rounds_html}
                        </tbody>
                    </table>
                </div>
                <div class="mt-4 space-y-2">
                    {''.join([f'<p class="text-sm text-gray-600">• {h}</p>' for h in company.get('highlights', [])[:3]])}
                </div>
            </div>
            """
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-8 text-center">{content.get('title', 'Cap Table Comparison')}</h2>
    <div class="flex gap-8">
        {companies_html}
    </div>
</div>
        """
    
    def _html_sankey_slide(self, content: Dict[str, Any], slide_idx: int, chart_images: Dict[str, str] = None) -> str:
        """Generate Sankey diagram slide HTML"""
        chart_data = content.get('chart_data', {})
        
        # Check if we have a pre-rendered Sankey image
        chart_key = f"slide_{slide_idx}_chart_0"
        chart_html = ""
        
        if chart_images and chart_key in chart_images:
            # Use pre-rendered high-quality PNG
            img_base64 = chart_images[chart_key]
            chart_html = f'<img src="data:image/png;base64,{img_base64}" style="width: 100%; height: 500px; object-fit: contain;" />'
            logger.info(f"[HTML_SANKEY] Using pre-rendered image for slide {slide_idx}")
        else:
            # Fall back to D3 container
            chart_html = f'<div id="sankey-{slide_idx}" style="height: 500px; width: 100%;"></div>'
            logger.info(f"[HTML_SANKEY] Using D3 container for slide {slide_idx}")
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-6">{content.get('title', 'Ownership Flow')}</h2>
    {chart_html}
</div>
        """
    
    def _html_side_by_side_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate side-by-side comparison slide HTML"""
        left_content = content.get('left', {})
        right_content = content.get('right', {})
        
        def generate_side_html(side_content, side_idx):
            html = f'<div class="flex-1">'
            
            # Subtitle
            if side_content.get('subtitle'):
                html += f'<h3 class="text-xl font-bold text-gray-900 mb-4 text-center">{side_content.get("subtitle")}</h3>'
            
            # Chart if present
            if side_content.get('chart_data'):
                chart_type = side_content.get('chart_data', {}).get('type', 'bar').lower()
                html += f'<div class="chart-container mb-4" data-chart-type="{chart_type}"><canvas id="chart-{slide_idx}-{side_idx}"></canvas></div>'
            
            # Metrics
            if side_content.get('metrics'):
                html += '<div class="bg-gray-50 rounded-lg p-4 mb-4">'
                for key, value in list(side_content.get('metrics', {}).items())[:5]:
                    html += f'<div class="flex justify-between py-1"><span class="text-sm text-gray-600">{key}:</span><span class="text-sm font-semibold">{value}</span></div>'
                html += '</div>'
            
            # Bullets
            if side_content.get('bullets'):
                html += '<ul class="list-disc list-inside text-sm space-y-1">'
                for bullet in side_content.get('bullets', [])[:4]:
                    html += f'<li class="text-gray-700">{bullet}</li>'
                html += '</ul>'
            
            html += '</div>'
            return html
        
        left_html = generate_side_html(left_content, 0) if left_content else '<div class="flex-1"></div>'
        right_html = generate_side_html(right_content, 1) if right_content else '<div class="flex-1"></div>'
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-8 text-center">{content.get('title', 'Side-by-Side Comparison')}</h2>
    <div class="flex gap-8">
        {left_html}
        {right_html}
    </div>
</div>
        """
    
    def _html_exit_scenarios_comprehensive_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate comprehensive exit scenarios slide HTML with charts"""
        charts_html = ""
        
        # Render charts if available
        if content.get("charts"):
            charts_html = '<div class="grid grid-cols-1 md:grid-cols-2 gap-4 mb-6">'
            for chart_idx, chart in enumerate(content["charts"]):
                chart_id = f"chart-{slide_idx}-{chart_idx}"
                charts_html += f"""
                <div class="bg-white border border-gray-200 rounded-lg p-4">
                    <h4 class="text-sm font-semibold text-gray-700 mb-3">{chart.get('title', '')}</h4>
                    <canvas id="{chart_id}" style="max-height: 300px;"></canvas>
                </div>
                """
            charts_html += '</div>'
        
        # Render reality check table if available
        reality_table_html = ""
        if content.get("reality_check"):
            reality_table_html = """
            <div class="bg-red-50 border-2 border-red-200 rounded-lg p-4 mb-6">
                <h3 class="text-lg font-bold text-red-900 mb-3">🚨 The $150M Problem - Reality Check</h3>
                <table class="w-full text-sm">
                    <thead>
                        <tr class="border-b border-red-200">
                            <th class="text-left py-2">Company</th>
                            <th class="text-left">Exit Value</th>
                            <th class="text-left">Our Return</th>
                            <th class="text-left">MOIC</th>
                            <th class="text-left">DPI Impact</th>
                        </tr>
                    </thead>
                    <tbody>
            """
            for row in content["reality_check"]:
                reality_table_html += f"""
                    <tr class="border-b border-red-100">
                        <td class="py-2 font-semibold">{row['company']}</td>
                        <td>{row['exit_value']}</td>
                        <td class="{'text-red-600' if 'Lost' in row.get('reality', '') else ''}">{row['our_proceeds']}</td>
                        <td class="{'text-red-600 font-bold' if float(row['moic'].replace('x','')) < 1 else ''}">{row['moic']}</td>
                        <td>{row['dpi_contribution']}</td>
                    </tr>
                """
            reality_table_html += """
                    </tbody>
                </table>
            </div>
            """
        
        # Render companies data
        companies_html = ""
        companies = content.get("companies", {})
        for company_name, company_data in companies.items():
            breakpoints_html = ""
            if company_data.get("breakpoints"):
                bp = company_data["breakpoints"]
                
                # Helper to format breakpoint values
                def fmt_bp(key, default=0):
                    val = bp.get(key, default)
                    if val is None or val == 0:
                        return None
                    return f"${safe_divide(ensure_numeric(val), 1e6, 0):.0f}M"
                
                # Build breakpoints display with both scenarios
                bp_items = []
                
                # Without pro rata
                if fmt_bp("exit_for_1x_return_no_pro_rata"):
                    bp_items.append(f'<div>• 1x (no pro rata): {fmt_bp("exit_for_1x_return_no_pro_rata")}</div>')
                if fmt_bp("exit_for_2x_return_no_pro_rata"):
                    bp_items.append(f'<div>• 2x (no pro rata): {fmt_bp("exit_for_2x_return_no_pro_rata")}</div>')
                if fmt_bp("exit_for_3x_return_no_pro_rata"):
                    bp_items.append(f'<div>• 3x (no pro rata): {fmt_bp("exit_for_3x_return_no_pro_rata")}</div>')
                
                # With pro rata
                if fmt_bp("exit_for_1x_return_with_pro_rata"):
                    bp_items.append(f'<div>• 1x (w/ pro rata): <span class="text-blue-700 font-semibold">{fmt_bp("exit_for_1x_return_with_pro_rata")}</span></div>')
                if fmt_bp("exit_for_2x_return_with_pro_rata"):
                    bp_items.append(f'<div>• 2x (w/ pro rata): <span class="text-blue-700 font-semibold">{fmt_bp("exit_for_2x_return_with_pro_rata")}</span></div>')
                if fmt_bp("exit_for_3x_return_with_pro_rata"):
                    bp_items.append(f'<div>• 3x (w/ pro rata): <span class="text-blue-700 font-semibold">{fmt_bp("exit_for_3x_return_with_pro_rata")}</span></div>')
                
                # Fallback to legacy fields if new ones don't exist
                if not bp_items:
                    if fmt_bp("target_3x_exit"):
                        bp_items.append(f'<div>• 3x Target: {fmt_bp("target_3x_exit")}</div>')
                    if fmt_bp("target_2x_exit"):
                        bp_items.append(f'<div>• 2x Target: {fmt_bp("target_2x_exit")}</div>')
                
                breakpoints_html = f"""
                <div class="bg-white rounded p-3 mb-3">
                    <h5 class="text-xs font-semibold text-gray-600 mb-2">Key Breakpoints (Exit Values in $M)</h5>
                    <div class="space-y-1 text-sm">
                        {''.join(bp_items)}
                    </div>
                </div>
                """ if bp_items else ""
            
            entry_exit_html = f"""
            <div class="grid grid-cols-2 gap-3 mb-3">
                <div class="bg-white rounded p-2">
                    <h5 class="text-xs font-semibold text-gray-600 mb-1">Entry</h5>
                    <div class="text-xs">
                        <div>${safe_divide(ensure_numeric(company_data.get('entry_economics', {}).get('investment', 0)), 1e6, 0):.1f}M</div>
                        <div>{company_data.get('entry_economics', {}).get('entry_ownership', 0):.1f}%</div>
                    </div>
                </div>
                <div class="bg-white rounded p-2">
                    <h5 class="text-xs font-semibold text-gray-600 mb-1">Exit</h5>
                    <div class="text-xs">
                        <div>No F: {company_data.get('exit_economics', {}).get('ownership_no_followon', 0):.1f}%</div>
                        <div>W/ F: {company_data.get('exit_economics', {}).get('ownership_with_followon', 0):.1f}%</div>
                    </div>
                </div>
            </div>
            """
            
            companies_html += f"""
            <div class="bg-gray-50 rounded-lg p-4 border-l-4 border-gray-800">
                <h3 class="font-semibold text-base text-gray-900 mb-3">{company_name}</h3>
                {entry_exit_html}
                {breakpoints_html}
            </div>
            """
        
        # Render insights if available
        insights_html = ""
        if content.get("insights"):
            insights_html = '<div class="bg-yellow-50 border border-yellow-200 rounded-lg p-4 mb-6"><ul class="space-y-2">'
            for insight in content["insights"]:
                insights_html += f'<li class="text-sm">{insight}</li>'
            insights_html += '</ul></div>'
        
        return f"""
<div class="slide bg-white p-12">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Exit Scenarios & Ownership Analysis')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', 'PWERM Analysis with Liquidation Breakpoints')}</p>
    
    {reality_table_html}
    
    {charts_html}
    
    {insights_html}
    
    <div class="grid grid-cols-2 gap-6">
        {companies_html}
    </div>
</div>
        """
    
    def _serialize_chart_config(self, config: Dict[str, Any]) -> str:
        """Serialize chart config with proper handling of JavaScript functions"""
        import json
        
        # Convert to JSON string
        json_str = json.dumps(config)
        
        # Replace function strings with actual JavaScript functions
        # Look for patterns like "function(...)" 
        import re
        
        # Pattern to match function strings
        pattern = r'"function\([^)]*\)[^"]*"'
        
        def replace_func(match):
            # Remove quotes around function
            func_str = match.group(0)[1:-1]  # Remove surrounding quotes
            return func_str
        
        json_str = re.sub(pattern, replace_func, json_str)
        
        return json_str
    
    def _generate_chart_scripts(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]]) -> str:
        """Generate Chart.js initialization scripts"""
        scripts = []
        
        # Handle both dict with slides key and direct list of slides
        if isinstance(deck_data, list):
            slides = deck_data
        else:
            slides = deck_data.get("slides", deck_data.get("deck_slides", []))
        
        for slide_idx, slide_data in enumerate(slides):
            content = slide_data.get("content", {})
            slide_type = slide_data.get("type", "")
            
            # Handle comprehensive exit scenarios with multiple charts
            if slide_type == "exit_scenarios_comprehensive" and content.get("charts"):
                for chart_idx, chart_data in enumerate(content["charts"]):
                    # Skip Sankey charts (they require D3, not Chart.js)
                    if chart_data.get('type') == 'sankey':
                        continue
                    chart_config = self._create_chart_config(chart_data)
                    scripts.append(f"""
                        var ctx{slide_idx}_{chart_idx} = document.getElementById('chart-{slide_idx}-{chart_idx}');
                        if (ctx{slide_idx}_{chart_idx}) {{
                            new Chart(ctx{slide_idx}_{chart_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                        }}
                    """)
            
            # Handle individual cap table slides specifically
            if slide_type == "cap_table":
                # Main chart data
                if content.get("chart_data"):
                    chart_data = content.get("chart_data")
                    chart_config = self._create_chart_config(chart_data)
                    scripts.append(f"""
                        var ctx{slide_idx} = document.getElementById('chart-{slide_idx}');
                        if (ctx{slide_idx}) {{
                            new Chart(ctx{slide_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                        }}
                    """)
                
                # Future chart data (waterfall/bar chart)
                if content.get("future_chart_data"):
                    future_chart_data = content.get("future_chart_data")
                    future_chart_config = self._create_chart_config(future_chart_data)
                    scripts.append(f"""
                        var futureCtx{slide_idx} = document.getElementById('future-chart-{slide_idx}');
                        if (futureCtx{slide_idx}) {{
                            new Chart(futureCtx{slide_idx}.getContext('2d'), {self._serialize_chart_config(future_chart_config)});
                        }}
                    """)
                
                # Future pie charts array
                future_pie_charts = content.get("future_pie_charts", [])
                for pie_idx, pie_chart in enumerate(future_pie_charts):
                    pie_data = pie_chart.get('data') or pie_chart
                    pie_chart_config = self._create_chart_config({
                        'type': 'pie',
                        'data': pie_data,
                        'title': pie_chart.get('title', f'Scenario {pie_idx + 1}')
                    })
                    scripts.append(f"""
                        var futurePieCtx{slide_idx}_{pie_idx} = document.getElementById('future-pie-{slide_idx}-{pie_idx}');
                        if (futurePieCtx{slide_idx}_{pie_idx}) {{
                            new Chart(futurePieCtx{slide_idx}_{pie_idx}.getContext('2d'), {self._serialize_chart_config(pie_chart_config)});
                        }}
                    """)
            
            # Handle fund impact slides
            elif slide_type == "fund_return_impact_enhanced" and content.get("chart_data"):
                chart_data = content.get("chart_data")
                chart_config = self._create_chart_config(chart_data)
                scripts.append(f"""
                    var fundImpactCtx{slide_idx} = document.getElementById('fund-impact-chart-{slide_idx}');
                    if (fundImpactCtx{slide_idx}) {{
                        new Chart(fundImpactCtx{slide_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                    }}
                """)
            
            # Handle path_to_100m_comparison slides
            elif slide_type == "path_to_100m_comparison" and content.get("chart_data"):
                chart_data = content.get("chart_data")
                chart_config = self._create_chart_config(chart_data)
                scripts.append(f"""
                    var ctx{slide_idx} = document.getElementById('chart-{slide_idx}');
                    if (ctx{slide_idx}) {{
                        new Chart(ctx{slide_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                    }}
                """)
            
            # Handle probability_cloud slides
            elif slide_type == "probability_cloud" and content.get("chart_data"):
                chart_data = content.get("chart_data")
                # Probability cloud uses special config
                chart_config = self._create_probability_cloud_config(chart_data)
                chart_id = f"prob-cloud-{slide_idx}"
                scripts.append(f"""
                    var ctx{slide_idx} = document.getElementById('{chart_id}');
                    if (ctx{slide_idx}) {{
                        new Chart(ctx{slide_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                    }}
                """)
            
            # Handle regular charts
            elif content.get("chart_data"):
                chart_data = content.get("chart_data")
                chart_type = chart_data.get('type', '').lower()
                
                # Handle D3.js-based charts (sankey, side_by_side_sankey, heatmap)
                if chart_type in ['sankey', 'side_by_side_sankey', 'heatmap']:
                    sankey_data = chart_data.get('data', {})
                    # Use D3.js for Sankey rendering
                    scripts.append(f"""
                        (function() {{
                            const container = document.getElementById('chart-{slide_idx}');
                            if (!container || typeof d3 === 'undefined') {{
                                console.warn('D3.js not available, falling back to bar chart');
                                // Fallback to bar chart if D3 not available
                                var ctx = container;
                                if (ctx) {{
                                    const fallbackChart = {self._serialize_chart_config(
                                        self._create_ownership_chart_from_sankey(
                                            sankey_data,
                                            content.get('title', 'Ownership Evolution')
                                        )
                                    )};
                                    new Chart(ctx.getContext('2d'), fallbackChart);
                                }}
                                return;
                            }}
                            
                            // Create SVG directly in the container
                            container.style.display = 'block';
                            const width = container.offsetWidth || 800;
                            const height = container.offsetHeight || 400;
                            
                            // Create SVG inside the container
                            const svg = d3.select(container)
                                .append('svg')
                                .attr('width', width)
                                .attr('height', height);
                            
                            // Render Sankey with D3
                            const data = {json.dumps(sankey_data)};
                            
                            const sankey = d3.sankey()
                                .nodeId(d => d.id)
                                .nodeAlign(d3.sankeyJustify)
                                .nodeWidth(15)
                                .nodePadding(10)
                                .extent([[10, 10], [width - 10, height - 10]]);
                            
                            const {{nodes, links}} = sankey(data);
                            
                            // Add links
                            svg.append('g')
                                .selectAll('path')
                                .data(links)
                                .join('path')
                                .attr('d', d3.sankeyLinkHorizontal())
                                .attr('stroke', d => d.color || '#999')
                                .attr('stroke-width', d => d.width)
                                .attr('fill', 'none')
                                .attr('opacity', 0.5);
                            
                            // Add nodes
                            svg.append('g')
                                .selectAll('rect')
                                .data(nodes)
                                .join('rect')
                                .attr('x', d => d.x0)
                                .attr('y', d => d.y0)
                                .attr('height', d => d.y1 - d.y0)
                                .attr('width', d => d.x1 - d.x0)
                                .attr('fill', d => d.color || '#69b3a2');
                            
                            // Add labels
                            svg.append('g')
                                .selectAll('text')
                                .data(nodes)
                                .join('text')
                                .attr('x', d => d.x0 < width / 2 ? d.x1 + 6 : d.x0 - 6)
                                .attr('y', d => (d.y1 + d.y0) / 2)
                                .attr('dy', '0.35em')
                                .attr('text-anchor', d => d.x0 < width / 2 ? 'start' : 'end')
                                .text(d => d.name)
                                .style('font-size', '11px');
                        }})();
                    """)
                else:
                    chart_config = self._create_chart_config(chart_data)
                    scripts.append(f"""
                        var ctx{slide_idx} = document.getElementById('chart-{slide_idx}');
                        if (ctx{slide_idx}) {{
                            new Chart(ctx{slide_idx}.getContext('2d'), {self._serialize_chart_config(chart_config)});
                        }}
                    """)
            
            # Handle cap table comparison with side-by-side Sankey
            if slide_type == 'cap_table_comparison':
                devices = content.get('devices', [])
                for device in devices:
                    if device.get('type') == 'side_by_side_sankey':
                        data = device.get('data', {})
                        # Create simple bar charts for each company showing ownership
                        for idx, (company_key, company_name) in enumerate([
                            ('company1_data', data.get('company1_name', 'Company 1')),
                            ('company2_data', data.get('company2_name', 'Company 2'))
                        ]):
                            sankey_data = data.get(company_key, {})
                            # Convert Sankey data to simple ownership chart
                            simple_config = self._create_ownership_chart_from_sankey(sankey_data, company_name)
                            scripts.append(f"""
                                var ctx{slide_idx}_{idx} = document.getElementById('chart-{slide_idx}-{idx}');
                                if (ctx{slide_idx}_{idx}) {{
                                    new Chart(ctx{slide_idx}_{idx}.getContext('2d'), {self._serialize_chart_config(simple_config)});
                                }}
                            """)
            
            # Handle side-by-side charts
            if slide_type == 'side_by_side':
                for side_idx, side_key in enumerate(['left', 'right']):
                    side_content = content.get(side_key, {})
                    side_chart_data = side_content.get('chart_data')
                    if side_chart_data:
                        side_config = self._create_chart_config(side_chart_data)
                        scripts.append(f"""
                            var ctx{slide_idx}_{side_idx} = document.getElementById('chart-{slide_idx}-{side_idx}');
                            if (ctx{slide_idx}_{side_idx}) {{
                                new Chart(ctx{slide_idx}_{side_idx}.getContext('2d'), {self._serialize_chart_config(side_config)});
                            }}
                        """)
        
        return '\n'.join(scripts)
    
    def _create_chart_config(self, chart_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create Chart.js configuration from our chart data"""
        if not chart_data:
            return self._create_empty_chart_config()
        
        chart_type = chart_data.get('type', 'bar')
        data = chart_data.get('data', {})
        
        # Validate data has required fields
        if not data.get('labels') or not data.get('datasets'):
            logger.warning(f"Chart data missing labels or datasets: {chart_data}")
            return self._create_empty_chart_config()
        
        # Handle Probability Cloud specially
        if chart_type == 'probability_cloud':
            return self._create_probability_cloud_config(chart_data)
        
        # Handle Sankey diagrams specially
        # Sankey requires D3, not Chart.js - return empty config (will be rendered via D3 in HTML)
        if chart_type == 'sankey':
            return {
                'type': 'sankey',
                'data': data,
                'options': chart_data.get('options', {}),
                '_renderViaD3': True  # Flag for special handling
            }
        
        # Handle Waterfall charts specially
        if chart_type == 'waterfall':
            return self._create_waterfall_config(chart_data)
        
        # Modern dark theme color palette - vibrant on dark background
        colors = [
            'rgba(0, 0, 0, 0.9)',        # Pure black
            'rgba(45, 45, 45, 0.9)',     # Dark charcoal
            'rgba(74, 74, 74, 0.9)',     # Medium charcoal
            'rgba(107, 107, 107, 0.9)',  # Steel gray
            'rgba(138, 138, 138, 0.9)',  # Light steel
            'rgba(176, 176, 176, 0.9)',  # Silver
        ]
        
        # Calculate max value for Y-axis formatting
        max_value = 0
        datasets = data.get('datasets', [])
        for dataset in datasets:
            dataset_values = dataset.get('data', [])
            numeric_values = [v for v in dataset_values if isinstance(v, (int, float))]
            if numeric_values:
                dataset_max = max(numeric_values)
                max_value = max(max_value, dataset_max)
        
        # Update datasets with monochrome colors
        for i, dataset in enumerate(datasets):
            if not dataset.get('backgroundColor'):
                dataset['backgroundColor'] = colors[i % len(colors)]
            if not dataset.get('borderColor'):
                dataset['borderColor'] = colors[i % len(colors)]
            dataset['borderWidth'] = 2
            
            # Add smooth curves for line charts
            if chart_type == 'line':
                dataset['tension'] = 0.4
                dataset['fill'] = False
        
        # Generate Y-axis formatter function (JavaScript)
        y_axis_formatter = self._get_js_formatter_function(max_value)
        
        config = {
            'type': chart_type,
            'data': data,
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'plugins': {
                    'legend': {
                        'position': 'bottom',
                        'labels': {
                            'padding': 16,
                            'font': {
                                'size': 12,
                                'family': "'Inter', sans-serif",
                                'weight': '500'
                            },
                            'color': 'hsl(224, 12%, 38%)',
                            'usePointStyle': True
                        }
                    },
                    'title': {
                        'display': bool(chart_data.get('title')),
                        'text': chart_data.get('title', ''),
                        'font': {
                            'size': 16,
                            'family': "'Inter', sans-serif",
                            'weight': '600'
                        },
                        'color': 'hsl(224, 28%, 12%)',
                        'padding': {'bottom': 20}
                    },
                    'tooltip': {
                        'backgroundColor': 'hsl(224, 26%, 16%)',
                        'titleFont': {
                            'family': "'Inter', sans-serif",
                            'size': 13,
                            'weight': '600'
                        },
                        'bodyFont': {
                            'family': "'Inter', sans-serif",
                            'size': 12
                        },
                        'padding': 12,
                        'cornerRadius': 6,
                        'callbacks': {
                            # JavaScript callback for tooltip formatting
                            'label': f'function(context) {{ const value = context.parsed.y; {y_axis_formatter} return context.dataset.label + ": " + formatValue(value); }}'
                        }
                    }
                },
                'scales': {
                    'y': {
                        'beginAtZero': True,
                        'grid': {
                            'color': 'hsl(220, 13%, 88%)',
                            'drawBorder': False
                        },
                        'ticks': {
                            'font': {
                                'family': "'Inter', sans-serif",
                                'size': 11
                            },
                            'color': 'hsl(224, 12%, 38%)',
                            # JavaScript callback for Y-axis tick formatting
                            'callback': f'function(value) {{ {y_axis_formatter} return formatValue(value); }}'
                        },
                        'title': {
                            'display': True,
                            'text': chart_data.get('yAxisLabel', 'Value'),
                            'font': {
                                'family': "'Inter', sans-serif",
                                'size': 12,
                                'weight': '600'
                            },
                            'color': 'hsl(224, 28%, 12%)'
                        }
                    },
                    'x': {
                        'grid': {
                            'display': False
                        },
                        'ticks': {
                            'font': {
                                'family': "'Inter', sans-serif",
                                'size': 10
                            },
                            'color': 'hsl(224, 12%, 38%)',
                            'autoSkip': True,
                            'maxRotation': 45,
                            'minRotation': 0,
                            'maxTicksLimit': 12
                        }
                    }
                }
            }
        }
        
        # Special handling for pie/doughnut charts
        if chart_type in ['pie', 'doughnut']:
            del config['options']['scales']
        
        # Disable animations for PDF rendering
        config['options']['animation'] = False
        config['options']['responsive'] = False
        config['options']['maintainAspectRatio'] = False
        
        return config
    
    def _get_js_formatter_function(self, max_value: float) -> str:
        """Generate JavaScript formatter function for Y-axis ticks - $5M format"""
        if max_value >= 1_000_000:
            # Millions/Billions
            return """
                function formatValue(value) {
                    if (value === 0) return '$0';
                    const millions = value / 1000000;
                    if (millions >= 1000) {
                        const billions = millions / 1000;
                        if (billions >= 10) {
                            return '$' + billions.toFixed(0) + 'B';
                        } else {
                            return '$' + billions.toFixed(1) + 'B';
                        }
                    } else if (millions >= 10) {
                        return '$' + millions.toFixed(0) + 'M';
                    } else if (millions >= 1) {
                        return '$' + millions.toFixed(0) + 'M';
                    } else {
                        return '$' + millions.toFixed(1) + 'M';
                    }
                }
            """
        elif max_value >= 1000:
            # Thousands
            return """
                function formatValue(value) {
                    if (value === 0) return '$0';
                    return '$' + (value / 1000).toFixed(0) + 'K';
                }
            """
        else:
            # Small values
            return """
                function formatValue(value) {
                    return '$' + value.toFixed(0);
                }
            """
    
    def _create_empty_chart_config(self) -> Dict[str, Any]:
        """Create an empty chart configuration for missing data"""
        return {
            'type': 'bar',
            'data': {
                'labels': ['No Data'],
                'datasets': [{
                    'label': 'No data available',
                    'data': [0],
                    'backgroundColor': 'rgba(156, 163, 175, 0.3)',
                    'borderColor': 'rgba(156, 163, 175, 1)',
                    'borderWidth': 1
                }]
            },
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'animation': False,  # Disable animations for PDF rendering
                'plugins': {
                    'title': {
                        'display': True,
                        'text': 'Chart data not available',
                        'font': {
                            'size': 14,
                            'family': "'Inter', sans-serif"
                        },
                        'color': '#9ca3af'
                    },
                    'legend': {
                        'display': False
                    }
                }
            }
        }
    
    def _create_probability_cloud_config(self, chart_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create Chart.js configuration for probability cloud visualization"""
        data = chart_data.get('data', {})
        scenario_curves = data.get('scenario_curves', [])
        breakpoint_clouds = data.get('breakpoint_clouds', [])
        config = data.get('config', {})
        
        # Extract x-axis config
        x_config = config.get('x_axis', {})
        x_min = x_config.get('min', 10_000_000)
        x_max = x_config.get('max', 10_000_000_000)
        
        # Create logarithmic scale points for x-axis
        import math
        num_points = 20
        log_min = math.log10(x_min)
        log_max = math.log10(x_max)
        log_step = (log_max - log_min) / (num_points - 1)
        x_values = [10 ** (log_min + i * log_step) for i in range(num_points)]
        
        # Format x-axis labels
        labels = []
        for val in x_values:
            if val >= 1e9:
                labels.append(f'${val/1e9:.1f}B')
            elif val >= 1e6:
                labels.append(f'${val/1e6:.0f}M')
            else:
                labels.append(f'${val/1e3:.0f}K')
        
        # Create datasets for each scenario curve
        datasets = []
        
        # Add top scenarios sorted by probability
        sorted_scenarios = sorted(scenario_curves, 
                                key=lambda s: s.get('probability', 0), 
                                reverse=True)
        
        # Color palette for scenarios
        scenario_colors = {
            'IPO': 'rgba(16, 185, 129, ',  # Green
            'Unicorn': 'rgba(17, 24, 39, ',     # Darkest gray
            'Strong': 'rgba(31, 41, 55, ',      # Dark gray  
            'Strategic': 'rgba(55, 65, 81, ',   # Medium gray
            'Base': 'rgba(99, 102, 241, ',  # Indigo
            'PE': 'rgba(99, 102, 241, ',  # Indigo
            'Modest': 'rgba(245, 158, 11, ',  # Amber
            'Acquihire': 'rgba(245, 158, 11, ',  # Amber
            'Downside': 'rgba(239, 68, 68, ',  # Red
            'Distressed': 'rgba(239, 68, 68, ',  # Red
        }
        
        # Dynamically determine how many scenarios to show (all if <= 15, otherwise top 12)
        max_scenarios = min(len(sorted_scenarios), 15) if len(sorted_scenarios) <= 15 else 12
        
        # Add scenarios based on probability
        for idx, scenario in enumerate(sorted_scenarios[:max_scenarios]):
            return_curve = scenario.get('return_curve', {})
            exit_values = return_curve.get('exit_values', [])
            return_multiples = return_curve.get('return_multiples', [])
            
            if not exit_values or not return_multiples:
                continue
            
            # Interpolate returns for our x_values
            interpolated_returns = []
            for x_val in x_values:
                # Find the appropriate return value for this exit value
                if x_val <= exit_values[0]:
                    interpolated_returns.append(return_multiples[0])
                elif x_val >= exit_values[-1]:
                    interpolated_returns.append(return_multiples[-1])
                else:
                    # Linear interpolation
                    for i in range(len(exit_values) - 1):
                        if exit_values[i] <= x_val <= exit_values[i+1]:
                            # Interpolate
                            t = (x_val - exit_values[i]) / (exit_values[i+1] - exit_values[i])
                            interp_val = return_multiples[i] + t * (return_multiples[i+1] - return_multiples[i])
                            interpolated_returns.append(interp_val)
                            break
            
            # Determine color based on scenario name
            scenario_name = scenario.get('name', '')
            color_base = 'rgba(156, 163, 175, '  # Default gray
            for keyword, color in scenario_colors.items():
                if keyword in scenario_name:
                    color_base = color
                    break
            
            # Opacity based on probability
            opacity = 0.3 + safe_multiply(ensure_numeric(scenario.get('probability', 0.1)), 0.7, 0.1)
            color = f"{color_base}{opacity})"
            
            datasets.append({
                'label': scenario_name[:30],  # Limit label length
                'data': interpolated_returns,
                'borderColor': color,
                'backgroundColor': color.replace(str(opacity), '0.1'),
                'borderWidth': 2,
                'tension': 0.4,  # Smooth curves
                'fill': False,
                'pointRadius': 0,  # Hide points for cleaner look
                'pointHoverRadius': 4,
                'probability': scenario.get('probability', 0)  # Store for tooltip
            })
        
        # Add breakpoint vertical lines as annotations
        annotations = {}
        for cloud in breakpoint_clouds:
            median = cloud.get('median', 0)
            if median > 0:
                # Find closest x-axis index
                closest_idx = min(range(len(x_values)), 
                                key=lambda i: abs(x_values[i] - median))
                
                annotation_id = f"breakpoint_{cloud.get('type', '')}"
                annotations[annotation_id] = {
                    'type': 'line',
                    'xMin': closest_idx,
                    'xMax': closest_idx,
                    'borderColor': cloud.get('color', 'rgba(0, 0, 0, 0.3)'),
                    'borderWidth': 2,
                    'borderDash': [5, 3],
                    'label': {
                        'content': cloud.get('label', ''),
                        'enabled': True,
                        'position': 'start',
                        'rotation': -90,
                        'backgroundColor': 'rgba(255, 255, 255, 0.8)',
                        'color': cloud.get('color', '#000'),
                        'font': {
                            'size': 10
                        }
                    }
                }
        
        return {
            'type': 'line',
            'data': {
                'labels': labels,
                'datasets': datasets
            },
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'animation': False,  # Disable animations for PDF rendering
                'interaction': {
                    'mode': 'index',
                    'intersect': False
                },
                'plugins': {
                    'legend': {
                        'position': 'right',
                        'labels': {
                            'padding': 10,
                            'font': {
                                'size': 10
                            },
                            # Note: Filter function will be added via JavaScript string
                            'filter': 'function(item, chart) { return item.datasetIndex < 5; }'
                        }
                    },
                    'title': {
                        'display': True,
                        'text': chart_data.get('title', 'Probability-Weighted Return Scenarios'),
                        'font': {
                            'size': 16,
                            'weight': '600'
                        }
                    },
                    'annotation': {
                        'annotations': annotations
                    },
                    'tooltip': {
                        'callbacks': {
                            # Note: Callback function will be added via JavaScript string
                            'label': 'function(context) { const dataset = context.dataset; const value = context.parsed.y; const prob = dataset.probability || 0; return dataset.label + ": " + value.toFixed(1) + "x (" + (prob * 100).toFixed(0) + "% prob)"; }'
                        }
                    }
                },
                'scales': {
                    'x': {
                        'title': {
                            'display': True,
                            'text': 'Exit Valuation',
                            'font': {
                                'size': 12
                            }
                        },
                        'grid': {
                            'color': 'rgba(0, 0, 0, 0.05)'
                        },
                        'ticks': {
                            'maxTicksLimit': 8,
                            'font': {
                                'size': 10
                            },
                            'autoSkip': True
                        }
                    },
                    'y': {
                        'title': {
                            'display': True,
                            'text': 'Return Multiple',
                            'font': {
                                'size': 12
                            }
                        },
                        'beginAtZero': True,
                        'max': config.get('y_axis', {}).get('max', 10),
                        'grid': {
                            'color': 'rgba(0, 0, 0, 0.05)'
                        },
                        'ticks': {
                            # Note: Callback function will be added via JavaScript string
                            'callback': 'function(value) { return value + "x"; }',
                            'font': {
                                'size': 10
                            }
                        }
                    }
                }
            }
        }
    
    def _create_sankey_config(self, chart_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create configuration for Sankey diagrams (using alternative visualization)"""
        # Since Chart.js doesn't support Sankey natively, we'll convert to a flow visualization
        # This returns a special config that will be handled differently in the JavaScript
        sankey_data = chart_data.get('data', {})
        nodes = sankey_data.get('nodes', [])
        links = sankey_data.get('links', [])
        
        # Convert Sankey data to a stacked bar chart representation
        # Group by rounds/stages
        stages = {}
        for link in links:
            source_idx = link.get('source', 0)
            target_idx = link.get('target', 0)
            value = link.get('value', 0)
            
            if source_idx < len(nodes) and target_idx < len(nodes):
                source_name = nodes[source_idx].get('name', 'Source')
                target_name = nodes[target_idx].get('name', 'Target')
                
                # Extract round information from node names
                if 'Round' in source_name or 'Round' in target_name:
                    round_name = source_name if 'Round' in source_name else target_name
                    if round_name not in stages:
                        stages[round_name] = {'founders': 0, 'investors': 0, 'employees': 0}
                    
                    # Categorize the ownership
                    if 'Founder' in target_name:
                        stages[round_name]['founders'] = value
                    elif 'Investor' in target_name:
                        stages[round_name]['investors'] = value
                    elif 'Employee' in target_name:
                        stages[round_name]['employees'] = value
        
        # Convert to stacked bar chart format
        labels = list(stages.keys())
        datasets = [
            {
                'label': 'Founders',
                'data': [stages[label].get('founders', 0) for label in labels],
                'backgroundColor': 'rgba(99, 102, 241, 0.8)'
            },
            {
                'label': 'Investors',
                'data': [stages[label].get('investors', 0) for label in labels],
                'backgroundColor': 'rgba(55, 65, 81, 0.8)'
            },
            {
                'label': 'Employees',
                'data': [stages[label].get('employees', 0) for label in labels],
                'backgroundColor': 'rgba(107, 114, 128, 0.7)'
            }
        ]
        
        return {
            'type': 'bar',
            'data': {
                'labels': labels,
                'datasets': datasets
            },
            'options': {
                'responsive': True,
                'maintainAspectRatio': False,
                'animation': False,  # Disable animations for PDF rendering
                'plugins': {
                    'title': {
                        'display': True,
                        'text': 'Ownership Evolution by Round',
                        'font': {'size': 16, 'weight': '600'}
                    },
                    'legend': {
                        'position': 'bottom'
                    },
                    'tooltip': {
                        'mode': 'index',
                        'intersect': False
                    }
                },
                'scales': {
                    'x': {
                        'stacked': True,
                        'title': {
                            'display': True,
                            'text': 'Funding Round'
                        }
                    },
                    'y': {
                        'stacked': True,
                        'title': {
                            'display': True,
                            'text': 'Ownership %'
                        },
                        'max': 100
                    }
                }
            }
        }
    
    async def export_to_pdf_async(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]]) -> bytes:
        """Export deck to PDF by rendering actual Next.js page with fallback to sync path"""
        if not PLAYWRIGHT_AVAILABLE:
            error_msg = (
                "Playwright is required for PDF export but is not installed or not properly configured. "
                "Please install with: pip install playwright && playwright install chromium"
            )
            logger.error(error_msg)
            logger.warning("[PDF_EXPORT] Falling back to sync path...")
            return self.export_to_pdf(deck_data)
        
        try:
            from playwright.async_api import async_playwright
        except ImportError as e:
            error_msg = f"Failed to import Playwright: {e}. Please ensure Playwright is properly installed."
            logger.error(error_msg)
            logger.warning("[PDF_EXPORT] Falling back to sync path...")
            return self.export_to_pdf(deck_data)
        
        # Pre-render complex charts and inject images into deck data
        logger.info("[PDF_EXPORT] Pre-rendering complex charts...")
        deck_data_with_images = await self._inject_prerendered_charts(deck_data)
        
        # Store deck data temporarily for PDF rendering
        try:
            deck_id = deck_storage.store_deck(deck_data_with_images)
            logger.info(f"[PDF_EXPORT] Stored deck data with ID: {deck_id}")
        except Exception as e:
            logger.warning(f"[PDF_EXPORT] Failed to store deck data: {e}, falling back to sync path")
            return self.export_to_pdf_sync(deck_data)
        
        if not await _ensure_playwright_browser_async():
            raise Exception("Playwright Chromium browser is not available. Please run: playwright install chromium")
        
        try:
            async with async_playwright() as p:
                try:
                    browser = await p.chromium.launch(
                        headless=True,
                        args=[
                            '--no-sandbox',
                            '--disable-setuid-sandbox',
                            '--disable-dev-shm-usage',
                            '--disable-gpu',
                            '--single-process',
                            '--no-zygote'
                        ]
                    )
                except Exception as e:
                    error_msg = f"Failed to launch browser: {e}. Browser may need to be reinstalled. Try: playwright install chromium"
                    logger.error(error_msg)
                    raise Exception(error_msg)
                
                # Use 1024x768 viewport to match slide dimensions (4:3 aspect ratio)
                page = await browser.new_page(viewport={'width': 1024, 'height': 768})
                
                # Navigate to deck presentation page with PDF mode
                deck_url = f"http://localhost:3001/deck-agent?deckId={deck_id}&pdfMode=true"
                logger.info(f"[PDF_EXPORT] Navigating to: {deck_url}")
                
                try:
                    await page.goto(deck_url, wait_until='networkidle', timeout=30000)
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Failed to connect to frontend: {e}, falling back to sync path")
                    await browser.close()
                    return self.export_to_pdf_sync(deck_data)
                
                # Wait for React to hydrate and deck to load
                logger.info("[PDF_EXPORT] Waiting for page to load...")
                await page.wait_for_function("""
                () => {
                    // Wait for deck data to be loaded - check for slides in the DOM
                    const slideElements = document.querySelectorAll('[data-testid="deck-presentation"] > div, .slide, [class*="slide"]');
                    const hasSlides = slideElements.length > 0;
                    
                    // Also check for deck container
                    const deckContainer = document.querySelector('[data-testid="deck-presentation"]');
                    const hasContainer = !!deckContainer;
                    
                    // Check for actual content (not just empty divs)
                    const hasContent = document.body.innerText.length > 50;
                    
                    return hasContainer && (hasSlides || hasContent);
                }
                """, timeout=20000)
                
                # Additional wait for deck content to be fully rendered
                logger.info("[PDF_EXPORT] Waiting for deck content...")
                await page.wait_for_function("""
                () => {
                    // Check for actual slide content with meaningful text
                    const slideContainer = document.querySelector('[data-testid="deck-presentation"]');
                    if (!slideContainer) return false;
                    
                    // Check for slide titles or content
                    const hasTitles = slideContainer.querySelectorAll('h1, h2, .slide-title').length > 0;
                    const hasText = slideContainer.innerText.length > 100;
                    
                    // In PDF mode, should have multiple slides or at least one with content
                    const slides = slideContainer.querySelectorAll('div[style*="1024px"], .slide, [class*="page-break"]');
                    
                    return (hasTitles || hasText) && (slides.length > 0 || slideContainer.innerText.length > 100);
                }
                """, timeout=20000)
                
                # Wait for all charts to render with intelligent detection
                logger.info("[PDF_EXPORT] Waiting for charts to render...")
                
                # First, wait for Chart.js and D3.js to be available
                try:
                    await page.wait_for_function("typeof Chart !== 'undefined'", timeout=15000)
                    logger.info("[PDF_EXPORT] Chart.js loaded")
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Chart.js not loaded: {e}")
                
                try:
                    await page.wait_for_function("typeof d3 !== 'undefined' && d3 !== null", timeout=10000)
                    logger.info("[PDF_EXPORT] D3.js loaded")
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] D3.js not loaded (may not be needed): {e}")
                
                # Wait for chart containers to appear
                try:
                    await page.wait_for_function("""
                    () => {
                        // Wait for at least one chart container or confirm no charts exist
                        const chartContainers = document.querySelectorAll('[data-chart-type]');
                        const hasAnyCharts = chartContainers.length > 0 || 
                                           document.querySelectorAll('svg, canvas').length > 0;
                        return hasAnyCharts || document.body.innerText.length > 100;
                    }
                    """, timeout=10000)
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Timeout waiting for chart containers: {e}")
                
                # Progressive waiting: check for data-chart-ready attributes first (fastest)
                logger.info("[PDF_EXPORT] Checking for chart-ready attributes...")
                chart_ready_check = """
                () => {
                    const chartContainers = document.querySelectorAll('[data-chart-type]');
                    
                    // No charts = ready
                    if (chartContainers.length === 0) {
                        // Fallback: check for any canvas/SVG that might be a chart
                        const allCanvases = document.querySelectorAll('canvas');
                        const allSvgs = document.querySelectorAll('svg');
                        if (allCanvases.length === 0 && allSvgs.length === 0) {
                            return { ready: true, reason: 'no_charts', renderedCount: 0, totalCharts: 0 };
                        }
                        // If we have canvases/SVGs but no containers, assume they're ready
                        // (they might be from other chart libraries)
                        return { ready: true, reason: 'standalone_charts', renderedCount: allCanvases.length + allSvgs.length, totalCharts: allCanvases.length + allSvgs.length };
                    }
                    
                    let readyCount = 0;
                    let totalCharts = chartContainers.length;
                    
                    // Check data-chart-ready attribute (fastest method)
                    for (let container of chartContainers) {
                        if (container.getAttribute('data-chart-ready') === 'true') {
                            readyCount++;
                            continue;
                        }
                        
                        // Fallback: check for actual content
                        const chartType = container.getAttribute('data-chart-type') || '';
                        
                        // D3-based charts (SVG)
                        if (['sankey', 'sunburst', 'heatmap', 'side_by_side_sankey', 'probability_cloud'].includes(chartType)) {
                            const svg = container.querySelector('svg');
                            if (svg && svg.children.length >= 2) {
                                const hasContent = svg.querySelector('path, circle, rect, line, text, g') !== null;
                                if (hasContent) {
                                    readyCount++;
                                    container.setAttribute('data-chart-ready', 'true');
                                }
                            }
                        }
                        // Recharts (SVG)
                        else if (['waterfall', 'bubble', 'funnel', 'radialBar', 'timeline_valuation'].includes(chartType)) {
                            const svg = container.querySelector('.recharts-wrapper svg') || container.querySelector('svg');
                            if (svg && svg.children.length >= 2) {
                                readyCount++;
                                container.setAttribute('data-chart-ready', 'true');
                            }
                        }
                        // Canvas charts (Chart.js)
                        else {
                            const canvas = container.querySelector('canvas');
                            if (canvas) {
                                try {
                                    const ctx = canvas.getContext('2d');
                                    if (ctx) {
                                        // Check a larger sample area for better detection
                                        const sampleWidth = Math.min(canvas.width, 50);
                                        const sampleHeight = Math.min(canvas.height, 50);
                                        const imageData = ctx.getImageData(0, 0, sampleWidth, sampleHeight);
                                        const data = imageData.data;
                                        
                                        // Check if any pixels are non-transparent (alpha > 0)
                                        let hasContent = false;
                                        let nonTransparentPixels = 0;
                                        for (let i = 3; i < data.length; i += 4) {
                                            if (data[i] > 10) { // Alpha > 10 (not fully transparent)
                                                nonTransparentPixels++;
                                                if (nonTransparentPixels > 5) { // At least 5 non-transparent pixels
                                                    hasContent = true;
                                                    break;
                                                }
                                            }
                                        }
                                        
                                        if (hasContent) {
                                            readyCount++;
                                            container.setAttribute('data-chart-ready', 'true');
                                        }
                                    }
                                } catch (e) {
                                    // Security error or canvas not ready - wait a bit more
                                    // Don't mark as ready yet, will retry
                                }
                            }
                        }
                    }
                    
                    return { 
                        ready: readyCount === totalCharts,
                        renderedCount: readyCount,
                        totalCharts: totalCharts
                    };
                }
                """
                
                # Progressive waiting: check every 500ms, max 20 seconds
                max_attempts = 40  # 40 attempts x 500ms = 20 seconds
                all_charts_ready = False
                
                for attempt in range(max_attempts):
                    try:
                        result = await page.evaluate(chart_ready_check)
                        if result.get('ready'):
                            logger.info(f"[PDF_EXPORT] All {result.get('totalCharts', 0)} charts rendered successfully")
                            all_charts_ready = True
                            break
                        else:
                            rendered = result.get('renderedCount', 0)
                            total = result.get('totalCharts', 0)
                            if attempt % 4 == 0:  # Log every 2 seconds
                                logger.info(f"[PDF_EXPORT] Charts rendering: {rendered}/{total} (attempt {attempt + 1}/{max_attempts})")
                    except Exception as e:
                        logger.warning(f"[PDF_EXPORT] Error checking chart status: {e}")
                    
                    await page.wait_for_timeout(500)
                
                if not all_charts_ready:
                    # Final check
                    try:
                        result = await page.evaluate(chart_ready_check)
                        rendered = result.get('renderedCount', 0)
                        total = result.get('totalCharts', 0)
                        logger.warning(f"[PDF_EXPORT] Timeout waiting for all charts. Proceeding with {rendered}/{total} charts ready.")
                    except Exception as e:
                        logger.warning(f"[PDF_EXPORT] Error in final chart check: {e}")
                
                # Additional wait to ensure all animations and transitions complete
                logger.info("[PDF_EXPORT] Waiting for animations to complete...")
                await page.wait_for_timeout(3000)  # Wait 3 seconds for any animations
                
                # Wait for any pending animations/renders - multiple frames for complex charts
                try:
                    await page.evaluate("""() => new Promise(resolve => {
                        requestAnimationFrame(() => {
                            requestAnimationFrame(() => {
                                requestAnimationFrame(resolve);
                            });
                        });
                    })""")
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Error waiting for animation frame: {e}")
                
                # Final verification: check that charts have actual content
                try:
                    final_check = await page.evaluate("""
                    () => {
                        const chartContainers = document.querySelectorAll('[data-chart-type]');
                        let allHaveContent = true;
                        
                        for (let container of chartContainers) {
                            const chartType = container.getAttribute('data-chart-type') || '';
                            
                            // Check D3/SVG charts
                            if (['sankey', 'sunburst', 'heatmap', 'side_by_side_sankey', 'probability_cloud'].includes(chartType)) {
                                const svg = container.querySelector('svg');
                                if (!svg || svg.children.length < 2) {
                                    allHaveContent = false;
                                    break;
                                }
                                // Check for actual SVG elements
                                const hasElements = svg.querySelector('path, circle, rect, line, text, g') !== null;
                                if (!hasElements) {
                                    allHaveContent = false;
                                    break;
                                }
                            }
                            // Check canvas charts
                            else {
                                const canvas = container.querySelector('canvas');
                                if (canvas) {
                                    try {
                                        const ctx = canvas.getContext('2d');
                                        if (ctx) {
                                            const imageData = ctx.getImageData(0, 0, Math.min(canvas.width, 50), Math.min(canvas.height, 50));
                                            const data = imageData.data;
                                            let hasContent = false;
                                            for (let i = 3; i < data.length; i += 4) {
                                                if (data[i] > 10) {
                                                    hasContent = true;
                                                    break;
                                                }
                                            }
                                            if (!hasContent) {
                                                allHaveContent = false;
                                                break;
                                            }
                                        }
                                    } catch (e) {
                                        // Assume content exists if we can't check
                                    }
                                }
                            }
                        }
                        
                        return allHaveContent || chartContainers.length === 0;
                    }
                    """)
                    
                    if not final_check:
                        logger.warning("[PDF_EXPORT] Some charts may not have content, but proceeding with PDF generation")
                        await page.wait_for_timeout(2000)  # Extra wait
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Error in final chart content check: {e}")
                
                # Take a screenshot for debugging
                try:
                    screenshot_path = f"/tmp/deck_export_debug_{deck_id}.png"
                    await page.screenshot(path=screenshot_path, full_page=True)
                    logger.info(f"[PDF_EXPORT] Debug screenshot saved: {screenshot_path}")
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Could not save debug screenshot: {e}")
                
                # Log final chart status for debugging
                try:
                    final_status = await page.evaluate("""
                    () => {
                        const chartContainers = document.querySelectorAll('[data-chart-type]');
                        const readyCharts = document.querySelectorAll('[data-chart-ready="true"]');
                        const svgs = document.querySelectorAll('svg');
                        const canvases = document.querySelectorAll('canvas');
                        
                        return {
                            totalCharts: chartContainers.length,
                            readyCharts: readyCharts.length,
                            svgCount: svgs.length,
                            canvasCount: canvases.length,
                            allReady: chartContainers.length === readyCharts.length
                        };
                    }
                    """)
                    logger.info(f"[PDF_EXPORT] Final chart status: {final_status}")
                except Exception as e:
                    logger.warning(f"[PDF_EXPORT] Error getting final chart status: {e}")
                
                # Generate PDF with custom dimensions matching slide size (1024x768)
                logger.info("[PDF_EXPORT] Generating PDF...")
                try:
                    # Use custom page size: 1024x768 pixels = 10.67x8 inches at 96 DPI
                    pdf_bytes = await page.pdf(
                        width='10.67in',
                        height='8in',
                        print_background=True,
                        margin={'top': '0in', 'bottom': '0in', 'left': '0in', 'right': '0in'},
                        prefer_css_page_size=True
                    )
                    
                    if not pdf_bytes or len(pdf_bytes) < 1000:
                        error_msg = f"PDF generation produced invalid output ({len(pdf_bytes) if pdf_bytes else 0} bytes)"
                        logger.error(f"[PDF_EXPORT] {error_msg}")
                        logger.warning("[PDF_EXPORT] Falling back to sync path...")
                        await browser.close()
                        return self.export_to_pdf_sync(deck_data)
                    
                    logger.info(f"[PDF_EXPORT] PDF generated successfully: {len(pdf_bytes)} bytes")
                    
                    # Clean up deck data
                    try:
                        cleanup_response = await httpx.AsyncClient().delete(f"{backend_url}/api/deck-storage/{deck_id}")
                        if cleanup_response.status_code == 200:
                            logger.info(f"[PDF_EXPORT] Cleaned up deck data: {deck_id}")
                    except Exception as e:
                        logger.warning(f"[PDF_EXPORT] Could not clean up deck data: {e}")
                    
                    return pdf_bytes
                    
                except (ConnectionError, TimeoutError, Exception) as e:
                    error_msg = f"Error during PDF generation: {str(e)}"
                    logger.error(f"[PDF_EXPORT] {error_msg}")
                    
                    # Capture screenshot on error for debugging
                    try:
                        error_screenshot_path = f"/tmp/deck_export_error_{deck_id}.png"
                        await page.screenshot(path=error_screenshot_path, full_page=True)
                        logger.error(f"[PDF_EXPORT] Error screenshot saved: {error_screenshot_path}")
                    except Exception as screenshot_error:
                        logger.warning(f"[PDF_EXPORT] Could not save error screenshot: {screenshot_error}")
                    
                    # Fallback to sync path for connection errors, timeouts, or empty PDFs
                    logger.warning("[PDF_EXPORT] Falling back to sync path...")
                    try:
                        await browser.close()
                    except:
                        pass
                    return self.export_to_pdf_sync(deck_data)
                finally:
                    try:
                        await browser.close()
                    except:
                        pass
                
        finally:
            # Clean up stored deck data
            deck_storage.delete_deck(deck_id)
            logger.info(f"[PDF_EXPORT] Cleaned up deck data: {deck_id}")
    
    async def _prerender_complex_charts(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]]) -> Dict[str, str]:
        """Pre-render complex charts to high-quality PNG images"""
        
        # Handle both dict with slides key and direct list of slides
        if isinstance(deck_data, list):
            slides = deck_data
        else:
            slides = deck_data.get("slides", deck_data.get("deck_slides", []))
        
        chart_images = {}
        
        for slide_idx, slide_data in enumerate(slides):
            content = slide_data.get("content", {})
            
            # Handle both single chart_data and multiple charts
            charts_to_render = []
            
            # Check for single chart_data
            chart_data = content.get("chart_data")
            if chart_data:
                charts_to_render.append((0, chart_data))
            
            # Check for multiple charts
            charts = content.get("charts", [])
            for chart_idx, chart_data in enumerate(charts):
                charts_to_render.append((chart_idx, chart_data))
            
            # Render each chart
            for chart_idx, chart_data in charts_to_render:
                # Get chart type from content or chart_data
                chart_type = content.get("chart_type", "")
                if not chart_type:
                    chart_type = chart_data.get("type", "")
                
                # Check if this chart type should be pre-rendered
                if chart_renderer.should_prerender_chart(chart_type):
                    chart_key = f"slide_{slide_idx}_chart_{chart_idx}"
                    
                    try:
                        logger.info(f"[PRERENDER] Rendering {chart_type} chart for slide {slide_idx}")
                        
                        # Render chart to high-quality PNG
                        img_base64 = await chart_renderer.render_tableau_chart(
                            chart_type, chart_data, width=800, height=400
                        )
                        
                        if img_base64:
                            chart_images[chart_key] = img_base64
                            logger.info(f"[PRERENDER] Successfully rendered {chart_type} chart")
                        else:
                            logger.warning(f"[PRERENDER] Failed to render {chart_type} chart")
                            
                    except Exception as e:
                        logger.error(f"[PRERENDER] Error rendering {chart_type} chart: {e}")
        
        logger.info(f"[PRERENDER] Pre-rendered {len(chart_images)} complex charts")
        return chart_images

    async def _inject_prerendered_charts(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]]) -> Union[Dict[str, Any], List[Dict[str, Any]]]:
        """Pre-render complex charts and inject images directly into deck data"""
        import copy
        
        # Deep copy to avoid modifying original data
        deck_data_copy = copy.deepcopy(deck_data)
        
        # Handle both dict with slides key and direct list of slides
        if isinstance(deck_data_copy, list):
            slides = deck_data_copy
        else:
            slides = deck_data_copy.get("slides", deck_data_copy.get("deck_slides", []))
        
        for slide_idx, slide_data in enumerate(slides):
            content = slide_data.get("content", {})
            
            # Handle both single chart_data and multiple charts
            charts_to_render = []
            
            # Check for single chart_data
            chart_data = content.get("chart_data")
            if chart_data:
                charts_to_render.append((0, chart_data, "chart_data"))
            
            # Check for multiple charts
            charts = content.get("charts", [])
            for chart_idx, chart_data in enumerate(charts):
                charts_to_render.append((chart_idx, chart_data, "charts"))
            
            # Render each chart and inject image
            for chart_idx, chart_data, chart_key in charts_to_render:
                chart_type = chart_data.get("type", "")
                
                # Check if this chart type should be pre-rendered
                if chart_renderer.should_prerender_chart(chart_type):
                    try:
                        logger.info(f"[INJECT_PRERENDER] Rendering {chart_type} chart for slide {slide_idx}")
                        
                        # Render chart to high-quality PNG
                        img_base64 = await chart_renderer.render_tableau_chart(
                            chart_type, chart_data, width=800, height=400
                        )
                        
                        if img_base64:
                            # Inject pre-rendered image into chart_data
                            if chart_key == "chart_data":
                                slide_data["content"]["chart_data"] = {
                                    "type": "image",
                                    "src": f"data:image/png;base64,{img_base64}",
                                    "alt": f"{chart_type} chart",
                                    "title": chart_data.get("title", ""),
                                    "original_type": chart_type,
                                    "original_data": chart_data
                                }
                            else:  # charts array
                                slide_data["content"]["charts"][chart_idx] = {
                                    "type": "image",
                                    "src": f"data:image/png;base64,{img_base64}",
                                    "alt": f"{chart_type} chart",
                                    "title": chart_data.get("title", ""),
                                    "original_type": chart_type,
                                    "original_data": chart_data
                                }
                            
                            logger.info(f"[INJECT_PRERENDER] Successfully injected {chart_type} chart image")
                        else:
                            logger.warning(f"[INJECT_PRERENDER] Failed to render {chart_type} chart")
                            
                    except Exception as e:
                        logger.error(f"[INJECT_PRERENDER] Error rendering {chart_type} chart: {e}")
        
        logger.info(f"[INJECT_PRERENDER] Injected pre-rendered charts into deck data")
        return deck_data_copy

    def export_to_pdf(self, deck_data: Union[Dict[str, Any], List[Dict[str, Any]]]) -> bytes:
        """Sync wrapper for PDF export - detects if we're in an async context"""
        import asyncio
        import inspect
        
        # Check if we're in an async context
        try:
            loop = asyncio.get_running_loop()
            # We're in an async context, need to run in a new thread
            import concurrent.futures
            import threading
            
            result = None
            exception = None
            
            def run_sync():
                nonlocal result, exception
                try:
                    # Create a new event loop for this thread
                    new_loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(new_loop)
                    try:
                        result = new_loop.run_until_complete(self.export_to_pdf_async(deck_data))
                    finally:
                        new_loop.close()
                except Exception as e:
                    exception = e
            
            thread = threading.Thread(target=run_sync)
            thread.start()
            thread.join(timeout=30)
            
            if exception:
                raise exception
            if result is None:
                raise Exception("PDF generation timed out")
            return result
            
        except RuntimeError:
            # No event loop, we can use sync API directly
            if not PLAYWRIGHT_AVAILABLE:
                raise Exception("Playwright is required for PDF export")
            
            # Note: Chart pre-rendering is handled in export_to_pdf_async
            # For sync mode, we'll generate HTML without pre-rendered charts
            chart_images = {}
            
            html = self._generate_html_deck(deck_data, chart_images)
            
            # Save for debugging
            with open("/tmp/deck_debug.html", "w") as f:
                f.write(html)
            
            # Ensure Chromium browser is installed
            if not _ensure_playwright_browser():
                raise Exception("Playwright Chromium browser is not available. Please run: playwright install chromium")
            
            with sync_playwright() as p:
                try:
                    browser = p.chromium.launch(
                        headless=True,
                        args=[
                            '--no-sandbox',
                            '--disable-setuid-sandbox',
                            '--disable-dev-shm-usage',
                            '--disable-gpu',
                            '--single-process',
                            '--no-zygote'
                        ]
                    )
                except Exception as e:
                    error_msg = f"Failed to launch browser: {e}. Browser may need to be reinstalled. Try: playwright install chromium"
                    logger.error(error_msg)
                    raise Exception(error_msg)
                
                # Use 1024x768 viewport to match slide dimensions (4:3 aspect ratio)
                page = browser.new_page(viewport={'width': 1024, 'height': 768})
                
                # Enable console logging for debugging
                page.on("console", lambda msg: logger.info(f"[PDF_CONSOLE] {msg.text}"))
                
                page.set_content(html)
                page.wait_for_load_state('networkidle')
                
                # Wait for Chart.js and D3 to load with retry mechanism
                logger.info("[PDF_GEN] Waiting for Chart.js and D3 to load...")
                chart_loaded = False
                d3_loaded = False
                
                for attempt in range(3):
                    try:
                        page.wait_for_function("typeof Chart !== 'undefined'", timeout=15000)
                        chart_loaded = True
                        logger.info("[PDF_GEN] Chart.js loaded successfully")
                        
                        # Also wait for D3 for Sankey charts (check if d3 is available globally from CDN)
                        page.wait_for_function("typeof d3 !== 'undefined' && d3 !== null", timeout=10000)
                        d3_loaded = True
                        logger.info("[PDF_GEN] D3 loaded successfully")
                        break
                    except Exception as e:
                        logger.warning(f"[PDF_GEN] Chart.js/D3 load attempt {attempt + 1} failed: {e}")
                        if attempt < 2:
                            page.wait_for_timeout(2000)
                            page.reload()
                        else:
                            logger.error("[PDF_GEN] Chart.js/D3 failed to load after 3 attempts")
                
                if not chart_loaded:
                    logger.warning("[PDF_GEN] Proceeding without Chart.js - charts may not render")
                if not d3_loaded:
                    logger.warning("[PDF_GEN] Proceeding without D3 - Sankey charts may not render")
                
                # Wait for charts to be initialized with progressive checking
                logger.info("[PDF_GEN] Waiting for charts to render...")
                
                # Progressive wait with early exit if charts are ready
                max_wait = 20000  # Maximum 20 seconds for complex charts
                check_interval = 500  # Check every 500ms
                waited = 0
                
                while waited < max_wait:
                    # Check if any charts exist and are rendering (canvas + SVG)
                    charts_ready = page.evaluate("""() => {
                        // Check for chart containers with data-chart-type attribute
                        const chartContainers = document.querySelectorAll('[data-chart-type]');
                        if (chartContainers.length === 0) {
                            // Fallback: check for canvas or SVG elements
                            const canvases = document.querySelectorAll('canvas');
                            const svgs = document.querySelectorAll('svg');
                            return canvases.length === 0 && svgs.length === 0; // No charts to wait for
                        }
                        
                        // Check each chart container
                        for (let container of chartContainers) {
                            const chartType = container.dataset.chartType?.toLowerCase() || '';
                            
                            // D3-based charts (SVG) - sankey, sunburst, heatmap
                            if (['sankey', 'sunburst', 'heatmap'].includes(chartType)) {
                                const svg = container.querySelector('svg');
                                if (!svg || svg.children.length < 2) return false;
                            }
                            // Recharts (SVG) - waterfall, bubble, funnel, radialBar, treemap, composed
                            else if (['waterfall', 'bubble', 'funnel', 'radialbar', 'treemap', 'composed'].includes(chartType)) {
                                const svg = container.querySelector('.recharts-wrapper svg') || container.querySelector('svg');
                                if (!svg || svg.children.length < 2) return false;
                            }
                            // Chart.js (Canvas) - bar, line, pie, etc.
                            else {
                                const canvas = container.querySelector('canvas');
                                if (!canvas) return false;
                                
                                const ctx = canvas.getContext('2d');
                                if (!ctx) return false;
                                
                                try {
                                    const imageData = ctx.getImageData(0, 0, Math.min(canvas.width, 10), Math.min(canvas.height, 10));
                                    const data = imageData.data;
                                    
                                    // Check for non-transparent pixels
                                    let hasContent = false;
                                    for (let i = 3; i < data.length; i += 4) {
                                        if (data[i] > 0) {
                                            hasContent = true;
                                            break;
                                        }
                                    }
                                    if (!hasContent) return false;
                                } catch(e) {
                                    // Security error - assume rendered if canvas exists
                                }
                            }
                        }
                        return true; // All charts ready
                    }""")
                    
                    if charts_ready:
                        logger.info(f"[PDF_GEN] Charts ready after {waited}ms")
                        break
                    
                    page.wait_for_timeout(check_interval)
                    waited += check_interval
                
                # Additional wait for chart animations to complete
                page.wait_for_timeout(2000)
                
                # Check if charts are actually rendered by looking for chart containers
                chart_info = page.evaluate("""() => {
                    const containers = document.querySelectorAll('[data-chart-type]');
                    const canvases = document.querySelectorAll('canvas');
                    const svgs = document.querySelectorAll('svg');
                    return {
                        containers: containers.length,
                        canvases: canvases.length,
                        svgs: svgs.length
                    };
                }""")
                logger.info(f"[PDF_GEN] Found {chart_info.get('containers', 0)} chart containers, {chart_info.get('canvases', 0)} canvas elements, {chart_info.get('svgs', 0)} SVG elements")
                
                # CRITICAL: Check that all charts (canvas + SVG) have actual content
                chart_content_check = """
                () => {
                    const chartContainers = document.querySelectorAll('[data-chart-type]');
                    if (chartContainers.length === 0) {
                        // Fallback: check for any canvas or SVG
                        const canvases = document.querySelectorAll('canvas');
                        const svgs = document.querySelectorAll('svg');
                        return { 
                            ready: canvases.length === 0 && svgs.length === 0, 
                            reason: 'no_charts',
                            renderedCount: 0,
                            totalCharts: 0
                        };
                    }
                    
                    let renderedCount = 0;
                    let totalCharts = chartContainers.length;
                    
                    for (let container of chartContainers) {
                        const chartType = container.dataset.chartType?.toLowerCase() || '';
                        let isRendered = false;
                        
                        // D3-based charts (SVG)
                        if (['sankey', 'sunburst', 'heatmap'].includes(chartType)) {
                            const svg = container.querySelector('svg');
                            if (svg && svg.children.length >= 2) {
                                isRendered = true;
                            }
                        }
                        // Recharts (SVG)
                        else if (['waterfall', 'bubble', 'funnel', 'radialbar', 'treemap', 'composed'].includes(chartType)) {
                            const svg = container.querySelector('.recharts-wrapper svg') || container.querySelector('svg');
                            if (svg && svg.children.length >= 2) {
                                isRendered = true;
                            }
                        }
                        // Chart.js (Canvas)
                        else {
                            const canvas = container.querySelector('canvas');
                            if (canvas) {
                                const ctx = canvas.getContext('2d');
                                if (ctx) {
                                    try {
                                        const imageData = ctx.getImageData(0, 0, Math.min(canvas.width, 10), Math.min(canvas.height, 10));
                                        const data = imageData.data;
                                        
                                        // Check if any pixels are non-transparent
                                        for (let i = 3; i < data.length; i += 4) {
                                            if (data[i] > 0) {
                                                isRendered = true;
                                                break;
                                            }
                                        }
                                    } catch(e) {
                                        // Security error - assume rendered if canvas exists
                                        isRendered = true;
                                    }
                                }
                            }
                        }
                        
                        if (isRendered) renderedCount++;
                    }
                    
                    return { 
                        ready: renderedCount === totalCharts,
                        renderedCount: renderedCount,
                        totalCharts: totalCharts
                    };
                }
                """
                
                # Try to wait for charts to have content (max 10 seconds)
                try:
                    for attempt in range(20):  # 20 attempts x 500ms = 10 seconds
                        result = page.evaluate(chart_content_check)
                        if result.get('ready'):
                            logger.info(f"[PDF_GEN] All {result.get('totalCharts', 0)} charts rendered successfully")
                            break
                        else:
                            logger.info(f"[PDF_GEN] Charts rendering: {result.get('renderedCount', 0)}/{result.get('totalCharts', 0)}")
                            page.wait_for_timeout(500)
                    else:
                        logger.warning(f"[PDF_GEN] Timeout waiting for all charts to render, proceeding anyway")
                except Exception as e:
                    logger.warning(f"[PDF_GEN] Error checking chart content: {e}, proceeding anyway")
                
                # Additional wait for chart animations to complete
                page.wait_for_timeout(2000)
                
                # Wait for any pending animations/renders
                page.evaluate('() => new Promise(resolve => requestAnimationFrame(resolve))')
                
                logger.info("[PDF_GEN] Generating PDF...")
                # Use custom size to match 1024x768 viewport (4:3 aspect ratio)
                # 1024px = 10.67in at 96 DPI, 768px = 8in
                pdf_bytes = page.pdf(
                    width='10.67in',
                    height='8in',
                    print_background=True,
                    margin={'top': '0in', 'right': '0in', 'bottom': '0in', 'left': '0in'},
                    display_header_footer=False,
                    prefer_css_page_size=True,
                    scale=1.0
                )
                
                logger.info(f"[PDF_GEN] PDF generated successfully ({len(pdf_bytes)} bytes)")
                browser.close()
                return pdf_bytes
    
    def _export_to_pdf_reportlab(self, deck_data: Dict[str, Any]) -> bytes:
        """Fallback PDF export using ReportLab (old method)"""
        try:
            # Import reportlab components only when needed
            from reportlab.lib.pagesizes import landscape, LETTER
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            from reportlab.lib.units import inch
            
            # Create PDF buffer
            buffer = io.BytesIO()
            
            # Create PDF document with presentation layout (landscape)
            doc = SimpleDocTemplate(
                buffer,
                pagesize=landscape(LETTER),  # Landscape for presentation format
                rightMargin=72,  # Wider margins for presentation
                leftMargin=72,
                topMargin=54,
                bottomMargin=54,
            )
            
            # Container for the 'Flowable' objects
            elements = []
            
            # Define modern styles
            styles = getSampleStyleSheet()
            
            # Main title style - presentation format (centered, larger)
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=40,  # Larger for presentation
                fontName='Helvetica-Bold',
                textColor=colors.HexColor('#080808'),  # Deep black
                spaceAfter=24,
                alignment=TA_CENTER,  # Center for presentation
                leading=44
            )
            
            # Subtitle style - presentation format
            subtitle_style = ParagraphStyle(
                'CustomSubtitle',
                parent=styles['Heading2'],
                fontSize=24,  # Larger for presentation
                fontName='Helvetica',
                textColor=colors.HexColor('#2D2D2D'),  # Dark gray
                spaceAfter=20,
                alignment=TA_CENTER,  # Center for presentation
                leading=28
            )
            
            # Body text - presentation format (larger, better spacing)
            body_style = ParagraphStyle(
                'ModernBody',
                parent=styles['BodyText'],
                fontSize=14,  # Larger for readability
                fontName='Helvetica',
                textColor=colors.HexColor('#404040'),  # Medium gray
                spaceAfter=12,
                leading=20,
                alignment=TA_LEFT
            )
            
            # Metric label style - presentation format
            metric_label_style = ParagraphStyle(
                'MetricLabel',
                fontSize=12,
                fontName='Helvetica',
                textColor=colors.HexColor('#666666'),  # Muted gray
                spaceAfter=6
            )
            
            # Metric value style - presentation format
            metric_value_style = ParagraphStyle(
                'MetricValue',
                fontSize=28,  # Larger for impact
                fontName='Helvetica-Bold',
                textColor=colors.HexColor('#080808'),  # Deep black
                spaceAfter=16
            )
            
            # Chart title style
            chart_title_style = ParagraphStyle(
                'ChartTitle',
                fontSize=14,
                fontName='Helvetica-Bold',
                textColor=colors.HexColor('#1E293B'),  # Slate-800
                spaceAfter=8,
                alignment=TA_CENTER
            )
            
            # Chart analysis style
            chart_analysis_style = ParagraphStyle(
                'ChartAnalysis',
                fontSize=10,
                fontName='Helvetica',
                textColor=colors.HexColor('#2D2D2D'),  # Dark gray
                spaceAfter=6,
                alignment=TA_LEFT,
                leftIndent=12,
                rightIndent=12,
                borderColor=colors.HexColor('#E2E8F0'),  # Slate-200
                borderWidth=1,
                borderPadding=8,
                backColor=colors.HexColor('#F8FAFC')  # Slate-50
            )
            
            # Process each slide
            slides = deck_data.get("slides", [])
            for i, slide_data in enumerate(slides):
                if i > 0:
                    elements.append(PageBreak())
                
                # Pass all styles
                all_styles = {
                    'title': title_style,
                    'subtitle': subtitle_style,
                    'body': body_style,
                    'metric_label': metric_label_style,
                    'metric_value': metric_value_style,
                    'chart_title': chart_title_style,
                    'chart_analysis': chart_analysis_style,
                    'default': styles
                }
                self._add_pdf_slide(slide_data, elements, all_styles)
            
            # Build PDF
            doc.build(elements)
            
            # Get bytes
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error exporting to PDF: {e}")
            raise
    
    def _get_chart_color(self, index: int) -> str:
        """Get consistent chart color based on index - visually distinct neo-noir shades"""
        colors = [
            'rgba(0, 0, 0, 0.9)',        # Pure black
            'rgba(45, 45, 45, 0.9)',     # Dark charcoal
            'rgba(74, 74, 74, 0.9)',     # Medium charcoal
            'rgba(107, 107, 107, 0.9)',  # Steel gray
            'rgba(138, 138, 138, 0.9)',  # Light steel
            'rgba(176, 176, 176, 0.9)',  # Silver
        ]
        return colors[index % len(colors)]
    
    def _create_sankey_table_representation(self, chart_data: Dict[str, Any], all_styles: Dict) -> None:
        """Convert sankey diagram to table representation for PDF"""
        data = chart_data.get("data", {})
        nodes = data.get("nodes", [])
        links = data.get("links", [])
        
        if not nodes or not links:
            return
        
        # Create table data
        table_data = [["From", "To", "Value"]]
        
        for link in links:
            source_name = nodes[link.get("source", 0)].get("name", "Unknown")
            target_name = nodes[link.get("target", 0)].get("name", "Unknown")
            value = link.get("value", 0)
            
            # Format value
            if isinstance(value, (int, float)):
                if value >= 1_000_000:
                    formatted_value = f"${value/1_000_000:.1f}M"
                elif value >= 1_000:
                    formatted_value = f"${value/1_000:.0f}K"
                else:
                    formatted_value = f"${value:.0f}"
            else:
                formatted_value = str(value)
            
            table_data.append([source_name, target_name, formatted_value])
        
        # Create table
        from reportlab.platypus import Table, TableStyle
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#080808')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#FFFFFF')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8F8F8')),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#E0E0E0')),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#FFFFFF'), colors.HexColor('#F8F8F8')]),
        ]))
        
        return table
    
    def _add_pdf_slide(self, slide_data: Dict[str, Any], elements: List, all_styles: Dict):
        """Add a slide to the PDF document in presentation format"""
        from reportlab.platypus import Paragraph, Spacer
        
        content = slide_data.get("content", {})
        slide_type = slide_data.get("type", "content")
        styles = all_styles['default']
        
        # Add slide number indicator
        slide_number = slide_data.get("order", 0) + 1
        
        # Add title with presentation styling
        title = content.get("title", "")
        if title:
            elements.append(Paragraph(title, all_styles['title']))
            elements.append(Spacer(1, 12))
        
        # Add subtitle
        subtitle = content.get("subtitle", "")
        if subtitle:
            elements.append(Paragraph(subtitle, all_styles['subtitle']))
            elements.append(Spacer(1, 16))
        
        # Add chart with analysis if present
        chart_data = content.get("chart_data") or slide_data.get("chart")
        if chart_data:
            try:
                # Add chart title if present
                chart_title = chart_data.get("options", {}).get("title", {}).get("text", "")
                if chart_title:
                    elements.append(Paragraph(chart_title, all_styles['chart_title']))
                    elements.append(Spacer(1, 8))
                
                # Add the chart visualization
                self._add_pdf_chart(elements, chart_data, all_styles)
                
                # Add chart analysis
                self._add_chart_analysis(elements, chart_data, all_styles)
                
            except Exception as e:
                logger.warning(f"Could not add chart to PDF: {e}")
                # Fall back to table representation
                self._add_pdf_chart_as_table(elements, chart_data, all_styles)
        
        # Add bullets with modern styling
        bullets = content.get("bullets", [])
        if bullets:
            for bullet in bullets:
                # Use modern bullet point with better spacing
                bullet_text = f"<para leftIndent=20><bullet>&bull;</bullet> {bullet}</para>"
                elements.append(Paragraph(bullet_text, all_styles['body']))
                elements.append(Spacer(1, 8))
            elements.append(Spacer(1, 12))
        
        # Add metrics with modern card-style layout
        metrics = content.get("metrics", {})
        if metrics:
            # Create a modern metrics grid
            metric_data = []
            for key, value in metrics.items():
                # Format the value nicely
                if isinstance(value, (int, float)):
                    if value >= 1_000_000:
                        formatted_value = f"${value/1_000_000:.1f}M"
                    elif value >= 1_000:
                        formatted_value = f"${value/1_000:.0f}K"
                    else:
                        formatted_value = f"{value:.1f}%" if value < 100 else str(value)
                else:
                    formatted_value = str(value)
                
                metric_data.append([
                    Paragraph(key.replace('_', ' ').title(), all_styles['metric_label']),
                    Paragraph(formatted_value, all_styles['metric_value'])
                ])
            
            if metric_data:
                # Create modern table without heavy borders
                from reportlab.platypus import Table, TableStyle
                from reportlab.lib.units import inch
                table = Table(metric_data, colWidths=[3*inch, 2*inch])
                table.setStyle(TableStyle([
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 0), (-1, -1), 11),
                    ('LEFTPADDING', (0, 0), (-1, -1), 0),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 12),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LINEBELOW', (0, 0), (-1, -2), 0.5, colors.HexColor('#E2E8F0')),  # Subtle lines
                ]))
                elements.append(table)
                elements.append(Spacer(1, 20))
        
        # Add body text with modern styling
        body = content.get("body", "")
        if body:
            elements.append(Paragraph(body, all_styles['body']))
            elements.append(Spacer(1, 16))
    
    def _add_chart_analysis(self, elements: List, chart_data: Dict[str, Any], all_styles: Dict):
        """Add intelligent analysis of chart data"""
        data = chart_data.get("data", {})
        chart_type = chart_data.get("type", "bar").lower()
        datasets = data.get("datasets", [])
        labels = data.get("labels", [])
        
        if not datasets or not labels:
            return
        
        analysis_points = []
        
        # Analyze based on chart type
        if chart_type == "bar" and datasets:
            # Find highest and lowest values
            first_dataset = datasets[0].get("data", [])
            if first_dataset and len(first_dataset) > 0:
                # Filter out non-numeric values that could cause max() to fail
                numeric_values = [v for v in first_dataset if isinstance(v, (int, float)) and not isinstance(v, bool)]
                if numeric_values:
                    max_val = max(numeric_values)
                    min_val = min(numeric_values)
                    max_idx = first_dataset.index(max_val)
                    min_idx = first_dataset.index(min_val)
                
                # Format values nicely
                def format_value(val):
                    if val >= 1_000_000:
                        return f"${val/1_000_000:.1f}M"
                    elif val >= 1_000:
                        return f"${val/1_000:.0f}K"
                    else:
                        return f"{val:.0f}"
                
                if labels:
                    analysis_points.append(
                        f"• Highest: {labels[max_idx]} at {format_value(max_val)}"
                    )
                    analysis_points.append(
                        f"• Lowest: {labels[min_idx]} at {format_value(min_val)}"
                    )
                    
                    # Calculate growth if sequential
                    if len(first_dataset) > 1:
                        growth = ((first_dataset[-1] / first_dataset[0]) - 1) * 100
                        analysis_points.append(
                            f"• Growth: {growth:.1f}% from {labels[0]} to {labels[-1]}"
                        )
        
        elif chart_type == "line" and datasets:
            # Analyze trend
            for i, dataset in enumerate(datasets):
                dataset_data = dataset.get("data", [])
                if len(dataset_data) > 1:
                    # Calculate CAGR if time series
                    start_val = dataset_data[0]
                    end_val = dataset_data[-1]
                    if start_val > 0:
                        periods = len(dataset_data) - 1
                        cagr = ((end_val / start_val) ** (1/periods) - 1) * 100
                        dataset_label = dataset.get("label", f"Series {i+1}")
                        analysis_points.append(
                            f"• {dataset_label}: {cagr:.1f}% CAGR"
                        )
        
        elif chart_type == "pie" and datasets:
            # Show top segments
            first_dataset = datasets[0].get("data", [])
            if first_dataset and labels:
                total = sum(first_dataset)
                segments = list(zip(labels, first_dataset))
                segments.sort(key=lambda x: x[1], reverse=True)
                
                for label, value in segments[:3]:  # Top 3
                    percentage = (value / total * 100) if total > 0 else 0
                    analysis_points.append(
                        f"• {label}: {percentage:.1f}% of total"
                    )
        
        # Add analysis section if we have points
        if analysis_points:
            elements.append(Spacer(1, 8))
            
            # Add analysis box
            analysis_text = "<b>Key Insights:</b><br/>" + "<br/>".join(analysis_points)
            analysis_para = Paragraph(analysis_text, all_styles.get('chart_analysis', all_styles['body']))
            
            # Create a bordered box for analysis
            from reportlab.platypus import Table, TableStyle
            analysis_table = Table([[analysis_para]], colWidths=[6.5*inch])
            analysis_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#F8FAFC')),  # Light background
                ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0')),  # Border
                ('LEFTPADDING', (0, 0), (-1, -1), 12),
                ('RIGHTPADDING', (0, 0), (-1, -1), 12),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ]))
            elements.append(analysis_table)
            elements.append(Spacer(1, 12))
    
    def _add_pdf_chart(self, elements: List, chart_data: Dict[str, Any], all_styles: Dict):
        """Add a chart to PDF using reportlab with better styling"""
        from reportlab.graphics.shapes import Drawing
        from reportlab.graphics.charts.barcharts import VerticalBarChart
        from reportlab.graphics.charts.piecharts import Pie
        from reportlab.graphics.charts.linecharts import HorizontalLineChart
        from reportlab.lib import colors as rlcolors
        
        data = chart_data.get("data", {})
        chart_type = chart_data.get("type", "bar").lower()
        
        # Handle grouped bar charts (convert to regular bar chart for PDF)
        if chart_type == 'grouped_bar':
            chart_type = 'bar'
            # Flatten grouped data for PDF compatibility
            if 'series' in data:
                labels = data.get('categories', [])
                datasets = []
                for series in data['series']:
                    datasets.append({
                        'label': series['name'],
                        'data': series['data'],
                        'backgroundColor': self._get_chart_color(len(datasets))
                    })
                data = {'labels': labels, 'datasets': datasets}
        
        # Handle waterfall charts - convert to stacked bar for PDF compatibility
        if chart_type == 'waterfall':
            chart_type = 'bar'
            # Convert waterfall data to stacked bar format
            if 'data' in data and isinstance(data['data'], list):
                labels = [f"Step {i+1}" for i in range(len(data['data']))]
                datasets = [{
                    'label': 'Waterfall',
                    'data': data['data'],
                    'backgroundColor': self._get_chart_color(0)
                }]
                data = {'labels': labels, 'datasets': datasets}
        
        # Handle sankey diagrams - convert to simple bar chart for PDF
        if chart_type == 'sankey':
            chart_type = 'bar'
            # Convert sankey data to bar format
            if 'nodes' in data and 'links' in data:
                # Extract node names and create simple bar chart
                nodes = data.get('nodes', [])
                labels = [node.get('name', f'Node {i}') for i, node in enumerate(nodes[:10])]  # Limit to 10 nodes
                datasets = [{
                    'label': 'Value',
                    'data': [1] * len(labels),  # Placeholder values
                    'backgroundColor': self._get_chart_color(0)
                }]
                data = {'labels': labels, 'datasets': datasets}
        
        labels = data.get("labels", [])
        datasets = data.get("datasets", [])
        
        if not labels or not datasets:
            return
        
        # Larger drawing for better visibility
        drawing = Drawing(500, 250)
        
        if chart_type == "pie" and datasets:
            # Create larger pie chart with better positioning
            pie = Pie()
            pie.x = 150
            pie.y = 25
            pie.width = 200
            pie.height = 200
            
            # Use first dataset for pie chart
            pie.data = datasets[0].get("data", [])
            pie.labels = labels
            
            # Add colors
            pie.slices.strokeWidth = 0.5
            pie_colors = [
                rlcolors.HexColor('#4F46E5'),
                rlcolors.HexColor('#7C3AED'),
                rlcolors.HexColor('#EC4899'),
                rlcolors.HexColor('#F59E0B'),
                rlcolors.HexColor('#10B981'),
                rlcolors.HexColor('#3B82F6'),
            ]
            for i in range(len(pie.data)):
                pie.slices[i].fillColor = pie_colors[i % len(pie_colors)]
            
            drawing.add(pie)
            
        elif chart_type == "line":
            # Create larger line chart with better positioning
            lc = HorizontalLineChart()
            lc.x = 65
            lc.y = 50
            lc.height = 150
            lc.width = 400
            
            # Format data for line chart
            chart_values = []
            for dataset in datasets:
                chart_values.append(dataset.get("data", []))
            
            lc.data = chart_values
            lc.categoryAxis.categoryNames = labels
            lc.valueAxis.valueMin = 0
            # Safely calculate max value
            try:
                max_val = max(max(d) for d in chart_values if d) * 1.1
            except (ValueError, TypeError):
                max_val = 100
            lc.valueAxis.valueMax = max_val
            lc.lines[0].strokeColor = rlcolors.blue
            lc.lines[0].strokeWidth = 2
            
            drawing.add(lc)
            
        else:  # Default to bar chart
            # Create larger bar chart with better positioning
            bc = VerticalBarChart()
            bc.x = 65
            bc.y = 50
            bc.height = 150
            bc.width = 400
            
            # Format data for bar chart
            chart_values = []
            for dataset in datasets:
                chart_values.append(dataset.get("data", []))
            
            bc.data = chart_values
            bc.categoryAxis.categoryNames = labels
            bc.valueAxis.valueMin = 0
            # Safely calculate max value
            try:
                max_val = max(max(d) for d in chart_values if d) * 1.1
            except (ValueError, TypeError):
                max_val = 100
            bc.valueAxis.valueMax = max_val
            
            # Set bar colors
            bar_colors = [
                rlcolors.HexColor('#4F46E5'),
                rlcolors.HexColor('#7C3AED'),
                rlcolors.HexColor('#EC4899'),
            ]
            for i, dataset in enumerate(datasets):
                bc.bars[i].fillColor = bar_colors[i % len(bar_colors)]
            
            drawing.add(bc)
        
        # Add axis labels if available
        x_label = chart_data.get("options", {}).get("scales", {}).get("x", {}).get("title", {}).get("text", "")
        y_label = chart_data.get("options", {}).get("scales", {}).get("y", {}).get("title", {}).get("text", "")
        
        if x_label or y_label:
            from reportlab.graphics.shapes import String
            if x_label:
                x_label_obj = String(250, 15, x_label, fontSize=10, textAnchor='middle')
                drawing.add(x_label_obj)
            if y_label:
                # Rotate y-axis label
                from reportlab.graphics.shapes import rotate
                y_label_obj = String(25, 150, y_label, fontSize=10, textAnchor='middle')
                drawing.add(rotate(y_label_obj, 90))
        
        elements.append(drawing)
        
        from reportlab.platypus import Spacer, Paragraph
        from reportlab.lib.enums import TA_CENTER
        from reportlab.lib.styles import ParagraphStyle
        elements.append(Spacer(1, 12))
        
        if chart_data.get("title"):
            chart_title_style = ParagraphStyle(
                'ChartTitle',
                fontSize=10,
                alignment=TA_CENTER,
                textColor=rlcolors.HexColor('#666666')
            )
            elements.append(Paragraph(str(chart_data.get("title")), chart_title_style))
            elements.append(Spacer(1, 12))
    
    def _add_pdf_chart_as_table(self, elements: List, chart_data: Dict[str, Any], all_styles: Dict):
        """Add chart data as table when chart rendering fails"""
        data = chart_data.get("data", {})
        labels = data.get("labels", [])
        datasets = data.get("datasets", [])
        
        if not labels or not datasets:
            return
        
        # Create table data
        table_data = [["Category"] + [ds.get("label", f"Series {i+1}") for i, ds in enumerate(datasets)]]
        
        for i, label in enumerate(labels):
            row = [str(label)]
            for dataset in datasets:
                values = dataset.get("data", [])
                if i < len(values):
                    value = values[i]
                    if isinstance(value, (int, float)):
                        row.append(f"{value:,.0f}")
                    else:
                        row.append(str(value))
                else:
                    row.append("")
            table_data.append(row)
        
        # Create and style table
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black)
        ]))
        
        elements.append(table)
        elements.append(Spacer(1, 12))
    
    def _html_followon_strategy_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate follow-on strategy slide HTML"""
        companies = content.get('companies', {})
        strategy_html = ""
        
        for company_name, strategy in companies.items():
            strategy_html += f"""
            <div class="bg-white rounded-lg p-6 mb-4 border border-gray-200">
                <h3 class="text-xl font-bold mb-4">{company_name}</h3>
                <div class="grid grid-cols-2 gap-4">
                    <div>
                        <p class="text-sm text-gray-600">Entry</p>
                        <p class="text-lg font-semibold">{strategy.get('entry_investment', 'TBD')}</p>
                        <p class="text-sm">{strategy.get('entry_ownership', 'TBD')}</p>
                    </div>
                    <div>
                        <p class="text-sm text-gray-600">Exit (with follow-on)</p>
                        <p class="text-lg font-semibold">{strategy.get('total_invested', 'TBD')}</p>
                        <p class="text-sm">{strategy.get('exit_ownership', 'TBD')}</p>
                    </div>
                </div>
                <div class="mt-4 p-3 bg-gray-50 rounded">
                    <p class="text-sm">{strategy.get('recommendation', 'Analysis pending')}</p>
                </div>
            </div>
            """
        
        return f"""
<div class="slide content-slide">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Follow-on Strategy Decision Framework')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    {strategy_html if strategy_html else '<p class="text-gray-500">No follow-on strategy data available</p>'}
</div>
        """
    
    def _html_fund_impact_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate fund impact slide HTML"""
        chart_html = ""
        
        # Add chart if available
        if content.get('chart_data'):
            chart_id = f"fund-impact-chart-{slide_idx}"
            chart_html = f"""
            <div class="bg-white rounded-lg p-4 mb-6 border border-gray-200">
                <canvas id="{chart_id}" style="max-height: 400px;"></canvas>
            </div>
            """
        
        # Add metrics table
        metrics_html = ""
        if content.get('fund_metrics'):
            metrics = content['fund_metrics']
            metrics_html = f"""
            <div class="grid grid-cols-3 gap-4 mb-6">
                <div class="bg-blue-50 rounded-lg p-4">
                    <p class="text-sm text-blue-600 font-semibold">Fund Size</p>
                    <p class="text-2xl font-bold text-blue-900">{metrics.get('fund_size', '$234M')}</p>
                </div>
                <div class="bg-green-50 rounded-lg p-4">
                    <p class="text-sm text-green-600 font-semibold">Deployed</p>
                    <p class="text-2xl font-bold text-green-900">{metrics.get('deployed', '$125M')}</p>
                </div>
                <div class="bg-gray-100 rounded-lg p-4 border-l-3 border-gray-700">
                    <p class="text-sm text-gray-600 font-semibold">DPI Target</p>
                    <p class="text-2xl font-semibold text-gray-900">{metrics.get('dpi_target', '3.0x')}</p>
                </div>
            </div>
            """
        
        return f"""
<div class="slide content-slide">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Impact on $234M Fund Returns')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', '')}</p>
    {metrics_html}
    {chart_html}
</div>
        """
    
    def _html_risk_analysis_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate risk analysis slide HTML"""
        risks = content.get('risks', [])
        risks_html = ""
        
        for risk in risks:
            severity_color = {
                'high': 'red',
                'medium': 'yellow',
                'low': 'green'
            }.get(risk.get('severity', 'medium'), 'yellow')
            
            risks_html += f"""
            <div class="bg-white rounded-lg p-4 mb-3 border-l-4 border-{severity_color}-500">
                <div class="flex justify-between items-start mb-2">
                    <h4 class="font-semibold text-gray-900">{risk.get('title', 'Risk')}</h4>
                    <span class="px-2 py-1 text-xs font-semibold bg-{severity_color}-100 text-{severity_color}-800 rounded">
                        {risk.get('severity', 'Medium').upper()}
                    </span>
                </div>
                <p class="text-sm text-gray-700 mb-2">{risk.get('description', '')}</p>
                <p class="text-sm text-gray-600 italic">Mitigation: {risk.get('mitigation', 'TBD')}</p>
            </div>
            """
        
        return f"""
<div class="slide content-slide">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Risk Analysis')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', 'Key risks and mitigation strategies')}</p>
    {risks_html if risks_html else '<p class="text-gray-500">Risk analysis in progress</p>'}
</div>
        """
    
    def _html_probability_cloud_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate probability cloud visualization slide HTML"""
        # Check if we have chart data for the probability cloud
        chart_data = content.get('chart_data')
        chart_html = ""
        
        if chart_data:
            chart_id = f"prob-cloud-{slide_idx}"
            
            # Check if it's a proper probability cloud chart
            if chart_data.get('type') == 'probability_cloud':
                # Use canvas for Chart.js rendering
                chart_html = f"""
                <div class="bg-white rounded-lg p-6 mb-6">
                    <canvas id="{chart_id}" style="max-height: 500px;"></canvas>
                </div>
                """
            else:
                # Fallback to regular chart
                chart_html = f"""
                <div class="bg-white rounded-lg p-6 mb-6">
                    <canvas id="{chart_id}" style="max-height: 400px;"></canvas>
                </div>
                """
        
        # Add description
        description_html = ""
        if content.get('description'):
            description_html = f"""
            <div class="bg-blue-50 border-l-4 border-blue-500 rounded-lg p-4 mb-6">
                <p class="text-gray-700">{content.get('description')}</p>
            </div>
            """
        
        return f"""
<div class="slide content-slide">
    <h2 class="text-3xl font-bold text-gray-900 mb-2">{content.get('title', 'Exit Probability Cloud')}</h2>
    <p class="text-gray-600 mb-6">{content.get('subtitle', 'Return distribution across exit scenarios')}</p>
    {description_html}
    {chart_html}
    <div class="text-sm text-gray-500 mt-4">
        <p>• Shows return scenarios with probability distributions</p>
        <p>• Includes defensive breakpoints and decision zones</p>
        <p>• Dynamic calculation based on cap table evolution</p>
    </div>
</div>
        """
    
    def _html_breakpoint_analysis_slide(self, content: Dict[str, Any], slide_idx: int) -> str:
        """Generate breakpoint analysis slide HTML with improved layout and chart support"""
        # Chart visualization if available
        chart_html = ""
        if content.get("chart_data"):
            chart_data = content["chart_data"]
            chart_type = chart_data.get("type", "bar")
            chart_id = f"chart-{slide_idx}"
            
            # Generate chart container with proper initialization
            chart_html = f"""
            <div class="mb-4">
                <div class="chart-container" data-chart-type="{chart_type}" id="{chart_id}">
                    <canvas id="{chart_id}"></canvas>
                </div>
            </div>
            """
        
        # Reality check table - make it more compact
        reality_table_html = ""
        if content.get("reality_check"):
            reality_table_html = """
            <div class="bg-red-50 border-2 border-red-200 rounded-lg p-3 mb-4">
                <h3 class="text-base font-bold text-red-900 mb-2">Reality Check - Key Exit Scenarios</h3>
                <div class="overflow-x-auto">
                    <table class="w-full text-xs">
                        <thead>
                            <tr class="border-b border-red-200">
                                <th class="text-left py-1.5 pr-2">Company</th>
                                <th class="text-left py-1.5 pr-2">Exit Value</th>
                                <th class="text-left py-1.5 pr-2">Our Return</th>
                                <th class="text-left py-1.5 pr-2">MOIC</th>
                                <th class="text-left py-1.5">DPI Impact</th>
                            </tr>
                        </thead>
                        <tbody>
            """
            for row in content["reality_check"]:
                # Parse MOIC to check if it's below 1
                moic_str = str(row.get('moic', '1x'))
                try:
                    moic_val = float(moic_str.replace('x', ''))
                    moic_class = "text-red-600 font-bold" if moic_val < 1 else "text-gray-700"
                except:
                    moic_class = "text-gray-700"
                
                reality_table_html += f"""
                            <tr class="border-b border-red-100">
                                <td class="py-1.5 pr-2 font-semibold">{row.get('company', '')}</td>
                                <td class="py-1.5 pr-2">{row.get('exit_value', '')}</td>
                                <td class="py-1.5 pr-2">{row.get('our_proceeds', '')}</td>
                                <td class="py-1.5 pr-2 {moic_class}">{row.get('moic', '')}</td>
                                <td class="py-1.5">{row.get('dpi_contribution', '')}</td>
                            </tr>
                """
            reality_table_html += """
                        </tbody>
                    </table>
                </div>
            </div>
            """
        
        # Companies breakpoints - improved layout
        companies_html = ""
        companies = content.get("companies", {})
        if companies:
            companies_html = '<div class="grid grid-cols-2 gap-3 mb-4">'
            for company_name, company_data in list(companies.items())[:2]:
                if company_data.get("breakpoints"):
                    bp = company_data["breakpoints"]
                    
                    # Helper function to safely get and format breakpoint values
                    def format_bp(key, default=0):
                        val = bp.get(key, default)
                        if val is None or val == 0:
                            return "N/A"
                        return f"${val/1e6:.0f}M"
                    
                    companies_html += f"""
                    <div class="bg-gray-50 rounded-lg p-3 border-l-4 border-gray-800">
                        <h3 class="font-semibold text-sm text-gray-900 mb-2">{company_name}</h3>
                        <div class="space-y-2.5 text-xs">
                            <div class="border-b border-gray-300 pb-2 mb-2">
                                <div class="text-xs font-bold text-gray-700 mb-1.5">Without Pro Rata:</div>
                                <div class="space-y-1.5 pl-2">
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">1x Return:</span>
                                        <span class="font-semibold text-gray-900 text-right">{format_bp("exit_for_1x_return_no_pro_rata")}</span>
                                    </div>
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">2x Return:</span>
                                        <span class="font-semibold text-gray-900 text-right">{format_bp("exit_for_2x_return_no_pro_rata")}</span>
                                    </div>
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">3x Return:</span>
                                        <span class="font-semibold text-gray-900 text-right">{format_bp("exit_for_3x_return_no_pro_rata")}</span>
                                    </div>
                                </div>
                            </div>
                            <div class="pt-1">
                                <div class="text-xs font-bold text-gray-700 mb-1.5">With Pro Rata:</div>
                                <div class="space-y-1.5 pl-2">
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">1x Return:</span>
                                        <span class="font-semibold text-blue-700 text-right">{format_bp("exit_for_1x_return_with_pro_rata")}</span>
                                    </div>
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">2x Return:</span>
                                        <span class="font-semibold text-blue-700 text-right">{format_bp("exit_for_2x_return_with_pro_rata")}</span>
                                    </div>
                                    <div class="flex justify-between items-center gap-2">
                                        <span class="text-gray-600 flex-shrink-0">3x Return:</span>
                                        <span class="font-semibold text-blue-700 text-right">{format_bp("exit_for_3x_return_with_pro_rata")}</span>
                                    </div>
                                </div>
                            </div>
                            <div class="text-xs text-gray-500 italic pt-1.5 border-t border-gray-200 mt-2">
                                All values in millions ($M)
                            </div>
                        </div>
                    </div>
                    """
            companies_html += '</div>'
        
        # Insights - more compact
        insights_html = ""
        if content.get("insights"):
            insights = content["insights"]
            if isinstance(insights, list):
                insights_html = '<div class="bg-yellow-50 border border-yellow-200 rounded-lg p-3"><h4 class="text-sm font-semibold text-yellow-900 mb-2">Key Insights</h4><ul class="space-y-1">'
                for insight in insights[:4]:  # Limit to 4 for better layout
                    insights_html += f'<li class="text-xs text-yellow-800">• {insight}</li>'
                insights_html += '</ul></div>'
            else:
                insights_html = f'<div class="bg-yellow-50 border border-yellow-200 rounded-lg p-3"><h4 class="text-sm font-semibold text-yellow-900 mb-2">Key Insights</h4><p class="text-xs text-yellow-800">{insights}</p></div>'
        
        return f"""
<div class="slide content-slide">
    <h2 class="text-2xl font-bold text-gray-900 mb-1">{content.get('title', 'Breakpoint Analysis')}</h2>
    <p class="text-sm text-gray-600 mb-4">{content.get('subtitle', 'Key inflection points for returns')}</p>
    {chart_html}
    {reality_table_html}
    {companies_html}
    {insights_html}
</div>
        """