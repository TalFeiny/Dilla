#!/usr/bin/env python3

from pptx import Presentation
from pptx.chart.data import CategoryChartData
import json

def extract_detailed_content(filepath):
    """Extract detailed content including chart data."""
    prs = Presentation(filepath)
    
    # Focus on key slides with issues
    key_slides = {}
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        all_text = []
        chart_info = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_text.append(text)
                    if not slide_title and idx == 1:
                        slide_title = text
            
            # Extract chart data
            if shape.has_chart:
                chart = shape.chart
                chart_data = {
                    "chart_type": str(chart.chart_type),
                    "title": chart.chart_title.text_frame.text if chart.has_title else "No title"
                }
                
                # Try to extract series data
                try:
                    series_data = []
                    for series in chart.series:
                        series_info = {
                            "name": series.name,
                            "values": []
                        }
                        # Get values if accessible
                        try:
                            for point in series.points:
                                series_info["values"].append(str(point.data_label.text_frame.text) if point.has_data_label else "N/A")
                        except:
                            pass
                        series_data.append(series_info)
                    chart_data["series"] = series_data
                except:
                    chart_data["series"] = "Unable to extract"
                    
                chart_info.append(chart_data)
        
        # Store detailed info for problematic slides
        if idx in [2, 4, 5, 6, 8, 9]:  # Key slides to analyze
            key_slides[f"Slide {idx}"] = {
                "all_text": all_text,
                "charts": chart_info,
                "issues": []
            }
    
    # Identify specific issues
    # Slide 2 - Executive Summary
    if "Slide 2" in key_slides:
        text = " ".join(key_slides["Slide 2"]["all_text"])
        if "$119,000" in text or "119" in text:
            key_slides["Slide 2"]["issues"].append("Inven revenue incorrectly shown as $119K (should be higher for Series A)")
        if "$63.5M" in text:
            key_slides["Slide 2"]["issues"].append("Inven valuation of $63.5M seems low for Series A")
            
    # Slide 4 - Company Comparison Table  
    if "Slide 4" in key_slides:
        text = " ".join(key_slides["Slide 4"]["all_text"])
        if "$119,000" in text:
            key_slides["Slide 4"]["issues"].append("CRITICAL: Inven revenue shown as $119K - likely data extraction error")
        if "$2,000,000" in text:
            key_slides["Slide 4"]["issues"].append("Farsight revenue of $2M seems low for $80M valuation")
            
    # Slide 6 - Business Model (empty)
    if "Slide 6" in key_slides:
        if not key_slides["Slide 6"]["charts"]:
            key_slides["Slide 6"]["issues"].append("NO CONTENT: Business Model slide is completely empty")
            
    # Slide 9 - Cap Table (empty)
    if "Slide 9" in key_slides:
        if not key_slides["Slide 9"]["charts"]:
            key_slides["Slide 9"]["issues"].append("NO CONTENT: Cap Table Comparison slide is empty (should have Sankey diagrams)")
    
    return key_slides

if __name__ == "__main__":
    result = extract_detailed_content("/Users/admin/Downloads/deck-1758994559903.pptx")
    print(json.dumps(result, indent=2, default=str))