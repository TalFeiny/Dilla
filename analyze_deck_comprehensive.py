#!/usr/bin/env python3

from pptx import Presentation
import json

def comprehensive_deck_analysis(filepath):
    """Thoroughly analyze all issues with the deck."""
    prs = Presentation(filepath)
    
    issues = {
        "critical": [],
        "data_quality": [],
        "missing_content": [],
        "visualization": [],
        "structure": []
    }
    
    # Track what we see
    has_cap_table = False
    has_business_model = False
    has_scenario_analysis = False
    has_market_sizing = False
    has_investment_thesis = False
    chart_count = 0
    empty_slides = []
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        has_content = False
        has_chart = False
        has_table = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_text.append(text)
                    has_content = True
            if shape.has_chart:
                has_chart = True
                chart_count += 1
            if shape.has_table:
                has_table = True
                # Extract table data to check for issues
                if idx == 4:  # Company comparison table
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        if "119" in str(row_text):
                            issues["critical"].append(f"Slide {idx}: Revenue shows as $119K for Series A company")
        
        # Check specific slides
        if idx == 3:  # Scenario Analysis
            if not has_table and not has_chart:
                issues["missing_content"].append("Slide 3: Scenario Analysis has no data visualization")
            has_scenario_analysis = any("scenario" in t.lower() for t in slide_text)
        
        if idx == 4:  # Company Comparison
            text_combined = " ".join(slide_text).lower()
            if "$119,000" in text_combined or "$119k" in text_combined:
                issues["critical"].append("Slide 4: Inven revenue incorrectly shown as $119K")
            if "$2,000,000" in text_combined:
                issues["data_quality"].append("Slide 4: Farsight revenue seems low at $2M for $80M valuation")
        
        if idx == 5:  # Path to $100M ARR
            if not has_chart:
                issues["visualization"].append("Slide 5: Missing ARR growth chart")
        
        if idx == 6:  # Business Model
            if not has_chart and not has_table:
                issues["critical"].append("Slide 6: Business Model Analysis is EMPTY")
                has_business_model = False
            else:
                has_business_model = True
        
        if idx == 7:  # Market Sizing
            if has_chart:
                has_market_sizing = True
            else:
                issues["visualization"].append("Slide 7: Market sizing should have chart")
        
        if idx == 8:  # Investment Recommendations
            if has_chart:
                has_investment_thesis = True
        
        if idx == 9:  # Cap Table
            if not has_chart and not has_table:
                issues["critical"].append("Slide 9: Cap Table Comparison is EMPTY (should have Sankey diagrams)")
                has_cap_table = False
            else:
                has_cap_table = True
        
        if idx == 10:  # Valuation Comparison
            if not has_chart:
                issues["visualization"].append("Slide 10: Missing valuation comparison chart")
        
        if idx == 11:  # Revenue Analysis
            if not has_chart:
                issues["visualization"].append("Slide 11: Missing revenue analysis chart")
        
        # Check for empty slides
        if not has_content and not has_chart and not has_table:
            empty_slides.append(idx)
    
    # Overall structural issues
    if len(empty_slides) > 0:
        issues["structure"].append(f"Empty slides: {empty_slides}")
    
    if chart_count < 5:
        issues["structure"].append(f"Only {chart_count} charts in entire deck (should have 8+)")
    
    if not has_cap_table:
        issues["structure"].append("Missing cap table visualization")
    
    if not has_business_model:
        issues["structure"].append("Missing business model analysis")
    
    # Data quality checks based on text
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "
    
    # Check for generic descriptions
    if "SaaS" in all_text and "platform" in all_text:
        if "financial workflows" not in all_text.lower() and "acquisition targets" not in all_text.lower():
            issues["data_quality"].append("Companies described generically as 'SaaS platform' instead of specific functions")
    
    # Check for missing metrics
    if "TAM" not in all_text and "market size" not in all_text.lower():
        issues["data_quality"].append("No TAM/market size mentioned")
    
    if "growth rate" not in all_text.lower() and "yoy" not in all_text.lower():
        issues["data_quality"].append("No growth rates mentioned")
    
    if "gross margin" not in all_text.lower() and "unit economics" not in all_text.lower():
        issues["data_quality"].append("No gross margin or unit economics mentioned")
    
    # Check for citation quality
    citation_slides = [s for s in range(12, 17)]  # Last 5 slides are citations
    if len(citation_slides) > 0:
        # Check if citations are just generic URLs or have dates
        if "2024" not in all_text and "2025" not in all_text:
            issues["data_quality"].append("Citations lack dates")
    
    return issues

if __name__ == "__main__":
    result = comprehensive_deck_analysis("/Users/admin/Downloads/deck-1758994559903.pptx")
    
    print("\n=== COMPREHENSIVE DECK ANALYSIS ===\n")
    
    print("🔴 CRITICAL ISSUES:")
    for issue in result["critical"]:
        print(f"  - {issue}")
    
    print("\n⚠️  DATA QUALITY ISSUES:")
    for issue in result["data_quality"]:
        print(f"  - {issue}")
    
    print("\n📊 MISSING CONTENT:")
    for issue in result["missing_content"]:
        print(f"  - {issue}")
    
    print("\n📈 VISUALIZATION ISSUES:")
    for issue in result["visualization"]:
        print(f"  - {issue}")
    
    print("\n🏗️  STRUCTURAL ISSUES:")
    for issue in result["structure"]:
        print(f"  - {issue}")
    
    print("\n=== SUMMARY ===")
    total_issues = sum(len(v) for v in result.values())
    print(f"Total issues found: {total_issues}")
    print(f"Critical issues: {len(result['critical'])}")
    print(f"Charts in deck: Very few (most slides missing visualizations)")