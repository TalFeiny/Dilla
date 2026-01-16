#!/usr/bin/env python3

import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

def analyze_deck(filepath):
    """Analyze a PowerPoint deck and extract its contents."""
    try:
        prs = Presentation(filepath)
        analysis = {
            "slide_count": len(prs.slides),
            "slides": []
        }
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_info = {
                "slide_number": idx,
                "title": "",
                "text_content": [],
                "has_chart": False,
                "has_table": False,
                "has_image": False,
                "shape_count": len(slide.shapes)
            }
            
            # Extract text from all shapes
            for shape in slide.shapes:
                # Check for title
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if not slide_info["title"] and shape == slide.shapes[0]:
                            slide_info["title"] = text
                        else:
                            slide_info["text_content"].append(text)
                
                # Check for charts
                if shape.has_chart:
                    slide_info["has_chart"] = True
                    
                # Check for tables
                if shape.has_table:
                    slide_info["has_table"] = True
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    slide_info["table_data"] = table_data
                    
                # Check for images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info["has_image"] = True
            
            analysis["slides"].append(slide_info)
        
        return analysis
        
    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    filepath = "/Users/admin/Downloads/deck-1758994559903.pptx"
    result = analyze_deck(filepath)
    print(json.dumps(result, indent=2))